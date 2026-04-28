#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génération PPTX professionnel — Audit Sécurité M2 Cybersécurité
Couleurs : Noir / Bleu marine / Blanc
25 slides avec explications en bas de chaque slide
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

# ── Couleurs ──────────────────────────────────────────────────────────────────
BLACK    = RGBColor(0x06, 0x0e, 0x1a)
NAVY_D   = RGBColor(0x0a, 0x16, 0x28)
NAVY     = RGBColor(0x0d, 0x1f, 0x3c)
NAVY_M   = RGBColor(0x14, 0x28, 0x50)
NAVY_L   = RGBColor(0x1a, 0x35, 0x66)
WHITE    = RGBColor(0xff, 0xff, 0xff)
WHITE_DIM= RGBColor(0xb8, 0xc8, 0xd8)
GREY     = RGBColor(0x6a, 0x80, 0x98)

BASE = os.path.dirname(os.path.abspath(__file__))

def img(name):
    """Retourne le chemin absolu d'une image"""
    paths = [
        os.path.join(BASE, name),
        os.path.join(BASE, 'WAZUCAPT', os.path.basename(name)),
        os.path.join(BASE, 'CAPMONITORINGSECU', os.path.basename(name)),
    ]
    for p in paths:
        if os.path.exists(p):
            return p
    return None

WAZUH = {
    'agents':    img('WAZUCAPT/wazuh_agents.png'),
    'agents01':  img('WAZUCAPT/wazuh_01_agents.png'),
    'overview':  img('WAZUCAPT/wazuh_02_overview.png'),
    'fim':       img('WAZUCAPT/wazuh_03_fim.png'),
    'cve':       img('WAZUCAPT/wazuh_04_cve.png'),
    'mitre':     img('WAZUCAPT/wazuh_05_mitre.png'),
}
CAP = {
    'siem':      img('CAPMONITORINGSECU/brave_screenshot_localhost.png'),
    'rgpd':      img('CAPMONITORINGSECU/brave_screenshot_localhost (1).png'),
    'rbac':      img('CAPMONITORINGSECU/brave_screenshot_localhost (2).png'),
    'waf':       img('CAPMONITORINGSECU/brave_screenshot_localhost (3).png'),
    'score':     img('CAPMONITORINGSECU/brave_screenshot_localhost (4).png'),
}

# ── Dimensions 16:9 ───────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # blank layout

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=NAVY_D):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color=NAVY, border_color=None, alpha=None):
    from pptx.util import Pt
    shp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return tb

def txbox_multi(slide, lines, x, y, w, h, size=11, color=WHITE_DIM):
    """lines = list of (text, bold, size_override)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, sz) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz if sz else size)
        run.font.bold = bold
        run.font.color.rgb = WHITE if bold else color
        run.font.name = 'Calibri'

def add_image(slide, path, x, y, w, h):
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, x, y, w, h)
        except Exception as e:
            print(f"  [WARN] Image non chargée {path}: {e}")

def slide_num(slide, n, total=25):
    txbox(slide, f'{n:02d} / {total}',
          W - Inches(1.5), Inches(0.1), Inches(1.4), Inches(0.3),
          size=9, color=GREY, align=PP_ALIGN.RIGHT)

def logo_txt(slide):
    txbox(slide, 'Cybersécurité M2 — YNOV Campus',
          Inches(0.3), Inches(0.1), Inches(4), Inches(0.3),
          size=8, color=GREY)

def label(slide, text, x=Inches(0.5), y=Inches(0.55)):
    txbox(slide, text, x, y, Inches(8), Inches(0.25),
          size=8, color=GREY, bold=True)

def title(slide, text, x=Inches(0.5), y=Inches(0.85), size=26):
    txbox(slide, text, x, y, Inches(12.5), Inches(1.0),
          size=size, bold=True, color=WHITE)

def divider(slide, x=Inches(0.5), y=Inches(1.85), w=Inches(1.0)):
    r = slide.shapes.add_shape(1, x, y, w, Pt(2))
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.fill.background()

def notes_zone(slide, text):
    """Zone bleu très sombre en bas avec explication"""
    notes_y = H - Inches(1.55)
    notes_h = Inches(1.45)
    r = slide.shapes.add_shape(1, 0, notes_y, W, notes_h)
    r.fill.solid(); r.fill.fore_color.rgb = BLACK
    r.line.color.rgb = RGBColor(0x1e, 0x40, 0x80)
    r.line.width = Pt(0.5)
    # Label
    txbox(slide, 'EXPLICATION — NOTE DE L\'ORATEUR',
          Inches(0.4), notes_y + Pt(4), Inches(5), Inches(0.25),
          size=7, color=GREY, bold=True)
    # Text
    txbox(slide, text,
          Inches(0.4), notes_y + Inches(0.28), W - Inches(0.8), notes_h - Inches(0.35),
          size=9, color=WHITE_DIM, wrap=True)

def kpi_bar(slide, kpis, y=Inches(5.6), h=Inches(0.65)):
    """kpis = list of (value, label)"""
    n = len(kpis)
    cell_w = (W - Inches(1.0)) / n
    x0 = Inches(0.5)
    for i, (val, lbl) in enumerate(kpis):
        cx = x0 + i * cell_w
        r = slide.shapes.add_shape(1, cx, y, cell_w - Pt(3), h)
        r.fill.solid(); r.fill.fore_color.rgb = NAVY_M
        r.line.color.rgb = RGBColor(0x2a, 0x45, 0x70)
        r.line.width = Pt(0.5)
        txbox(slide, val, cx + Inches(0.1), y + Pt(2), cell_w - Inches(0.2), Inches(0.3),
              size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, lbl, cx + Inches(0.05), y + Inches(0.32), cell_w - Inches(0.1), Inches(0.28),
              size=7, color=GREY, align=PP_ALIGN.CENTER)

def info_list(slide, items, x, y, w, h_each=Inches(0.38)):
    """items = list of (bold_prefix, text)"""
    cur_y = y
    for bp, txt in items:
        r = slide.shapes.add_shape(1, x, cur_y, w, h_each - Pt(2))
        r.fill.solid(); r.fill.fore_color.rgb = NAVY_M
        r.line.color.rgb = RGBColor(0x2a, 0x45, 0x70)
        r.line.width = Pt(0.3)
        full = f'{bp}  {txt}' if bp else txt
        txbox(slide, full, x + Inches(0.15), cur_y + Pt(2), w - Inches(0.2), h_each - Pt(6),
              size=9.5, color=WHITE_DIM)
        cur_y += h_each

def table_hdr(slide, cols, x, y, w, row_h=Inches(0.3)):
    """cols = [(text, width_fraction)]"""
    total_w = w
    cx = x
    r_bg = slide.shapes.add_shape(1, x, y, total_w, row_h)
    r_bg.fill.solid(); r_bg.fill.fore_color.rgb = NAVY_L
    r_bg.line.fill.background()
    for (txt, frac) in cols:
        cw = Inches(frac * (w.inches))
        txbox(slide, txt.upper(), cx + Pt(4), y + Pt(3), cw - Pt(6), row_h - Pt(4),
              size=7.5, bold=True, color=GREY)
        cx += cw

def table_row(slide, vals, col_fracs, x, y, w, row_h=Inches(0.28), even=False):
    total_w = w
    fill = NAVY_M if even else RGBColor(0x0f, 0x22, 0x3a)
    r_bg = slide.shapes.add_shape(1, x, y, total_w, row_h)
    r_bg.fill.solid(); r_bg.fill.fore_color.rgb = fill
    r_bg.line.fill.background()
    cx = x
    for i, (val, frac) in enumerate(zip(vals, col_fracs)):
        cw = Inches(frac * (w.inches))
        txbox(slide, val, cx + Pt(4), y + Pt(2), cw - Pt(6), row_h - Pt(2),
              size=8.5, color=WHITE_DIM)
        cx += cw

def add_screenshot(slide, path, x, y, w, h, caption=''):
    """Ajoute une capture avec bordure et caption"""
    # cadre
    r = slide.shapes.add_shape(1, x - Pt(2), y - Pt(14), w + Pt(4), Pt(14))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY_L
    r.line.fill.background()
    cap_text = 'CAPTURE WAZUH / MONITORING' if not caption else caption[:60]
    txbox(slide, cap_text, x, y - Pt(12), w, Pt(12), size=6.5, color=GREY)
    # image
    r2 = slide.shapes.add_shape(1, x, y, w, h)
    r2.fill.solid(); r2.fill.fore_color.rgb = NAVY_M
    r2.line.color.rgb = RGBColor(0x2a, 0x45, 0x70); r2.line.width = Pt(0.5)
    add_image(slide, path, x, y, w, h)
    # caption bas
    if caption:
        txbox(slide, caption, x, y + h + Pt(1), w, Pt(12),
              size=6.5, color=GREY, italic=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 01 — COVER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide()
bg(s, BLACK)
# Gradient overlay top-left
r = s.shapes.add_shape(1, 0, 0, Inches(6), Inches(4))
r.fill.solid(); r.fill.fore_color.rgb = NAVY_M
r.line.fill.background()
from pptx.util import Pt as PT
logo_txt(s); slide_num(s, 1)
txbox(s, 'RAPPORT DE SÉCURITÉ — VERSION 2.0 — 6 AVRIL 2026',
      Inches(0.5), Inches(0.9), Inches(10), Inches(0.3),
      size=8, color=GREY, bold=True)
r2 = s.shapes.add_shape(1, Inches(0.5), Inches(1.25), Inches(0.8), Pt(2))
r2.fill.solid(); r2.fill.fore_color.rgb = WHITE; r2.line.fill.background()
txbox(s, 'Audit de Sécurité', Inches(0.5), Inches(1.35), Inches(9), Inches(0.8),
      size=38, bold=True, color=WHITE)
txbox(s, 'Plateforme GestionScolaire YNOV Campus',
      Inches(0.5), Inches(2.15), Inches(10), Inches(0.5),
      size=22, bold=False, color=WHITE_DIM)
txbox(s, 'Déploiement SIEM Wazuh 4.7.4  ·  Conformité RGPD  ·  WAF OWASP  ·  MITRE ATT&CK  ·  CVE  ·  FIM',
      Inches(0.5), Inches(2.7), Inches(12), Inches(0.35),
      size=10.5, color=GREY)
# Meta
metas = [('Auteur','Anass AKKER'),('Établissement','YNOV Campus'),
          ('Niveau','Master 2 Cybersécurité'),('Date','6 Avril 2026')]
for i,(lbl2,val) in enumerate(metas):
    mx = Inches(0.5 + i * 3.0)
    txbox(s, lbl2.upper(), mx, Inches(3.2), Inches(2.8), Inches(0.2),
          size=7, color=GREY, bold=True)
    txbox(s, val, mx, Inches(3.42), Inches(2.8), Inches(0.28),
          size=11, bold=True, color=WHITE)
kpi_bar(s, [('2 500+','Événements FIM'),('17','CVE Détectées'),
             ('1 600+','Alertes MITRE ATT&CK'),('100%','Couverture Agents'),
             ('5 / 6','Articles RGPD Validés')], y=Inches(3.9), h=Inches(0.75))
notes_zone(s, "Ce rapport de sécurité M2 présente l'audit complet de la plateforme GestionScolaire déployée à YNOV Campus. L'objectif est de démontrer la mise en œuvre d'une sécurité multicouche : SIEM Wazuh 4.7.4 (infrastructure), WAF applicatif (OWASP Top 10), RBAC (contrôle d'accès), et conformité RGPD. Les résultats clés : 2 500+ événements FIM, 17 CVE identifiées (8 High + 9 Medium), 1 600+ alertes MITRE ATT&CK, couverture agents 100%, score applicatif 100/100. La sécurité est abordée de manière Defense in Depth : système, réseau, applicatif, données.")
print("Slide 01 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 02 — SOMMAIRE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 2)
label(s, 'NAVIGATION'); title(s, 'Sommaire', size=24)
divider(s)
toc = [
    ('01','Tableau de Bord KPI — Indicateurs Clés'),
    ('02','Statut Global des 8 Modules de Sécurité'),
    ('03','Chronologie de l\'Audit — 6 Phases'),
    ('04','Architecture Wazuh 4.7.4 — Docker'),
    ('05','Règles Personnalisées YNOV-APP (100000–100005)'),
    ('06','8 Problèmes Techniques Résolus — Partie 1/2'),
    ('07','8 Problèmes Techniques Résolus — Partie 2/2'),
    ('08','Bilan 8/8 Problèmes Résolus'),
    ('09','Infrastructure Agents — Couverture 100%'),
    ('10','Dashboard Agents Overview'),
    ('11','Vue d\'Ensemble Agent 002 macOS 15.7.4'),
    ('12','FIM — File Integrity Monitoring (2 500+)'),
    ('13','CVE — 17 Vulnérabilités Identifiées'),
    ('14','MITRE ATT&CK — 1 600+ Alertes'),
    ('15','Conformité RGPD × Wazuh (5/6 articles)'),
    ('16','Dashboard Accès Sécurisé — Score 100/100'),
    ('17','Dashboard RGPD Conformité Live'),
    ('18','RBAC & Audit Logs Firestore'),
    ('19','WAF OWASP Top 10 Actif'),
    ('20','SIEM Logs Temps Réel'),
    ('21','Plan d\'Action & Recommandations'),
    ('22','Bilan Final'),
    ('23','Conclusion'),
]
cols = 2; rows_per_col = 12
for i, (n, t) in enumerate(toc):
    col = i // rows_per_col
    row = i % rows_per_col
    cx = Inches(0.5 + col * 6.4)
    cy = Inches(2.05 + row * 0.34)
    rw = Inches(6.1); rh = Inches(0.31)
    rb = s.shapes.add_shape(1, cx, cy, rw, rh)
    rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.3)
    txbox(s, n, cx + Inches(0.1), cy + Pt(2), Inches(0.4), rh - Pt(4), size=8, color=GREY, bold=True)
    txbox(s, t, cx + Inches(0.55), cy + Pt(2), rw - Inches(0.65), rh - Pt(4), size=8.5, color=WHITE_DIM)
notes_zone(s, "Ce rapport suit une progression logique : résultats globaux → méthodologie → incidents techniques résolus → analyse modules Wazuh → conformité RGPD → captures monitoring applicatif → plan d'action. Chaque section est illustrée par des captures réelles du système. La structure reflète une démarche d'audit professionnelle : mesurer, analyser, documenter, corriger.")
print("Slide 02 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 03 — KPI
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 3)
label(s, 'SECTION 1 — RÉSULTATS')
title(s, 'Tableau de Bord KPI — Indicateurs Clés de Sécurité', size=22)
divider(s)
kpis5 = [
    ('2 500+', 'Événements FIM\nRule 550 Lvl 7 · root 99.44%'),
    ('17', 'CVE Identifiées\n8 High + 9 Medium'),
    ('1 600+', 'Alertes MITRE ATT&CK\nImpact T1499 ×1594'),
    ('100%', 'Couverture Agents\n2 agents actifs'),
    ('0 / 10', 'SCA Passed\nunix_audit — surface max'),
]
x0 = Inches(0.4); kw = Inches(2.45); ky = Inches(2.05); kh = Inches(1.1)
for i,(val,lbl) in enumerate(kpis5):
    cx = x0 + i*(kw + Inches(0.12))
    rb = s.shapes.add_shape(1, cx, ky, kw, kh)
    rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
    txbox(s, val, cx, ky+Pt(4), kw, Inches(0.45),
          size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    lines = lbl.split('\n')
    txbox(s, lines[0], cx+Inches(0.05), ky+Inches(0.5), kw-Inches(0.1), Inches(0.22),
          size=7.5, color=GREY, align=PP_ALIGN.CENTER)
    if len(lines)>1:
        txbox(s, lines[1], cx+Inches(0.05), ky+Inches(0.72), kw-Inches(0.1), Inches(0.3),
              size=7, color=GREY, align=PP_ALIGN.CENTER, italic=True)
info_list(s, [
    ('Stack technique :', 'Firebase Firestore · React 18 · Express.js · Node.js · Docker 26.1 · Wazuh 4.7.4'),
    ('Authentification :', 'JWT HS256 · bcrypt rounds=12 · RBAC 4 rôles (admin, sous-admin, comptable, étudiant/parent)'),
    ('Conformité RGPD :', '5/6 articles validés — Art. 5, 17, 25, 33, 35 ✓ — Art. 32 (chiffrement E2E) en cours'),
    ('Score applicatif :', 'Dashboard monitoring interne : 100/100 — WAF + JWT + RBAC + AuditLogs — 0 violations détectées'),
], x=Inches(0.4), y=Inches(3.3), w=Inches(12.5))
notes_zone(s, "Les 5 KPI reflètent une surveillance complète. 2 500+ événements FIM signalent des modifications massives de binaires système (surtout root), révélant un risque de modification malveillante. Les 17 CVE dont 8 critiques nécessitent un patch urgent. 1 600+ alertes MITRE ATT&CK montrent la détection en temps réel des tactiques d'attaque. Le score SCA 0/10 confirme que le hardening système est la priorité n°1. Le score applicatif 100/100 valide la sécurité de la couche web (WAF, JWT, RBAC, AuditLogs).")
print("Slide 03 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 04 — STATUT MODULES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 4)
label(s, 'SECTION 1 — SYNTHÈSE')
title(s, 'Statut Global des Modules de Sécurité', size=22)
divider(s)
cols4 = [('Module',3.5),('Résultat',1.8),('Niveau',1.8),('Statut',1.0),('Observation',4.2)]
fracs4 = [c[1]/12.3 for c in cols4]
tw = Inches(12.3); tx = Inches(0.5); ty = Inches(2.05)
table_hdr(s, cols4, tx, ty, tw)
rows4 = [
    ('File Integrity Monitoring (FIM)','2 500+ événements','Rule 550 Lvl 7','✓ ACTIF','/bin/bash · /bin/cp · /bin/df modifiés (root)'),
    ('CVE / Vulnerability Assessment','17 CVE','8 High · 9 Medium','✓ ACTIF','CVE-2019-5736 CVSS3=8.6 (docker critique)'),
    ('MITRE ATT&CK','1 600+ alertes','Impact + Defense Evasion','✓ ACTIF','T1499 ×1594 · Defense Evasion ×6'),
    ('SCA — Security Config Assessment','0/10 Passed','unix_audit CIS','⚠ CRITIQUE','Surface d\'attaque maximale détectée'),
    ('WAF Applicatif (waf.js)','4 catégories','OWASP Top 10','✓ ACTIF','SQLi · XSS · Path Traversal · CMD Inject'),
    ('Conformité RGPD','5/6 articles','Art. 32 en cours','~ PARTIEL','Art. 5, 17, 25, 33, 35 validés'),
    ('Audit Logs Firestore','9 types d\'événements','Immuables','✓ ACTIF','Traçabilité complète admin/étudiant'),
    ('RBAC','4 rôles définis','JWT HS256 + bcrypt','✓ ACTIF','admin · sous-admin · comptable · étudiant'),
]
for i, row in enumerate(rows4):
    table_row(s, row, fracs4, tx, ty + Inches(0.3)*(i+1), tw, even=(i%2==0))
notes_zone(s, "Ce tableau synthétise l'état de tous les modules déployés. 6 modules sont pleinement actifs (FIM, CVE, MITRE, WAF, AuditLogs, RBAC). 1 module est critique (SCA à 0/10 — aucune bonne pratique de configuration système appliquée). 1 module est partiel (RGPD — Art. 32 sur le chiffrement at-rest en attente). La colonne 'Observation' détaille les éléments spécifiques détectés, démontrant une analyse approfondie et non une simple vérification de présence des outils.")
print("Slide 04 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 05 — CHRONOLOGIE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 5)
label(s, 'SECTION 1 — DÉMARCHE')
title(s, 'Chronologie de l\'Audit de Sécurité — 6 Phases', size=22)
divider(s)
phases = [
    ('Phase 1','Développement sécurisé','Architecture React/Firebase/Node.js · JWT HS256 · bcrypt rounds=12 · RBAC 4 rôles · AuditLogs Firestore immuables · WAF middleware'),
    ('Phase 2','Déploiement Wazuh 4.7.4','Installation Docker Compose (Manager + Indexer + Dashboard) · Ports 1514/55000/443 · Enrôlement agents OSSEC'),
    ('Phase 3','Configuration modules','FIM ossec.conf · NVD Feed CVE · SCA unix_audit · MITRE ATT&CK rules · Règles YNOV-APP (100000–100005)'),
    ('Phase 4','Résolution 8 problèmes','Agent déconnecté · FIM vide · NVD 403 · bcrypt DLL · Firebase timeout · field action · SIP macOS · PPTX CRC-32'),
    ('Phase 5','Collecte métriques','2500+ FIM · 17 CVE · 1600+ MITRE · SCA 0/10 · RGPD 5/6 · Dashboard applicatif Score 100/100'),
    ('Phase 6','Rapport & Recommandations','Rédaction rapport v2.0 · Plan d\'action 6 priorités · Hardening SCA · Patch CVE · Art. 32 RGPD'),
]
for i,(ph,name,desc) in enumerate(phases):
    cy = Inches(2.05 + i*0.54)
    rb = s.shapes.add_shape(1, Inches(0.4), cy, Inches(12.5), Inches(0.50))
    rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.3)
    txbox(s, ph, Inches(0.55), cy+Pt(2), Inches(0.85), Inches(0.25), size=7.5, bold=True, color=GREY)
    txbox(s, name, Inches(1.45), cy+Pt(2), Inches(2.2), Inches(0.25), size=9, bold=True, color=WHITE)
    txbox(s, desc, Inches(3.75), cy+Pt(2), Inches(9.0), Inches(0.35), size=8.5, color=WHITE_DIM)
notes_zone(s, "La démarche est rigoureuse et professionnelle : Security by Design dès la conception, installation SIEM progressive, configuration module par module, diagnostic et résolution des problèmes techniques en environnement réel. La phase 4 (8 problèmes) démontre la capacité de diagnostic multicouche. La phase 5 produit des métriques mesurables et vérifiables par captures d'écran. L'ensemble constitue un audit complet conforme aux standards professionnels (ISO 27001, NIST CSF).")
print("Slide 05 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 06 — ARCHITECTURE WAZUH
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 6)
label(s, 'SECTION 3 — ARCHITECTURE')
title(s, 'Architecture Wazuh 4.7.4 — Déploiement Docker', size=22)
divider(s)
archi = [
    ('Wazuh Manager','Conteneur Docker · Port 1514 (agent) · Port 55000 (API REST) · Collecte & corrélation · Moteur règles YNOV-APP'),
    ('Wazuh Indexer (OpenSearch)','Moteur d\'indexation embarqué · Stockage alertes · Port 9200 · Full-text search logs Firestore/Docker/macOS'),
    ('Wazuh Dashboard','Interface Kibana fork · Port 443 HTTPS TLS · Visualisation FIM, CVE, MITRE, SCA · Accès authentifié'),
    ('Agent 001 — Docker Host','Surveillance conteneurs · Logs Docker Engine · FIM /var/lib/docker · Monitoring réseau conteneurs'),
    ('Agent 002 — macOS 15.7.4','Machine développement · FIM 2500+ events · CVE 17 via NVD · MITRE T1499 ×1594 · SCA 0/10 unix_audit'),
    ('Application GestionScolaire','React 18 + Firebase Firestore · Express.js/Node.js · WAF waf.js OWASP · JWT HS256 + bcrypt · AuditLogs · RBAC 4 rôles'),
]
for i,(comp,desc) in enumerate(archi):
    col = i//3; row = i%3
    cx = Inches(0.4 + col*6.55); cy = Inches(2.05 + row*1.05)
    rb = s.shapes.add_shape(1, cx, cy, Inches(6.25), Inches(0.95))
    rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
    txbox(s, comp, cx+Inches(0.15), cy+Pt(4), Inches(5.9), Inches(0.25),
          size=9.5, bold=True, color=WHITE)
    txbox(s, desc, cx+Inches(0.15), cy+Inches(0.32), Inches(5.9), Inches(0.55),
          size=8, color=WHITE_DIM)
notes_zone(s, "L'architecture Wazuh repose sur 3 composants Docker Compose : Manager (cerveau, collecte et corrèle), Indexer (OpenSearch, stockage et recherche), Dashboard (visualisation). 2 agents OSSEC enrôlés : Agent 001 surveille l'hôte Docker et ses conteneurs, Agent 002 surveille la machine macOS. Couverture complète : conteneurs, OS, application web. Tous les composants communiquent en TLS. L'application GestionScolaire génère des logs qui remontent via l'agent vers Wazuh pour corrélation avec les règles YNOV-APP personnalisées.")
print("Slide 06 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 07 — RÈGLES YNOV-APP
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 7)
label(s, 'SECTION 3 — RÈGLES PERSONNALISÉES')
title(s, 'Règles YNOV-APP — Détection Applicative Spécifique', size=22)
divider(s)
cols7 = [('Règle ID',1.2),('Niveau',1.5),('Description',2.8),('Déclencheur Pattern',4.0),('Catégorie',2.0)]
fracs7 = [c[1]/11.5 for c in cols7]
tw7 = Inches(11.5); tx7 = Inches(0.5); ty7 = Inches(2.05)
table_hdr(s, cols7, tx7, ty7, tw7)
rows7 = [
    ('100000','10 — Critique','Injection SQL','SELECT / UNION / DROP dans requête HTTP','WAF / SQLi'),
    ('100001','9 — Élevé','Attaque XSS','<script> / onerror / javascript:','WAF / XSS'),
    ('100002','9 — Élevé','Path traversal / LFI','../../ / %2e%2e dans URL','WAF / LFI'),
    ('100003','8 — Haut','Injection commande OS','; && | $( cmd ) dans paramètre','WAF / CMDi'),
    ('100004','7 — Moyen','Brute force auth','5+ échecs JWT en 60 secondes','Auth / Brute Force'),
    ('100005','6 — Moyen','Violation RBAC','403 sur endpoint privilégié','RBAC / Privilege'),
]
for i,row in enumerate(rows7):
    table_row(s, row, fracs7, tx7, ty7+Inches(0.3)*(i+1), tw7, even=(i%2==0))
txbox(s,'Déployées dans /var/ossec/etc/rules/local_rules.xml — rechargement dynamique sans redémarrage Manager · Format XML Wazuh natif',
      Inches(0.5), Inches(4.4), Inches(12), Inches(0.3), size=8.5, color=GREY, italic=True)
notes_zone(s, "Les règles YNOV-APP (IDs 100000 à 100005) sont spécifiquement créées pour l'application GestionScolaire. Contrairement aux règles génériques Wazuh, elles sont calibrées sur les patterns d'attaque des applications web Node.js/Firebase. Chaque règle peut déclencher une réponse active (blocage IP, notification admin). Les niveaux (6 à 10) permettent une priorisation automatique des alertes. Ces règles constituent la valeur ajoutée principale de l'intégration SIEM avec l'application : elles transforment Wazuh en outil de détection spécifique au contexte métier.")
print("Slide 07 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 08 — PROBLÈMES 1-4
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 8)
label(s, 'SECTION 2 — PROBLÈMES TECHNIQUES RÉSOLUS')
title(s, '8 Problèmes Techniques Résolus — Partie 1 / 2', size=22)
divider(s)
probs1 = [
    ('01','Agent Wazuh "Never Connected"',
     "Symptôme : Agent 002 affiché 'Never connected' malgré installation correcte.\nCause : Pare-feu macOS bloquait port 1514 UDP sortant vers Manager Docker.\nFix : Règle firewall macOS + ouverture port 1514 TCP/UDP + rechargement ossec-authd."),
    ('02','FIM — Aucun Résultat Dashboard',
     "Symptôme : Dashboard FIM vide malgré activation dans ossec.conf.\nCause : Chemins macOS non configurés dans syscheck <directories>.\nFix : Ajout <directories>/etc,/bin,/usr</directories> + <frequency>300</frequency> + restart agent."),
    ('03','NVD Feed CVE — Erreur 403 Forbidden',
     "Symptôme : Téléchargement CVE feed NVD NIST échoue avec HTTP 403.\nCause : API NVD exige désormais une clé API pour l'accès programmatique (depuis 2023).\nFix : Génération clé API NVD NIST + injection dans vulnerability-detector config Manager."),
    ('04','bcrypt — DLL Manquante (VCRUNTIME140)',
     "Symptôme : Node.js crash au démarrage 'VCRUNTIME140.dll not found'.\nCause : Module node-gyp bcrypt requiert Visual C++ Redistributables non installés.\nFix : Installation Visual C++ 2019 Redistributable x64 + npm rebuild bcrypt."),
]
for i,(n,name,desc) in enumerate(probs1):
    col = i%2; row = i//2
    cx = Inches(0.4 + col*6.5); cy = Inches(2.05 + row*1.55)
    rb = s.shapes.add_shape(1, cx, cy, Inches(6.2), Inches(1.45))
    rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
    txbox(s, f'Problème {n}', cx+Inches(0.15), cy+Pt(4), Inches(2), Inches(0.22),
          size=7.5, bold=True, color=GREY)
    txbox(s, name, cx+Inches(0.15), cy+Inches(0.28), Inches(5.85), Inches(0.25),
          size=9.5, bold=True, color=WHITE)
    txbox(s, desc, cx+Inches(0.15), cy+Inches(0.55), Inches(5.85), Inches(0.82),
          size=8, color=WHITE_DIM)
notes_zone(s, "Ces 4 premiers problèmes touchent l'infrastructure SIEM et l'environnement de déploiement. Problème 01 : classique sur macOS, le pare-feu natif bloque silencieusement le trafic UDP. Problème 02 : configuration spécifique à l'OS pour FIM. Problème 03 : changement de politique NVD NIST en 2023 (API authentifiée). Problème 04 : dépendances natives Node.js sur Windows. Chaque diagnostic a suivi la méthode : observer → analyser les logs → identifier la cause racine → corriger → valider.")
print("Slide 08 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 09 — PROBLÈMES 5-8
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 9)
label(s, 'SECTION 2 — PROBLÈMES TECHNIQUES RÉSOLUS')
title(s, '8 Problèmes Techniques Résolus — Partie 2 / 2', size=22)
divider(s)
probs2 = [
    ('05','Firebase Firestore — Timeout DEADLINE_EXCEEDED',
     "Symptôme : Requêtes Firestore échouent avec DEADLINE_EXCEEDED après 10s.\nCause : SDK Admin Firebase mal configuré (project ID / credentials service account absents).\nFix : Validation GOOGLE_APPLICATION_CREDENTIALS + timeout 30s + retry exponential backoff."),
    ('06','Champ Firestore \"action\" Réservé',
     "Symptôme : Enregistrement auditLogs Firestore échoue sur le champ 'action'.\nCause : Firestore rejette 'action' (conflit avec Security Rules internes non documenté).\nFix : Renommage 'action' → 'eventType' dans modèle AuditLog + migration documents existants."),
    ('07','SIP macOS — Accès /etc Refusé (EPERM)',
     "Symptôme : Agent Wazuh génère EPERM en lisant /etc/sudoers pour FIM.\nCause : System Integrity Protection (SIP) macOS bloque l'accès aux chemins système protégés.\nFix : Exclusion chemins SIP + surveillance /usr/local, /opt, ~/Documents (chemins accessibles)."),
    ('08','PPTX Corrompu — Erreur CRC-32',
     "Symptôme : Fichier PowerPoint généré corrompu à l'ouverture (CRC-32 mismatch).\nCause : Bibliothèque python-pptx (ancienne version) utilise compression ZIP incorrecte pour médias.\nFix : Mise à jour python-pptx 1.0.2 + génération structurée avec images correctement compressées."),
]
for i,(n,name,desc) in enumerate(probs2):
    col = i%2; row = i//2
    cx = Inches(0.4 + col*6.5); cy = Inches(2.05 + row*1.55)
    rb = s.shapes.add_shape(1, cx, cy, Inches(6.2), Inches(1.45))
    rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
    txbox(s, f'Problème {n}', cx+Inches(0.15), cy+Pt(4), Inches(2), Inches(0.22),
          size=7.5, bold=True, color=GREY)
    txbox(s, name, cx+Inches(0.15), cy+Inches(0.28), Inches(5.85), Inches(0.25),
          size=9.5, bold=True, color=WHITE)
    txbox(s, desc, cx+Inches(0.15), cy+Inches(0.55), Inches(5.85), Inches(0.82),
          size=8, color=WHITE_DIM)
notes_zone(s, "Les 4 problèmes suivants touchent des couches plus applicatives. Problème 05 : Firebase timeout fréquent quand les credentials ne sont pas correctement configurés dans les variables d'environnement. Problème 06 : contrainte non documentée de Firestore sur les noms de champs réservés. Problème 07 : intrinsèque à macOS Sequoia avec SIP — ne pas désactiver SIP, contourner proprement. Problème 08 : résolu avec python-pptx 1.0.2 qui gère correctement la compression des médias dans les archives ZIP OOXML.")
print("Slide 09 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 — BILAN 8 PROBLÈMES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 10)
label(s, 'SECTION 2 — BILAN')
title(s, 'Bilan Technique — 8 / 8 Problèmes Résolus avec Succès', size=22)
divider(s)
cols10 = [('#',0.5),('Problème',2.8),('Domaine',2.0),('Impact',3.2),('Statut',1.2)]
fracs10 = [c[1]/9.7 for c in cols10]
tw10 = Inches(12.3); tx10 = Inches(0.5)
table_hdr(s, cols10, tx10, Inches(2.05), tw10)
rows10 = [
    ('01','Agent Wazuh "Never Connected"','Infrastructure SIEM','Blocage complet du monitoring','✓ RÉSOLU'),
    ('02','FIM sans résultats dashboard','SIEM / FIM','Absence de détection filesystem','✓ RÉSOLU'),
    ('03','NVD Feed HTTP 403','CVE / Vulnerability','Pas de scan CVE automatique','✓ RÉSOLU'),
    ('04','bcrypt DLL manquante','Applicatif / Auth','Authentification non fonctionnelle','✓ RÉSOLU'),
    ('05','Firebase Firestore Timeout','Base de données','Perte de données auditLogs','✓ RÉSOLU'),
    ('06','Champ "action" réservé Firestore','Schema données','Échec enregistrement événements','✓ RÉSOLU'),
    ('07','SIP macOS /etc refusé (EPERM)','OS / Permissions','FIM incomplet sur chemins système','✓ RÉSOLU'),
    ('08','PPTX corrompu CRC-32','Outillage / Livrable','Présentation non ouvrable','✓ RÉSOLU'),
]
for i,row in enumerate(rows10):
    table_row(s, row, fracs10, tx10, Inches(2.35 + i*0.32), tw10, Inches(0.29), even=(i%2==0))
notes_zone(s, "La résolution de 100% des problèmes (8/8) démontre une expertise en débogage multicouches : infrastructure Docker, OS macOS, base de données Firebase, sécurité applicative Node.js. Chaque résolution a suivi la même démarche : isoler la cause racine (pas traiter le symptôme), documenter la solution, valider la correction. Aucun contournement non sécurisé n'a été utilisé (pas de désactivation SIP, pas de --no-verify). Cette rigueur de diagnostic est une compétence essentielle pour un ingénieur cybersécurité travaillant sur des environnements complexes et hétérogènes.")
print("Slide 10 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11 — AGENTS COUVERTURE + CAPTURE agents
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 11)
label(s, 'SECTION 3 — INFRASTRUCTURE')
title(s, 'Infrastructure Agents — Couverture 100%', size=22)
divider(s)
info_list(s, [
    ('Agent 001 :','Docker Host · Surveillance conteneurs · Logs Docker Engine · FIM /var/lib/docker'),
    ('Agent 002 :','macOS 15.7.4 · FIM 2500+ événements · CVE 17 via NVD · MITRE T1499 ×1594'),
    ('Protocole :','OSSEC agent · Port 1514 TCP/UDP · Authentification clé pré-partagée · Chiffrement TLS'),
    ('Statut :','2/2 agents actifs · Synchronisation temps réel · Couverture 100% de l\'infrastructure'),
], x=Inches(0.4), y=Inches(2.05), w=Inches(5.5), h_each=Inches(0.40))
add_screenshot(s, WAZUH['agents'], Inches(6.1), Inches(2.0), Inches(6.9), Inches(3.2),
               'WAZUH — Liste & Statut Agents (Agent 001 + Agent 002 · Status: Active)')
notes_zone(s, "Cette capture montre la liste des agents Wazuh enrôlés dans le Manager. Les 2 agents sont actifs et en communication continue. L'Agent 001 couvre l'infrastructure Docker, l'Agent 002 couvre la machine de développement macOS. La couverture 100% signifie qu'aucune partie de l'infrastructure n'échappe à la surveillance SIEM. Le protocole OSSEC-agent utilise une clé pré-partagée pour l'authentification, et TLS pour le chiffrement des communications agent → Manager. Les événements remontent en temps réel pour corrélation et indexation.")
print("Slide 11 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12 — DASHBOARD AGENTS OVERVIEW
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 12)
label(s, 'SECTION 3 — DASHBOARD AGENTS')
title(s, 'Dashboard Agents Overview — Tous Modules Actifs', size=22)
divider(s)
add_screenshot(s, WAZUH['agents01'], Inches(0.4), Inches(2.0), Inches(7.8), Inches(3.3),
               'WAZUH Dashboard — Agents Overview · 2 agents actifs · Modules: FIM, CVE, MITRE, SCA, Logcollector')
info_list(s, [
    ('FIM :','Activé sur les 2 agents · Surveillance filesystem temps réel'),
    ('CVE Scan :','Actif · Feed NVD NIST · Scan automatique'),
    ('MITRE ATT&CK :','Corrélation automatique tactiques/techniques'),
    ('SCA :','Audit configuration unix_audit actif'),
    ('Logcollector :','Collecte logs Express.js / Docker Engine'),
], x=Inches(8.4), y=Inches(2.05), w=Inches(4.6), h_each=Inches(0.40))
notes_zone(s, "Le dashboard Wazuh affiche en temps réel l'état de tous les agents et les modules actifs. On voit quels modules sont déployés (FIM, Vulnerability, MITRE, SCA, Logcollector), le nombre d'événements par agent, et les dernières alertes. Ce dashboard est l'interface centrale pour un opérateur SOC (Security Operations Center). Il permet de détecter immédiatement si un agent est déconnecté, si un module ne génère plus d'événements, ou si une anomalie survient. La richesse des informations affichées valide la qualité du déploiement Wazuh dans son ensemble.")
print("Slide 12 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13 — AGENT 002 VUE D'ENSEMBLE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 13)
label(s, 'SECTION 4 — AGENT 002 MACOS')
title(s, 'Vue d\'Ensemble Agent 002 — macOS 15.7.4', size=22)
divider(s)
add_screenshot(s, WAZUH['overview'], Inches(0.4), Inches(2.0), Inches(7.8), Inches(3.3),
               'AGENT 002 — Vue d\'ensemble complète · FIM + SCA + MITRE ATT&CK + Vulnerabilities actifs')
info_list(s, [
    ('OS :','macOS Sequoia 15.7.4 · Machine de développement principale'),
    ('FIM :','2 500+ événements · Rule 550 Niveau 7 · root 99.44%'),
    ('CVE :','17 vulnérabilités détectées · 8 High + 9 Medium'),
    ('MITRE :','1 600+ alertes · Impact T1499 ×1594 dominant'),
    ('SCA :','0/10 checks passed · unix_audit · surface d\'attaque max'),
], x=Inches(8.4), y=Inches(2.05), w=Inches(4.6), h_each=Inches(0.40))
notes_zone(s, "L'Agent 002 installé sur macOS 15.7.4 est l'agent le plus actif. Il concentre la majorité des événements car il surveille directement le système de fichiers de la machine de développement (activité intense : mises à jour OS, installations npm, modifications de configuration). La vue d'ensemble Wazuh montre simultanément tous les modules et leurs compteurs. C'est cette vue qui permet de constater d'un coup d'œil la santé globale de l'agent. Le volume d'événements (2500+ FIM, 1600+ MITRE, 17 CVE) démontre l'efficacité de la surveillance SIEM en environnement réel de développement.")
print("Slide 13 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14 — FIM
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 14)
label(s, 'SECTION 4 — FIM — FILE INTEGRITY MONITORING')
title(s, 'File Integrity Monitoring — 2 500+ Événements Détectés', size=22)
divider(s)
# KPI card gauche
rb = s.shapes.add_shape(1, Inches(0.4), Inches(2.05), Inches(5.5), Inches(0.85))
rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
txbox(s,'2 500+',Inches(0.5),Inches(2.1),Inches(5.3),Inches(0.5),size=28,bold=True,color=WHITE)
txbox(s,'ÉVÉNEMENTS FIM · Rule 550 Niveau 7',Inches(0.5),Inches(2.6),Inches(5.3),Inches(0.25),size=8,color=GREY)
cols14 = [('Indicateur',2.8),('Valeur',2.6)]
fracs14 = [c[1]/5.4 for c in cols14]
table_hdr(s, cols14, Inches(0.4), Inches(3.05), Inches(5.4))
rows14 = [
    ('Règle déclenchée','Rule 550 — Niveau 7'),
    ('Utilisateur dominant','root (99.44%)'),
    ('Type d\'événement','Modified (100%)'),
    ('Fichiers modifiés','/bin/bash · /bin/cp · /bin/df · /bin/echo'),
    ('Chemins surveillés','/bin · /usr · /etc (hors SIP)'),
]
for i,row in enumerate(rows14):
    table_row(s, row, fracs14, Inches(0.4), Inches(3.35)+Inches(0.3)*i, Inches(5.4), even=(i%2==0))
add_screenshot(s, WAZUH['fim'], Inches(6.1), Inches(2.0), Inches(7.0), Inches(3.3),
               'FIM DASHBOARD — 2 500+ événements · root 99.44% · Modified 100% · /bin/bash /bin/cp visibles')
notes_zone(s, "Le FIM surveille les modifications apportées aux fichiers système critiques. 2 500+ événements détectent des modifications massives des binaires macOS /bin (bash, cp, df, echo) par root. La Rule 550 Niveau 7 est la règle Wazuh standard pour les modifications de fichiers surveillés. Root à 99.44% est cohérent avec les mises à jour système macOS (binaires /bin mis à jour par root via mises à jour OS). Ce pattern doit rester surveillé car il peut aussi masquer une compromission de compte root ou un rootkit. Le FIM est essentiel pour détecter les modifications malveillantes de binaires système (technique T1548 MITRE).")
print("Slide 14 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 15 — CVE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 15)
label(s, 'SECTION 4 — CVE — VULNERABILITY ASSESSMENT')
title(s, 'CVE — 17 Vulnérabilités Identifiées (8 High + 9 Medium)', size=22)
divider(s)
add_screenshot(s, WAZUH['cve'], Inches(0.4), Inches(2.0), Inches(7.0), Inches(3.3),
               'CVE DASHBOARD — 17 CVE · 8 High + 9 Medium · Docker dominant · NVD NIST Feed')
# KPI
rb = s.shapes.add_shape(1, Inches(7.6), Inches(2.05), Inches(5.4), Inches(0.75))
rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
txbox(s,'17  CVE — 8 High + 9 Medium',Inches(7.7),Inches(2.1),Inches(5.2),Inches(0.55),size=18,bold=True,color=WHITE)
cols15 = [('CVE ID',2.0),('CVSS3',0.8),('Composant',1.5),('Sévérité',1.1)]
fracs15 = [c[1]/5.4 for c in cols15]
table_hdr(s, cols15, Inches(7.6), Inches(2.95), Inches(5.4))
rows15 = [
    ('CVE-2019-5736','8.6','runc / Docker','CRITIQUE'),
    ('CVE Docker Engine','7.x','Docker Engine','HIGH'),
    ('CVE Node.js deps','6.x–7.x','npm packages','Medium'),
    ('+14 CVE supplémentaires','—','Stack complète','High/Medium'),
]
for i,row in enumerate(rows15):
    table_row(s, row, fracs15, Inches(7.6), Inches(3.25)+Inches(0.3)*i, Inches(5.4), even=(i%2==0))
notes_zone(s, "Le module Vulnerability Detector compare les packages installés avec la base NVD NIST pour identifier les CVE. Les 17 CVE (8H+9M) constituent un risque réel. La CVE-2019-5736 (CVSS3=8.6) est la plus critique : elle affecte runc (runtime Docker) et permet à un attaquant depuis l'intérieur d'un conteneur d'écrire dans le binaire runc de l'hôte, permettant une évasion de conteneur et élévation de privilèges vers root sur l'hôte. Patch prioritaire : mise à jour Docker Engine vers ≥26.1.5 incluant runc ≥1.0.0-rc8. Les 16 autres CVE nécessitent un plan npm audit fix et mises à jour dépendances.")
print("Slide 15 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 16 — MITRE ATT&CK
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 16)
label(s, 'SECTION 4 — MITRE ATT&CK')
title(s, 'MITRE ATT&CK — 1 600+ Alertes Corrélées', size=22)
divider(s)
# KPI gauche
rb = s.shapes.add_shape(1, Inches(0.4), Inches(2.05), Inches(5.5), Inches(0.75))
rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
txbox(s,'1 600+  Alertes MITRE ATT&CK',Inches(0.5),Inches(2.1),Inches(5.3),Inches(0.55),size=18,bold=True,color=WHITE)
cols16 = [('Tactique',2.0),('Technique',1.5),('Occurrences',1.9)]
fracs16 = [c[1]/5.4 for c in cols16]
table_hdr(s, cols16, Inches(0.4), Inches(2.95), Inches(5.4))
rows16 = [
    ('Impact','T1499 — Endpoint DoS','×1594 (dominant)'),
    ('Defense Evasion','Multiples sous-techniques','×6'),
    ('Autres tactiques','—','Détectées'),
]
for i,row in enumerate(rows16):
    table_row(s, row, fracs16, Inches(0.4), Inches(3.25)+Inches(0.3)*i, Inches(5.4), even=(i%2==0))
info_list(s, [
    ('T1499 :','Endpoint DoS : modifications intensives binaires système macOS par root corrélées à cette technique ATT&CK (même pattern que FIM 2500+)'),
    ('Defense Evasion :','6 événements de tentatives de dissimulation d\'activité (modification logs, changement permissions fichiers)'),
], x=Inches(0.4), y=Inches(4.15), w=Inches(5.5))
add_screenshot(s, WAZUH['mitre'], Inches(6.1), Inches(2.0), Inches(7.0), Inches(3.3),
               'MITRE ATT&CK DASHBOARD — Impact T1499 ×1594 · Defense Evasion ×6 · Matrice ATT&CK Wazuh')
notes_zone(s, "MITRE ATT&CK est le framework universel de classification des comportements malveillants utilisé par les équipes SOC mondiales. Wazuh mappe automatiquement ses règles aux techniques ATT&CK. Les 1 594 événements T1499 (Endpoint DoS / Resource Exhaustion) correspondent aux modifications massives de fichiers binaires détectées par FIM — Wazuh interprète ce pattern comme une technique d'Impact potentielle. Les 6 événements Defense Evasion signalent des tentatives de dissimulation. Cette corrélation MITRE transforme des alertes brutes en intelligence de menace (Threat Intelligence) actionnable pour les équipes de réponse à incident (DFIR).")
print("Slide 16 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 17 — RGPD
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 17)
label(s, 'SECTION 5 — CONFORMITÉ RGPD')
title(s, 'Conformité RGPD × Wazuh — 5 / 6 Articles Validés', size=22)
divider(s)
rgpd_items = [
    ('Art. 5','Licéité du traitement','✓ VALIDÉ',
     'AuditLogs Firestore immuables · 9 types d\'événements tracés · Traçabilité complète des accès données personnelles'),
    ('Art. 17','Droit à l\'effacement','✓ VALIDÉ',
     'Endpoint DELETE /api/rgpd/erase · Anonymisation données étudiant/parent · Log d\'effacement conservé'),
    ('Art. 25','Privacy by Design','✓ VALIDÉ',
     'RBAC granulaire 4 rôles · Principe du moindre privilège · Séparation contextes admin/étudiant/parent'),
    ('Art. 33','Notification violation','✓ VALIDÉ',
     'Wazuh active response · Alertes e-mail admin sur événement critique · Délai 72h conforme RGPD'),
    ('Art. 35','AIPD — Analyse Impact','✓ VALIDÉ',
     'Documentation risques · Évaluation impact traitement données élèves · Mesures mitigatives documentées'),
    ('Art. 32','Sécurité du traitement','⚠ EN COURS',
     'JWT + bcrypt OK · TLS transit OK · Chiffrement at-rest Firestore partiel · E2E encryption absent'),
]
for i,(art,name,st,desc) in enumerate(rgpd_items):
    col = i//3; row = i%3
    cx = Inches(0.4 + col*6.5); cy = Inches(2.05 + row*1.08)
    rb = s.shapes.add_shape(1, cx, cy, Inches(6.2), Inches(0.98))
    rb.fill.solid()
    rb.fill.fore_color.rgb = NAVY_M if '✓' in st else RGBColor(0x10,0x20,0x38)
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70) if '✓' in st else RGBColor(0x35,0x50,0x70)
    rb.line.width = Pt(0.5)
    txbox(s, art, cx+Inches(0.12), cy+Pt(3), Inches(0.8), Inches(0.22), size=7.5, bold=True, color=GREY)
    txbox(s, name, cx+Inches(0.95), cy+Pt(3), Inches(2.5), Inches(0.22), size=9, bold=True, color=WHITE)
    txbox(s, st, cx+Inches(4.0), cy+Pt(3), Inches(2.1), Inches(0.22), size=8, bold=True,
          color=WHITE if '✓' in st else GREY, align=PP_ALIGN.RIGHT)
    txbox(s, desc, cx+Inches(0.12), cy+Inches(0.35), Inches(5.95), Inches(0.55), size=8, color=WHITE_DIM)
notes_zone(s, "La conformité RGPD a été évaluée article par article pour les 6 articles les plus pertinents dans une application de gestion scolaire traitant des données personnelles d'étudiants. 5 articles sont pleinement satisfaits grâce aux mécanismes techniques implémentés. L'Article 32 (sécurité du traitement) est partiellement satisfait : JWT/bcrypt et TLS sont en place, mais le chiffrement des données au repos dans Firestore n'est pas encore activé au niveau applicatif (chiffrement de champ AES-256). C'est la priorité 3 du plan d'action correctif. Un DPO peut s'appuyer sur ce tableau pour évaluer la conformité globale.")
print("Slide 17 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 18 — DASHBOARD SCORE 100
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 18)
label(s, 'ANNEXE A.5 — MONITORING APPLICATIF')
title(s, 'Dashboard Accès Sécurisé — Score Global 100 / 100', size=22)
divider(s)
add_screenshot(s, CAP['score'], Inches(0.4), Inches(2.0), Inches(8.0), Inches(3.3),
               'MONITORING — Dashboard Accès Sécurisé · Score 100/100 · Tous contrôles WAF, JWT, RBAC, RGPD passés')
rb = s.shapes.add_shape(1, Inches(8.7), Inches(2.05), Inches(4.3), Inches(0.75))
rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
txbox(s,'Score  100 / 100',Inches(8.8),Inches(2.1),Inches(4.1),Inches(0.55),size=20,bold=True,color=WHITE)
info_list(s, [
    ('WAF OWASP :','Actif et opérationnel · Toutes catégories bloquées'),
    ('JWT Auth :','Aucun bypass détecté · Tokens valides'),
    ('RBAC :','Contrôles d\'accès fonctionnels'),
    ('RGPD endpoints :','Conformes · Effacement opérationnel'),
    ('AuditLogs :','Surveillance immuable active'),
], x=Inches(8.7), y=Inches(3.0), w=Inches(4.3), h_each=Inches(0.37))
notes_zone(s, "Le dashboard de monitoring applicatif évalue l'état de sécurité de la plateforme GestionScolaire en temps réel. Il agrège les métriques de tous les contrôles de sécurité : WAF, JWT, RBAC, RGPD, AuditLogs. Le score 100/100 indique que tous les contrôles sont opérationnels et qu'aucune violation n'est détectée. Ce score est distinct du score SIEM Wazuh : il évalue la couche applicative (Express.js/Firebase), tandis que Wazuh évalue la couche infrastructure/système. Un score parfait sur les deux niveaux valide la stratégie Defense in Depth : chaque couche de sécurité est indépendante et opérationnelle.")
print("Slide 18 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 19 — DASHBOARD RGPD LIVE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 19)
label(s, 'ANNEXE A.2 — DASHBOARD RGPD')
title(s, 'Conformité RGPD Dashboard Live — Art. 15 / 17 / 33', size=22)
divider(s)
info_list(s, [
    ('Art. 15 — Accès :','Export données étudiant disponible en JSON via endpoint /api/rgpd/export'),
    ('Art. 17 — Effacement :','Endpoint DELETE /api/rgpd/erase opérationnel + log d\'effacement'),
    ('Art. 33 — Notification :','Alertes Wazuh configurées · Délai légal 72h respecté'),
    ('Traçabilité :','9 types d\'events auditLogs Firestore immuables enregistrés'),
    ('Minimisation :','Données pseudonymisées · Collecte minimale active'),
], x=Inches(0.4), y=Inches(2.05), w=Inches(5.5), h_each=Inches(0.40))
add_screenshot(s, CAP['rgpd'], Inches(6.1), Inches(2.0), Inches(7.0), Inches(3.3),
               'MONITORING — Conformité RGPD Art. 15/17/33 · Statut temps réel · Endpoints conformes')
notes_zone(s, "Ce dashboard RGPD affiche le statut de conformité des articles les plus techniques, directement mesurables par des indicateurs applicatifs. L'Article 15 (droit d'accès) est satisfait par un endpoint d'export JSON des données personnelles du compte connecté. L'Article 17 (droit à l'effacement) est satisfait par un endpoint de suppression/anonymisation avec confirmation. L'Article 33 (notification de violation) s'appuie sur les alertes Wazuh pour détecter et notifier toute violation potentielle dans le délai légal de 72 heures. Ce dashboard permet à un DPO (Data Protection Officer) de vérifier la conformité sans accès technique aux systèmes backend.")
print("Slide 19 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 20 — RBAC & AUDIT LOGS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 20)
label(s, 'ANNEXE A.3 — RBAC & AUDIT LOGS')
title(s, 'RBAC 4 Rôles & Audit Logs Firestore Immuables', size=22)
divider(s)
add_screenshot(s, CAP['rbac'], Inches(0.4), Inches(2.0), Inches(7.8), Inches(3.3),
               'MONITORING — RBAC 4 rôles · AuditLogs Firestore immuables · Journal 9 types d\'événements')
cols20 = [('Rôle',1.8),('Droits & Périmètre',3.6)]
fracs20 = [c[1]/5.4 for c in cols20]
table_hdr(s, cols20, Inches(8.4), Inches(2.05), Inches(5.4))
rows20 = [
    ('admin','Accès total · Gestion utilisateurs · Rapports · Configuration'),
    ('sous-admin','Gestion élèves · Classes · Tarifs · Inscriptions'),
    ('comptable','Factures · Paiements · Exports · Échelonnements'),
    ('étudiant/parent','Consultation propre dossier uniquement · Lecture seule'),
]
for i,row in enumerate(rows20):
    table_row(s, row, fracs20, Inches(8.4), Inches(2.35)+Inches(0.35)*i, Inches(5.4), Inches(0.32), even=(i%2==0))
txbox(s,'9 types d\'auditLogs immuables : login · logout · create · update · delete · export · rgpd_erase · role_change · access_denied',
      Inches(8.4), Inches(3.85), Inches(5.4), Inches(0.45), size=8.5, color=WHITE_DIM)
notes_zone(s, "Le RBAC (Role-Based Access Control) définit 4 rôles avec des permissions strictement séparées. Chaque accès, modification, ou tentative non autorisée génère un auditLog immuable dans Firestore. L'immutabilité est garantie par les règles Firestore Security Rules : seule la création est autorisée (pas de mise à jour ni suppression). Les 9 types d'événements couvrent le cycle de vie complet des données : authentification (login/logout), opérations CRUD (create/update/delete), exports, effacements RGPD, changements de rôle, et violations d'accès. Cette traçabilité est essentielle pour la conformité RGPD et les investigations forensiques post-incident (DFIR).")
print("Slide 20 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 21 — WAF OWASP
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 21)
label(s, 'ANNEXE A.4 — WAF OWASP TOP 10')
title(s, 'WAF OWASP Top 10 Actif — 4 Catégories d\'Attaques Bloquées', size=22)
divider(s)
cols21 = [('Attaque',2.0),('Référence OWASP',1.8),('Pattern Détecté',2.5),('Statut',1.1)]
fracs21 = [c[1]/7.4 for c in cols21]
table_hdr(s, cols21, Inches(0.4), Inches(2.05), Inches(7.4))
rows21 = [
    ('SQL Injection','A03:2021 — Injection','SELECT / UNION / DROP / INSERT dans param HTTP','✓ BLOQUÉ'),
    ('Cross-Site Scripting (XSS)','A03:2021 — Injection','<script> / onerror / javascript: dans param','✓ BLOQUÉ'),
    ('Path Traversal / LFI','A01:2021 — Broken Access','../../ / %2e%2e dans URL path','✓ BLOQUÉ'),
    ('Command Injection','A03:2021 — Injection','; && | $( cmd ) dans paramètre HTTP','✓ BLOQUÉ'),
]
for i,row in enumerate(rows21):
    table_row(s, row, fracs21, Inches(0.4), Inches(2.35)+Inches(0.38)*i, Inches(7.4), Inches(0.35), even=(i%2==0))
txbox(s,'WAF middleware waf.js — Analyse regex sur chaque requête HTTP · Réponse 403 immédiate · Log Wazuh Rules 100000–100003 · Intégré Express.js',
      Inches(0.4), Inches(3.95), Inches(7.4), Inches(0.35), size=8.5, color=WHITE_DIM, italic=True)
add_screenshot(s, CAP['waf'], Inches(8.0), Inches(2.0), Inches(5.1), Inches(3.3),
               'MONITORING — WAF OWASP Top 10 · SQLi · XSS · Path Traversal · CMD Injection bloqués')
notes_zone(s, "Le WAF est implémenté sous forme de middleware Express.js (waf.js) qui intercepte toutes les requêtes HTTP entrantes avant qu'elles n'atteignent les routes applicatives. Des expressions régulières détectent les patterns d'attaque OWASP Top 10 les plus courants. En cas de détection, la requête est immédiatement rejetée (HTTP 403) et un log est émis vers Wazuh (règles 100000–100003). Le dashboard WAF montre en temps réel les catégories d'attaques bloquées. Un WAF applicatif (couche 7) protège spécifiquement contre les injections dans les paramètres HTTP, que ne peut pas détecter un firewall réseau traditionnel (couche 3/4).")
print("Slide 21 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 22 — SIEM LOGS TEMPS RÉEL
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 22)
label(s, 'ANNEXE A.1 — SIEM LOGS TEMPS RÉEL')
title(s, 'SIEM Logs Temps Réel — Score 100/100 — Surveillance Normale', size=22)
divider(s)
add_screenshot(s, CAP['siem'], Inches(0.4), Inches(2.0), Inches(8.0), Inches(3.3),
               'SIEM LOGS — Score 100/100 · 0 auth events · 0 RBAC · 0 WAF blocks · Journal 20 derniers événements')
info_list(s, [
    ('Score :','100/100 — Aucune anomalie de sécurité détectée'),
    ('Auth events :','0 — Aucune tentative d\'authentification suspecte'),
    ('RBAC violations :','0 — Contrôles d\'accès respectés'),
    ('RGPD events :','0 — Aucune violation de données personnelles'),
    ('WAF blocks :','0 — Aucune attaque applicative en cours'),
    ('Statut :','Surveillance normale — Système sain'),
], x=Inches(8.7), y=Inches(2.05), w=Inches(4.3), h_each=Inches(0.37))
notes_zone(s, "Cet onglet SIEM Logs affiche les 20 derniers événements de sécurité enregistrés dans les auditLogs Firestore. Les compteurs à 0 (0 auth events, 0 RBAC violations, 0 WAF blocks, 0 RGPD events) indiquent un état de surveillance normale : aucune attaque en cours au moment de la capture. Le score 100/100 confirme que tous les contrôles de sécurité sont actifs. Ce tableau de bord simule ce que verrait un opérateur SOC niveau 1 en monitoring continu. L'absence d'événements dans le journal signifie que l'application fonctionne normalement, sans tentative d'intrusion détectée durant la période observée. En cas d'incident, les events apparaîtraient ici en temps réel.")
print("Slide 22 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 23 — PLAN D'ACTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 23)
label(s, 'SECTION 6 — PLAN D\'ACTION CORRECTIF')
title(s, 'Recommandations & Plan d\'Action Priorisé — 6 Priorités', size=22)
divider(s)
recos = [
    ('P1 — CRITIQUE','Patch CVE-2019-5736 — Mise à jour Docker Engine',
     'Mise à jour runc ≥1.0.0-rc8 · Docker Engine ≥26.1.5 · Isolement réseau conteneurs · Audit droits root · Test évasion post-patch'),
    ('P2 — HAUTE','Hardening SCA — Objectif ≥7/10 (actuellement 0/10)',
     'CIS Benchmark macOS Level 2 · Désactivation services inutiles · Firewall pf activé · SSH hardening · Audit sudoers · Rotation secrets /etc'),
    ('P3 — HAUTE','Chiffrement at-rest — Conformité Article 32 RGPD',
     'Chiffrement champs Firestore AES-256 (nom, prénom, email) · Gestion clés KMS Google Cloud · Transparent pour l\'application'),
    ('P4 — MOYENNE','Rotation JWT & Refresh Tokens sécurisés',
     'Durée JWT réduite à 15min · Refresh token 7 jours avec révocation · Blacklist tokens compromis Redis · Rotation automatique clés HS256'),
    ('P5 — MOYENNE','Patch toutes CVE identifiées (17 au total)',
     'npm audit fix --force · Mise à jour dépendances Node.js · Scan automatique hebdomadaire Dependabot/Snyk · CI/CD gate CVE critique'),
    ('P6 — BASSE','Monitoring continu & Alertes Wazuh automatisées',
     'Alertes e-mail/Slack pour Rule Niveau >7 · Rapport KPI hebdomadaire · Rapport SCA mensuel · Dashboard directeur SSI'),
]
for i,(prio,name,desc) in enumerate(recos):
    cy = Inches(2.05 + i*0.54)
    rb = s.shapes.add_shape(1, Inches(0.4), cy, Inches(12.5), Inches(0.50))
    rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.3)
    txbox(s, prio, Inches(0.55), cy+Pt(2), Inches(1.4), Inches(0.25), size=7.5, bold=True, color=GREY)
    txbox(s, name, Inches(2.0), cy+Pt(2), Inches(4.2), Inches(0.25), size=9, bold=True, color=WHITE)
    txbox(s, desc, Inches(6.3), cy+Pt(2), Inches(6.5), Inches(0.35), size=8.5, color=WHITE_DIM)
notes_zone(s, "Le plan d'action est organisé en 6 priorités selon criticité et effort. P1 (CVE-2019-5736 CVSS3=8.6) : non négociable, patch immédiat. P2 (SCA 0/10) : hardening système réalisable en 1-2 jours. P3 (Art. 32 RGPD) : obligation légale avant mise en production avec données réelles. P4-P5 : renforcent la posture à moyen terme. P6 : surveillance continue long terme. Ce plan suit la méthodologie standard de gestion des risques : Identifier → Évaluer → Traiter → Surveiller (ISO 27005, NIST RMF). Chaque action est mesurable et vérifiable par une nouvelle capture de métriques post-correction.")
print("Slide 23 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 24 — BILAN FINAL
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide(); bg(s); logo_txt(s); slide_num(s, 24)
label(s, 'SYNTHÈSE FINALE')
title(s, 'Bilan Final — Projet de Sécurité M2 Cybersécurité', size=22)
divider(s)
bilan = [
    ('FIM — Intégrité Fichiers','2 500+','événements · Rule 550 Niveau 7 · root 99.44% · /bin surveillé'),
    ('CVE — Vulnérabilités','17 identifiées','8 High + 9 Medium · CVE-2019-5736 CVSS3=8.6 · Plan patch P1'),
    ('MITRE ATT&CK','1 600+ alertes','Impact T1499 ×1594 · Defense Evasion ×6 · Corrélation auto'),
    ('SCA — Configuration','0 / 10','Surface d\'attaque maximale · Plan hardening CIS Level 2 · P2'),
    ('Conformité RGPD','5 / 6','Art. 5,17,25,33,35 validés · Art. 32 chiffrement en cours (P3)'),
    ('Score Sécurité Applicatif','100 / 100','WAF + JWT/bcrypt + RBAC + AuditLogs · 0 violations détectées'),
    ('Problèmes Techniques','8 / 8','100% résolus · Démarche diagnostique rigoureuse documentée'),
    ('Couverture SIEM','100%','2 agents actifs · Docker + macOS · Tous modules Wazuh opérationnels'),
]
for i,(lbl2,val,note) in enumerate(bilan):
    col = i//4; row = i%4
    cx = Inches(0.4 + col*6.5); cy = Inches(2.05 + row*1.1)
    rb = s.shapes.add_shape(1, cx, cy, Inches(6.2), Inches(1.0))
    rb.fill.solid(); rb.fill.fore_color.rgb = NAVY_M
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.5)
    txbox(s, lbl2, cx+Inches(0.15), cy+Pt(3), Inches(5.9), Inches(0.22), size=7.5, bold=True, color=GREY)
    txbox(s, val, cx+Inches(0.15), cy+Inches(0.28), Inches(5.9), Inches(0.35), size=16, bold=True, color=WHITE)
    txbox(s, note, cx+Inches(0.15), cy+Inches(0.66), Inches(5.9), Inches(0.28), size=7.5, color=WHITE_DIM)
notes_zone(s, "Ce bilan démontre la réalisation complète des objectifs du projet. 8 métriques clés attestent du niveau de sécurité : détection FIM (2500+), analyse CVE (17), corrélation MITRE (1600+), audit SCA (0/10 → plan correctif), conformité RGPD (5/6), score applicatif (100/100), résilience (8/8 problèmes résolus), couverture SIEM (100%). La combinaison SIEM Wazuh (infrastructure) + dashboard monitoring applicatif (couche web) constitue une approche Defense in Depth complète. Les points d'amélioration identifiés sont documentés avec un plan d'action précis, démontrant la maturité de l'approche de sécurité adoptée.")
print("Slide 24 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 25 — CONCLUSION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = new_slide()
bg(s, BLACK)
r = s.shapes.add_shape(1, 0, 0, Inches(7), H)
r.fill.solid(); r.fill.fore_color.rgb = NAVY_M; r.line.fill.background()
logo_txt(s); slide_num(s, 25)
txbox(s, 'CONCLUSION — RAPPORT DE SÉCURITÉ v2.0 — 6 AVRIL 2026',
      Inches(0.5), Inches(0.85), Inches(10), Inches(0.28), size=8, color=GREY, bold=True)
r2 = s.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(1.0), Pt(2))
r2.fill.solid(); r2.fill.fore_color.rgb = WHITE; r2.line.fill.background()
txbox(s, 'Projet de Sécurité', Inches(0.5), Inches(1.3), Inches(9), Inches(0.65),
      size=30, bold=True, color=WHITE)
txbox(s, 'Accompli avec Succès', Inches(0.5), Inches(1.95), Inches(9), Inches(0.5),
      size=22, color=WHITE_DIM)
concl = [
    'SIEM Wazuh 4.7.4 pleinement opérationnel — Infrastructure Docker + 2 agents, couverture 100%',
    'Détection multi-vecteur prouvée — 2500+ FIM, 17 CVE, 1600+ MITRE ATT&CK documentés',
    'Application sécurisée by design — JWT HS256 + bcrypt + WAF OWASP + RBAC + AuditLogs · Score 100/100',
    'Conformité RGPD 5/6 articles — Privacy by Design, effacement, notification, AIPD documentée',
    '8 problèmes techniques résolus à 100% — Démarche rigoureuse sur tous les niveaux du stack',
    'Plan d\'action correctif défini — 6 priorités (SCA, CVE, Art. 32, JWT) pour renforcement continu',
]
for i,c in enumerate(concl):
    cy = Inches(2.55 + i*0.45)
    rb = s.shapes.add_shape(1, Inches(0.4), cy, Inches(12.5), Inches(0.40))
    rb.fill.solid(); rb.fill.fore_color.rgb = RGBColor(0x0f,0x22,0x3a)
    rb.line.color.rgb = RGBColor(0x2a,0x45,0x70); rb.line.width = Pt(0.3)
    txbox(s, '—', Inches(0.55), cy+Pt(3), Inches(0.3), Inches(0.3), size=9, bold=True, color=GREY)
    txbox(s, c, Inches(0.9), cy+Pt(3), Inches(11.9), Inches(0.35), size=9, color=WHITE_DIM)
metas2 = [('Auteur','Anass AKKER'),('Formation','Master 2 Cybersécurité · YNOV Campus'),('Date','6 Avril 2026')]
for i,(lbl2,val) in enumerate(metas2):
    mx = Inches(0.4 + i*4.1)
    txbox(s, lbl2.upper(), mx, Inches(5.35), Inches(3.9), Inches(0.2), size=7, color=GREY, bold=True)
    txbox(s, val, mx, Inches(5.57), Inches(3.9), Inches(0.3), size=10, bold=True, color=WHITE)
notes_zone(s, "Ce projet démontre la mise en œuvre complète d'une stratégie Defense in Depth pour une application web de gestion scolaire. L'approche multicouche est illustrée par : couche infrastructure (Wazuh SIEM), couche OS (FIM, SCA, CVE), couche réseau (WAF), couche applicative (JWT, RBAC, bcrypt), et couche données (AuditLogs Firestore, RGPD). Les métriques (2500+ FIM, 17 CVE, 1600+ MITRE) constituent une preuve par les faits du niveau de menace réel. Les perspectives : automatisation du plan correctif, intégration SOC managé, certification ISO 27001. Ce projet démontre la capacité à concevoir, déployer, opérer et améliorer un système de sécurité de niveau professionnel.")
print("Slide 25 ✓")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SAUVEGARDE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out = os.path.join(BASE, 'PRESENTATION_SECURITE_M2.pptx')
prs.save(out)
print(f"\n✓ Fichier généré : {out}")
import os; sz = os.path.getsize(out)
print(f"  Taille : {sz/1024/1024:.1f} Mo")
