#!/usr/bin/env python3
"""
Répare les positions corrompues des shapes du slide 7 (RGPD §3.3)
Recalcule une grille 3x2 propre pour les 6 articles RGPD.
"""

from pptx import Presentation
from pptx.util import Emu

SRC = "/Users/anass/Downloads/frais-gestionScolaire 4/PRESENTATION_PFE_SECURITE_UPDATED.pptx"
DST = SRC

prs = Presentation(SRC)
slide = prs.slides[6]  # slide 7 (0-indexed)

# ─── Dimensions calculées ──────────────────────────────────────────────────
# Slide: 12188952 x 6858000 EMU (13.33" x 7.50")
# 3 colonnes de 3657600 EMU (4.00") chacune
# Marge gauche: 411480
# Gap inter-colonne: (12188952 - 411480 - 3×3657600) / 2 = 402336
#
# Rangée 1 top: 1600200
# Rangée 2 top: 1600200 + 1234440 + 300000 = 3134640
# Hauteur carte: 1234440, Hauteur header: 292608
# ──────────────────────────────────────────────────────────────────────────

COL = [411480, 4471416, 8531352]   # Left des 3 colonnes
ROW = [1600200, 3134640]           # Top des 2 rangées

CARD_H  = 1234440
HDR_H   = 292608
HDR_DY  = 38100    # offset top du texte header dans la carte
BODY_DX = 109728   # offset left du texte body
BODY_DY = 329184   # offset top du texte body
BODY_W  = 3474720
BODY_H  = 850392
CARD_W  = 3657600
HDR_W   = 3657600
HDR_TW  = 3657600
HDR_TH  = 256032

# Mapping: nom_shape → (col_idx, row_idx, role)
# role: 'rect_full', 'rect_hdr', 'tb_hdr', 'tb_body'
LAYOUT = {
    # Rangée 1
    "Rectangle 6":  (0, 0, "rect_full"),
    "Rectangle 7":  (0, 0, "rect_hdr"),
    "TextBox 8":    (0, 0, "tb_hdr"),
    "TextBox 9":    (0, 0, "tb_body"),

    "Rectangle 10": (1, 0, "rect_full"),
    "Rectangle 11": (1, 0, "rect_hdr"),
    "TextBox 12":   (1, 0, "tb_hdr"),
    "TextBox 13":   (1, 0, "tb_body"),

    "Rectangle 14": (2, 0, "rect_full"),
    "Rectangle 15": (2, 0, "rect_hdr"),
    "TextBox 16":   (2, 0, "tb_hdr"),
    "TextBox 17":   (2, 0, "tb_body"),

    # Rangée 2
    "Rectangle 18": (0, 1, "rect_full"),
    "Rectangle 19": (0, 1, "rect_hdr"),
    "TextBox 20":   (0, 1, "tb_hdr"),
    "TextBox 21":   (0, 1, "tb_body"),

    "Rectangle 22": (1, 1, "rect_full"),
    "Rectangle 23": (1, 1, "rect_hdr"),
    "TextBox 24":   (1, 1, "tb_hdr"),
    "TextBox 25":   (1, 1, "tb_body"),

    "Rectangle 26": (2, 1, "rect_full"),
    "Rectangle 27": (2, 1, "rect_hdr"),
    "TextBox 28":   (2, 1, "tb_hdr"),
    "TextBox 29":   (2, 1, "tb_body"),
}


def calc_pos(col, row, role):
    """Retourne (left, top, width, height) pour un shape donné."""
    cl = COL[col]
    rt = ROW[row]

    if role == "rect_full":
        return cl, rt, CARD_W, CARD_H
    elif role == "rect_hdr":
        return cl, rt, HDR_W, HDR_H
    elif role == "tb_hdr":
        return cl, rt + HDR_DY, HDR_TW, HDR_TH
    elif role == "tb_body":
        return cl + BODY_DX, rt + BODY_DY, BODY_W, BODY_H
    else:
        raise ValueError(f"role inconnu: {role}")


# ─── Appliquer les corrections ────────────────────────────────────────────
fixed = 0
for shape in slide.shapes:
    if shape.name in LAYOUT:
        col, row, role = LAYOUT[shape.name]
        left, top, w, h = calc_pos(col, row, role)

        old_l, old_t = shape.left, shape.top
        shape.left   = Emu(left)
        shape.top    = Emu(top)
        shape.width  = Emu(w)
        shape.height = Emu(h)

        print(f"  ✓ {shape.name:15s} ({role:10s}) col={col} row={row} → "
              f"L={left} T={top} W={w} H={h}  [was L={old_l} T={old_t}]")
        fixed += 1

print(f"\n✓ {fixed} shapes repositionnés")

prs.save(DST)
print(f"✅ Sauvegardé : {DST}")
