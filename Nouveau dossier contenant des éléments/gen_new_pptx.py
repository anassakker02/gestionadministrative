#!/usr/bin/env python3
"""
Génère une nouvelle présentation PPTX en remplaçant les images et textes
des slides Wazuh (22, 23, 25) et en ajoutant un slide Policy Monitoring (26).
"""

import copy
import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

BASE = "/Users/anass/Downloads/frais-gestionScolaire 4"
SRC  = os.path.join(BASE, "PRESENTATION_PFE_SECURITE_30SLIDES.pptx")
DST  = os.path.join(BASE, "PRESENTATION_PFE_SECURITE_UPDATED.pptx")

IMGS = {
    "security_events":    os.path.join(BASE, "SECURITEEVENTS.png"),
    "fim":                os.path.join(BASE, "INTEGRITYMONITORING.png"),
    "mitre":              os.path.join(BASE, "MITTRE ATTACK.png"),
    "policy_monitoring":  os.path.join(BASE, "POLICY MONITORING.png"),
}

shutil.copy2(SRC, DST)
prs = Presentation(DST)


# ─────────────────────────────────────────────────────────────────────────────
# Helper : remplace l'image d'un Picture shape
# ─────────────────────────────────────────────────────────────────────────────
def replace_picture(slide, shape_name, img_path):
    for shape in slide.shapes:
        if shape.shape_type == 13 and shape.name == shape_name:
            with open(img_path, "rb") as f:
                img_bytes = f.read()
            blip = shape._element.blipFill.blip
            r_embed = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if r_embed:
                part = slide.part
                img_part = part.rels[r_embed].target_part
                img_part._blob = img_bytes
                print(f"  ✓ Image remplacée : {shape_name} → {os.path.basename(img_path)}")
                return True
    print(f"  ✗ Shape '{shape_name}' non trouvé")
    return False


# ─────────────────────────────────────────────────────────────────────────────
# Helper : cherche un TextBox par son contenu (début) et remplace le texte
# ─────────────────────────────────────────────────────────────────────────────
def set_textbox_text(slide, old_prefix, new_text):
    for shape in slide.shapes:
        if shape.has_text_frame:
            current = shape.text_frame.text
            if current.startswith(old_prefix):
                tf = shape.text_frame
                # Garde le formatage du premier run
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                tf.paragraphs[0].runs[0].text = new_text
                print(f"  ✓ Texte mis à jour : '{old_prefix[:40]}' → '{new_text[:60]}'")
                return True
    return False


def set_textbox_by_index(slide, shape_name, new_text):
    """Remplace le texte d'un TextBox par son nom de shape."""
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            # Vide tout sauf le premier paragraphe/run
            full_text = tf.text
            # Reconstruire en gardant le style
            for i, para in enumerate(tf.paragraphs):
                if i == 0:
                    if para.runs:
                        para.runs[0].text = new_text
                        for r in para.runs[1:]:
                            r.text = ""
                    else:
                        para.add_run().text = new_text
                else:
                    for run in para.runs:
                        run.text = ""
            print(f"  ✓ [{shape_name}] → '{new_text[:70]}'")
            return True
    return False


# ─────────────────────────────────────────────────────────────────────────────
# Helper : remplace tous les textes d'un slide par mapping {shape_name: text}
# ─────────────────────────────────────────────────────────────────────────────
def update_slide_texts(slide, mapping):
    for shape in slide.shapes:
        if shape.name in mapping and shape.has_text_frame:
            new_text = mapping[shape.name]
            tf = shape.text_frame
            all_runs = []
            for para in tf.paragraphs:
                all_runs.extend(para.runs)
            if all_runs:
                all_runs[0].text = new_text
                for r in all_runs[1:]:
                    r.text = ""
            else:
                tf.paragraphs[0].add_run().text = new_text
            print(f"  ✓ [{shape.name}] = '{new_text[:70]}'")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 22 — SECURITY EVENTS (remplace Agent 002)
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== Slide 22 — Security Events ===")
slide22 = prs.slides[21]  # 0-indexed

replace_picture(slide22, "Picture 9", IMGS["security_events"])

update_slide_texts(slide22, {
    "TextBox 3": "WAZUH — SECURITY EVENTS",
    "TextBox 4": "Security Events — 436 420 Événements — T1565.001 · Rule 550 Niveau 7",
    "TextBox 7": "Security Events — 436 420 events · Top rule 550 (Integrity checksum) · T1565.001 Stored Data Manipulation",
    "TextBox 11": "Security Events — 436 420 events · Rule 550 · Level 7 · T1565.001 · FIM dominant",
    "TextBox 13": "Total Events :  436 420 — journalisation complète",
    "TextBox 15": "Top Rule :  Rule 550 Level 7 — Integrity checksum changed",
    "TextBox 17": "MITRE :  T1565.001 Stored Data Manipulation dominant",
    "TextBox 19": "Agent :  main-machine · macOS 15.7.4 · Active · 100% coverage",
    "TextBox 21": "NOTE",
    "TextBox 22": (
        "Le tableau Security Events centralise tous les événements de sécurité détectés par Wazuh. "
        "Avec 436 420 événements collectés, la Rule 550 (Integrity checksum changed, Level 7) est la plus déclenchée, "
        "mappée sur la technique MITRE T1565.001 (Stored Data Manipulation). "
        "Cela démontre une surveillance exhaustive des modifications de fichiers système sur l'agent macOS, "
        "garantissant la détection de tout rootkit ou backdoor éventuel."
    ),
    "TextBox 2": "22 / 30",
})

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 23 — FILE INTEGRITY MONITORING
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== Slide 23 — FIM ===")
slide23 = prs.slides[22]

replace_picture(slide23, "Picture 9", IMGS["fim"])

update_slide_texts(slide23, {
    "TextBox 3": "WAZUH — FILE INTEGRITY MONITORING",
    "TextBox 4": "FIM — 436 420 Événements · Rule 550 Level 7 · root 89.44% · /var/bin/afsa",
    "TextBox 7": "FIM Dashboard — Most active users: root 89.44% · Actions: modified 100% · /var/bin/afsa · /var/bin/apfs · /private/var/db/...",
    "TextBox 11": "FIM — root 89.44% des événements · modified 100% · /var/bin/afsa dominant",
    "TextBox 13": "Règle :  Rule 550 Niveau 7 — Integrity checksum changed",
    "TextBox 15": "root 89.44% :  Modifications binaires système macOS — /var/bin/*",
    "TextBox 17": "Fichiers :  /var/bin/afsa · /var/bin/apfs · /private/var/db/locationd",
    "TextBox 19": "Risque :  Rootkit/backdoor si modifications non légitimes des binaires",
    "TextBox 21": "NOTE",
    "TextBox 22": (
        "Le module FIM (File Integrity Monitoring) surveille les modifications de fichiers critiques via hash SHA-256. "
        "Les 436 420 événements montrent que root est responsable de 89.44% des modifications, "
        "toutes de type 'modified' (100%). Les fichiers les plus touchés sont /var/bin/afsa, /var/bin/apfs et "
        "/private/var/db/locationd, correspondant aux mises à jour système macOS légitimes. "
        "Tout écart non autorisé déclencherait Rule 550 (Level 7) — détection rootkit/backdoor en temps réel."
    ),
    "TextBox 2": "23 / 30",
})

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 25 — MITRE ATT&CK
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== Slide 25 — MITRE ATT&CK ===")
slide25 = prs.slides[24]

replace_picture(slide25, "Picture 9", IMGS["mitre"])

update_slide_texts(slide25, {
    "TextBox 3": "WAZUH — MITRE ATT&CK",
    "TextBox 4": "MITRE ATT&CK — T1565.001 Stored Data Manipulation dominant — Impact vs Defense Evasion",
    "TextBox 7": "MITRE ATT&CK — Top techniques: T1565.001 Stored Data Manipulation (Impact) · T1562 Disable or Modify Tools (Defense Evasion)",
    "TextBox 11": "MITRE ATT&CK — Stored Data Manipulation dominant · Impact tactic · T1562 Defense Evasion",
    "TextBox 13": "T1565.001 :  Stored Data Manipulation — tactic Impact (dominant)",
    "TextBox 15": "T1562 :  Disable or Modify Tools — tactic Defense Evasion",
    "TextBox 17": "Tactics :  Impact > Defense Evasion · Classification internationale SOC",
    "TextBox 19": "Corrélation :  PCI DSS 11.5 · RGPD Art.25 · Privacy by Design ✓",
    "TextBox 21": "NOTE",
    "TextBox 22": (
        "MITRE ATT&CK est le framework universel de classification des comportements malveillants. "
        "Wazuh mappe automatiquement ses règles aux techniques ATT&CK. "
        "La technique dominante est T1565.001 (Stored Data Manipulation) dans la tactique Impact, "
        "correspondant aux modifications de fichiers système détectées par FIM. "
        "T1562 (Disable or Modify Tools) en Defense Evasion représente les tentatives de contournement des outils de sécurité. "
        "Cette cartographie MITRE valide notre couverture CDC §3.3 et démontre une surveillance alignée sur les standards SOC internationaux."
    ),
    "TextBox 2": "25 / 30",
})


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 26 — POLICY MONITORING / ROOTCHECK (nouveau slide)
# Stratégie : dupliquer le slide 25 et modifier son contenu + image
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== Ajout Slide 26 — Policy Monitoring / Rootcheck ===")

def duplicate_slide(prs, slide_index):
    """Duplique un slide existant et l'insère à la fin."""
    template_slide = prs.slides[slide_index]
    slide_layout = template_slide.slide_layout

    # Crée un nouveau slide avec le même layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copie tous les éléments XML du slide template
    spTree = new_slide.shapes._spTree

    # Supprime les placeholders par défaut du layout
    for elem in list(spTree):
        spTree.remove(elem)

    # Copie l'arbre XML complet du template
    for elem in template_slide.shapes._spTree:
        spTree.append(copy.deepcopy(elem))

    # Copie les relations (images, etc.)
    for rel in template_slide.part.rels.values():
        if "image" in rel.reltype:
            target_part = rel.target_part
            new_slide.part.relate_to(target_part, rel.reltype)

    return new_slide


# Duplique slide 25 (index 24) → nouveau slide en fin
new_slide_26 = duplicate_slide(prs, 24)

# Remplace image
replace_picture(new_slide_26, "Picture 9", IMGS["policy_monitoring"])

# Met à jour les textes
update_slide_texts(new_slide_26, {
    "TextBox 3": "WAZUH — POLICY MONITORING (ROOTCHECK)",
    "TextBox 4": "Rootcheck — Détection Anomalies Système · Processus Cachés · Interface Promiscuité",
    "TextBox 7": "Policy Monitoring — Host-based anomaly detection · Trojaned files · Hidden processes · en3 promiscuous mode",
    "TextBox 11": "Rootcheck — Anomalies détectées : interface en3 promiscuous · processus cachés · trojaned files",
    "TextBox 13": "Trojaned files :  Détection de binaires système compromis (trojaned)",
    "TextBox 15": "Hidden processes :  Processus cachés détectés — tentative de dissimulation",
    "TextBox 17": "Promiscuity :  Interface en3 en mode promiscuous — capture réseau détectée",
    "TextBox 19": "Module :  Host-based Anomaly Detection · Rootcheck actif · temps réel",
    "TextBox 21": "NOTE",
    "TextBox 22": (
        "Le module Rootcheck de Wazuh effectue une détection d'anomalies basée sur l'hôte (HIDS). "
        "Il vérifie les signatures de rootkits connus, les fichiers trojans et les comportements suspects. "
        "Sur notre agent macOS, trois anomalies ont été détectées : (1) interface réseau en3 en mode promiscuous "
        "(capture de tout le trafic réseau — indicateur de sniffing), (2) processus cachés (tentative de dissimulation), "
        "(3) fichiers trojans potentiels. Ces alertes démontrent la capacité de Wazuh à détecter des compromissions avancées "
        "au niveau système, complétant la surveillance MITRE ATT&CK et FIM."
    ),
    "TextBox 2": "26 / 30",
})

print(f"\n✓ Slide 26 ajouté. Total slides: {len(prs.slides)}")


# ─────────────────────────────────────────────────────────────────────────────
# Réordonne les slides : met le nouveau slide 26 à la position 26
# (il a été ajouté en fin, on doit le déplacer à l'index 25)
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== Réordonnancement des slides ===")

xml_slides = prs.slides._sldIdLst
# Le nouveau slide est le dernier (index -1), on veut l'insérer à la position 25 (après slide 25)
slides_list = list(xml_slides)
last = slides_list[-1]
# Position cible : index 25 (entre slide 25 et l'ancien slide 26)
xml_slides.remove(last)
xml_slides.insert(25, last)
print(f"✓ Slide réordonné à la position 26")


# ─────────────────────────────────────────────────────────────────────────────
# Met à jour les numéros de slides pour les slides décalés (27-31 → 27-31)
# ─────────────────────────────────────────────────────────────────────────────
print("\n=== Mise à jour numéros de slides ===")
total = len(prs.slides)
print(f"Total final: {total} slides")

# Slide 26 (le nouveau, index 25) : déjà mis à 26/30 ci-dessus
# Les anciens slides 26-30 sont maintenant aux index 26-30
# On les met à jour : 27/30, 28/30, 29/30, 30/30, 31/30
old_nums = {
    26: "27",  # ancien slide 26 (MODULES) devient 27
    27: "28",  # DÉFENSE EN PROFONDEUR → 28
    28: "29",  # SYNTHÈSE FINALE → 29
    29: "30",  # BILAN GLOBAL → 30
    30: "31",  # Merci → 31
}

for idx, new_num in old_nums.items():
    if idx < total:
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if shape.name == "TextBox 2" and shape.has_text_frame:
                old_t = shape.text_frame.text
                new_t = f"{new_num} / {total}"
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = new_t
                else:
                    tf.paragraphs[0].add_run().text = new_t
                print(f"  Slide {idx+1}: '{old_t}' → '{new_t}'")

# Met à jour le slide 31 (Merci) avec le bon numéro
# Ajuste aussi le slide 22/23/24/25 si besoin (déjà fait)
# Slide 24 (CVE) → 24/31
slide24 = prs.slides[23]
for shape in slide24.shapes:
    if shape.name == "TextBox 2" and shape.has_text_frame:
        tf = shape.text_frame
        if tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = f"24 / {total}"

# ─────────────────────────────────────────────────────────────────────────────
# Sauvegarde
# ─────────────────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\n✅ Présentation sauvegardée : {DST}")
print(f"   Total slides : {total}")
