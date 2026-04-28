#!/usr/bin/env python3
"""
add_notes_presentation.py
Ajoute les Q&A en notes de présentateur sur les slides concernés.
Les notes sont visibles UNIQUEMENT par le présentateur (mode présentateur PowerPoint).
Le jury ne les voit pas sur l'écran projeté.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

PPTX_FILE = "PRESENTATION_PFE_SECURITE_UPDATED.pptx"

# ── Q&A par slide ─────────────────────────────────────────────────────────────
# Chaque slide reçoit ses questions/réponses en notes de présentateur

NOTES = {

    # ── Slides généraux / intro (1-5) ─────────────────────────────────────────
    1: """🎤 INTRODUCTION — Ce que tu peux dire pour lancer ta présentation :
"Ce projet couvre la mise en production d'une plateforme de gestion scolaire
avec un focus cybersécurité complet : WAF, SIEM Wazuh, monitoring temps réel, RGPD."

❓ Q possible du jury : "Qu'est-ce qui vous a motivé à ajouter autant de sécurité ?"
✅ R : "Les données traitées concernent des étudiants mineurs — données sensibles
au sens du RGPD. La sécurité n'est pas une option, c'est une obligation légale (Art.32)."
""",

    # ── WAF ───────────────────────────────────────────────────────────────────
    2: """🛡️ WAF MIDDLEWARE — Questions possibles :

❓ "Qu'avez-vous implémenté contre les injections ?"
✅ WAF middleware Express dans waf.js — 5 types : SQLi, XSS, Path Traversal,
   Command Injection, agents suspects. → HTTP 403 + log WAF_BLOCK Firestore.

❓ "Quels patterns SQL détectez-vous ?"
✅ SELECT/INSERT/UPDATE/DELETE/DROP/UNION/EXEC + OR 1=1 + SLEEP() + -- ; /* */
   5 regex indépendantes. decodeURIComponent() avant analyse pour encoder %xx.

❓ "Comment le WAF analyse une requête ?"
✅ Dans l'ordre : User-Agent → URL → Query params → Body JSON.
   Password exclu du scan pour éviter faux positifs (apostrophes dans mdp).

❓ "Que se passe-t-il lors d'une attaque ?"
✅ logWafBlock() → Firestore auditLogs {IP, chemin, type, pattern} → HTTP 403.
   Jamais de contact avec le code métier ni Firestore.
""",

    3: """🔐 AUTHENTIFICATION JWT — Questions possibles :

❓ "Comment fonctionne votre connexion ?"
✅ Email + password → bcrypt.compare() → si OK → 2 tokens JWT HS256 :
   access token 30min, refresh token 7j. Log USER_LOGIN_SUCCESS/FAILURE.

❓ "Durée des tokens JWT ?"
✅ Access token : 30 MINUTES. Refresh token : 7 JOURS.
   Secrets séparés : JWT_SECRET et REFRESH_TOKEN_SECRET.

❓ "Comment vérifiez-vous le token ?"
✅ auth.js : jwt.verify() → puis Firestore vérifie user.isActive === true.
   Compte désactivé = HTTP 403 même si token encore valide.

❓ "Que contient le payload JWT ?"
✅ Uniquement { id, email, role } — pas téléphone ni adresse (données PII).
""",

    4: """🔒 HACHAGE & CHIFFREMENT — Questions possibles :

❓ "Comment stockez-vous les mots de passe ?"
✅ bcrypt.hash(password, 10) — jamais en clair. bcrypt.compare() à la connexion.
   saltRounds=10 = 2^10 = 1024 itérations ≈ 100ms par hash.

❓ "saltRounds=10 c'est quoi ?"
✅ Facteur de coût — rend les attaques par dictionnaire non rentables.
   Même mot de passe = deux hashes DIFFÉRENTS grâce au salt aléatoire.

❓ "Le mot de passe est-il visible quelque part ?"
✅ NON. Jamais persisté, jamais loggué. WAF exclut explicitement le champ password.
   Dans les auditLogs : champ password = [REDACTED].

❓ "Quelles données chiffrez-vous ?"
✅ Téléphone + adresse → AES-256-CBC avant Firestore.
   Format stocké : ciphertext:iv. Clé dans ENCRYPTION_KEY (variable env).
   Déchiffrement uniquement après JWT valide.
""",

    5: """⏱️ RATE LIMITING & BRUTE FORCE — Questions possibles :

❓ "Comment protégez-vous contre le brute force ?"
✅ Rate limiter Express sur /api/auth/login :
   5 tentatives par IP / 15 min → 6e tentative = HTTP 429.
   Log AUTH_LOCKOUT dans auditLogs.

❓ "Montrez le résultat en pratique ?"
✅ Tentatives 1-5 → 401 "Email ou mot de passe incorrect"
   Tentative 6    → 429 "Too Many Requests — réessayez dans 15 min"

❓ "Un botnet avec 1000 IPs peut contourner ça ?"
✅ Oui, c'est une limite connue. Contre-mesure non encore implémentée :
   rate limiter par COMPTE en plus de l'IP + CAPTCHA progressif.
   Dans mon plan d'amélioration chapitre 6.
""",

    6: """📋 AUDITLOGS & TRAÇABILITÉ — Questions possibles :

❓ "Quels événements journalisez-vous ?"
✅ 9 types : AUTH_SUCCESS · AUTH_FAILURE · AUTH_LOCKOUT · LOGOUT · SESSION_EXPIRED
             ACCESS_DENIED · DATA_EXPORT · DATA_ANONYMIZE · WAF_BLOCK.

❓ "Comment les logs ne peuvent pas être falsifiés ?"
✅ Firestore Security Rules :
   allow update: if false;  // IMMUABLE
   allow delete: if false;  // IMMUABLE
   Même un admin ne peut pas modifier. Même le backend compromis.

❓ "Que contient un auditLog ?"
✅ userId + action + timestamp (serverTimestamp Google) + metadata :
   {email, IP, chemin, méthode, rôle, raison}
   Timestamp = côté serveur Google, pas manipulable côté client.
""",

    7: """👥 RBAC — CONTRÔLE D'ACCÈS — Questions possibles :

❓ "Combien de rôles et leurs différences ?"
✅ 6 rôles :
   admin      → tout faire
   sous-admin → pas de suppression, pas de création d'admin
   comptable  → finances uniquement (paiements/factures)
   enseignant → lecture seule
   etudiant   → ses données uniquement
   parent     → données de ses enfants uniquement

❓ "Où est vérifiée l'autorisation ?"
✅ DEUX niveaux indépendants :
   1. middleware authorize([roles]) → vérifie req.user.role avant le handler
   2. Firestore Security Rules → bloque accès direct BDD

❓ "Qu'est-ce que deny-by-default ?"
✅ Dernière règle firestore.rules :
   match /{document=**} { allow read, write: if false; }
   Toute collection non listée = refusée automatiquement.
""",

    8: """🖥️ DASHBOARD MONITORING SIEM — Questions possibles :

❓ "Que montre votre dashboard monitoring ?"
✅ Score sécurité /100 + alertes animées + 3 onglets :
   - Dashboard : auth / RGPD / RBAC / journalisation
   - WAF       : attaques bloquées + répartition + 10 dernières
   - SIEM      : journal 20 derniers events {IP, email, rôle, chemin}

❓ "Comment est calculé le score /100 ?"
✅ Départ 100, on soustrait :
   -20 max si auth_failures > 5
   -15 max si lockouts > 0
   -15 max si access_denied > 3
   -20 max si waf_blocks > 0
   80-100=vert · 60-79=orange · 0-59=rouge

❓ "À quelle fréquence se met à jour ?"
✅ setInterval(load, 60_000) → toutes les 60 secondes automatiquement.
   Données : API /monitoring/security → agrège auditLogs 24 dernières heures.
""",

    9: """🛡️ WAZUH SIEM — Questions possibles :

❓ "Pourquoi Wazuh en plus du monitoring applicatif ?"
✅ Monitoring applicatif = événements dans le CODE.
   Wazuh = infrastructure : OS, fichiers, CVE, config serveur.
   Deux couches complémentaires = défense en profondeur.

❓ "Combien d'événements collectés ?"
✅ 436 420 événements depuis le déploiement.
   7 règles custom YNOV (100010-100016) mappées OWASP + MITRE ATT&CK.

❓ "Qu'est-ce que le FIM ?"
✅ File Integrity Monitoring — inotify kernel Linux.
   Surveille /back/functions/src/ → détection < 5 secondes.
   Mappe MITRE ATT&CK T1565.001 (Stored Data Manipulation).

❓ "Quels modules Wazuh ?"
✅ Security Events · MITRE ATT&CK · GDPR · Integrity Monitoring (FIM)
   Vulnerabilities (CVE) · Policy Monitoring (CIS Benchmark)
""",

    10: """🔒 RGPD — Questions possibles :

❓ "Quels articles RGPD satisfaits ?"
✅ Art.5  (intégrité — auditLogs immuables)
   Art.15 (droit d'accès — GET /users/:id/export)
   Art.16 (rectification — PUT /users/:id)
   Art.17 (effacement — pseudonymisation userId)
   Art.30 (registre — Wazuh + auditLogs)
   Art.32 (sécurité technique — WAF/bcrypt/JWT/SIEM)
   TOTAL : 5/6 articles ✅

❓ "Art.17 vs logs immuables — contradiction ?"
✅ Les logs contiennent userId opaque, pas nom/prénom.
   Pour l'effacement : on remplace l'ID par un hash irréversible.
   Ligne reste (intégrité audit) mais n'est plus "donnée personnelle" Art.4.

❓ "Le 6e article manquant ?"
✅ Art.35 — DPIA (Data Protection Impact Assessment).
   Les mesures existent mais le document formel n'est pas produit.
""",

    11: """🔍 SCANNER DAST — Questions possibles :

❓ "Que fait votre scanner DAST ?"
✅ security_scan.js — 12 vraies requêtes HTTP contre l'API :
   auth, rate limiting, WAF (SQLi/XSS/Path/agents),
   headers HTTP, CORS, privilege escalation, payload oversized.

❓ "Résultat obtenu ?"
✅ 12/12 tests PASS — 0 FAIL
   Couverture OWASP : A01 · A02 · A03 · A05 · A07 · A09

❓ "Différence DAST vs SAST ?"
✅ SAST = code source statique (SonarQube).
   DAST = app en cours d'exécution (vraies requêtes HTTP).
   Mon scanner = DAST. ESLint = SAST léger.
""",

    12: """⚙️ CI/CD SÉCURITÉ — Questions possibles :

❓ "Comment la sécurité est intégrée dans votre pipeline ?"
✅ 6 stages : install → lint → test → build → deploy-staging → deploy-production
   Tests unitaires crypto (encryption.test.js) au stage TEST — bloquants.
   Secrets dans variables GitLab CI — jamais dans le code.

❓ "Les secrets sont-ils dans le code ?"
✅ NON. Variables GitLab CI : $FIREBASE_TOKEN · $JWT_SECRET.
   En local : .env dans le .gitignore.
   Amélioration : Google Secret Manager avec rotation automatique.
""",

    # ── Slides 13-33 : notes générales défense ────────────────────────────────
}

# Notes générales pour tous les slides sans note spécifique
GENERAL_NOTE = """🎤 NOTES PRÉSENTATEUR — SOUTENANCE PFE

📌 CHIFFRES CLÉS À RETENIR :
• 436 420 événements Wazuh monitorés
• 12/12 tests DAST passés (0 vulnérabilité critique)
• 7/7 démonstrations réussies
• 5/6 articles RGPD satisfaits
• 6 rôles RBAC · 9 types d'auditLogs · 7 règles Wazuh custom
• Access token : 30 min · Refresh token : 7 jours
• bcrypt saltRounds=10 = 1024 itérations
• WAF bloque : SQLi · XSS · Path Traversal · CMD Injection · 17 agents suspects

📌 PHRASES CLÉS À DIRE :
• "Défense en profondeur sur 2 niveaux : applicatif + infrastructure"
• "Les auditLogs sont immuables — allow update: if false au niveau Firestore"
• "Wazuh surveille l'infrastructure là où mon code ne peut pas voir"
• "Le mot de passe n'est jamais stocké en clair — bcrypt hash irréversible"
• "Le WAF intercepte avant tout traitement métier"

📌 SI LE JURY DEMANDE UNE DÉMONSTRATION :
Backend  : http://localhost:5001/gestionadminastration/us-central1/api/
Frontend : http://localhost:8081
Wazuh    : https://localhost  (admin / SecretPassword)

📌 RÉPONSE SI QUESTION DIFFICILE :
"C'est une excellente question — c'est précisément ce que j'ai identifié
dans mon plan d'amélioration. La solution serait [MFA / Redis / Argon2id / DPIA]."
"""


def set_notes(slide, text):
    """Ajoute ou remplace les notes d'un slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame

    # Vider tous les paragraphes existants sauf le premier
    txBody = tf._txBody
    paras = txBody.findall(qn('a:p'))
    for p in paras[1:]:
        txBody.remove(p)

    # Vider le premier paragraphe
    first_para = paras[0] if paras else etree.SubElement(txBody, qn('a:p'))
    for r in first_para.findall(qn('a:r')):
        first_para.remove(r)

    lines = text.strip().split('\n')

    # Écrire toutes les lignes
    for i, line in enumerate(lines):
        if i == 0:
            p_elem = first_para
        else:
            p_elem = etree.SubElement(txBody, qn('a:p'))
        r_elem = etree.SubElement(p_elem, qn('a:r'))
        rPr = etree.SubElement(r_elem, qn('a:rPr'), lang='fr-FR')
        rPr.set('sz', '1000')  # 10pt
        t_elem = etree.SubElement(r_elem, qn('a:t'))
        t_elem.text = line


def main():
    prs = Presentation(PPTX_FILE)
    total = len(prs.slides)
    print(f"📂 Présentation chargée : {total} slides")

    notes_added = 0
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        # Note spécifique pour ce slide
        if slide_num in NOTES:
            note_text = NOTES[slide_num]
        else:
            # Note générale pour les autres slides
            note_text = GENERAL_NOTE

        set_notes(slide, note_text)
        notes_added += 1

    prs.save(PPTX_FILE)
    print(f"✅ Notes ajoutées sur {notes_added} slides")
    print(f"✅ Sauvegardé : {PPTX_FILE}")
    print()
    print("📖 COMMENT VOIR LES NOTES :")
    print("   PowerPoint → Affichage → Mode Présentateur")
    print("   ou : Affichage → Notes")
    print("   Les notes sont VISIBLES SEULEMENT par le présentateur")
    print("   Le jury voit uniquement les slides projetés")


if __name__ == "__main__":
    main()
