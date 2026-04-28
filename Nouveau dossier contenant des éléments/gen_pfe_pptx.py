#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PFE Cybersécurité/Cyberdéfense — YNOV Campus 2026
Amine BAHOU / Anass Akker
30 slides — Noir / Bleu Marine / Blanc UNIQUEMENT
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette STRICTE : noir · bleu marine · blanc ──────────────────────────
BG1      = RGBColor(0x06, 0x0e, 0x1a)   # noir profond
BG2      = RGBColor(0x0a, 0x16, 0x28)   # marine très sombre
NAVY     = RGBColor(0x0d, 0x1f, 0x3c)   # bleu marine
NAVY_M   = RGBColor(0x14, 0x28, 0x50)   # marine moyen
NAVY_L   = RGBColor(0x1a, 0x35, 0x66)   # marine clair
NAVY_A   = RGBColor(0x1e, 0x40, 0x80)   # accent marine
WHITE    = RGBColor(0xff, 0xff, 0xff)
W80      = RGBColor(0xcc, 0xd8, 0xe8)   # blanc atténué
GREY     = RGBColor(0x5a, 0x72, 0x90)   # gris-bleu
BORD     = RGBColor(0x2a, 0x45, 0x70)   # bordure

BASE = os.path.dirname(os.path.abspath(__file__))
WAZUCAPT = os.path.join(BASE, 'WAZUCAPT')
CAPMON   = os.path.join(BASE, 'CAPMONITORINGSECU')

def img(folder, name):
    p = os.path.join(folder, name)
    return p if os.path.exists(p) else None

IMGS = {
    'agents':    img(WAZUCAPT, 'wazuh_agents.png'),
    'agents01':  img(WAZUCAPT, 'wazuh_01_agents.png'),
    'overview':  img(WAZUCAPT, 'wazuh_02_overview.png'),
    'fim':       img(WAZUCAPT, 'wazuh_03_fim.png'),
    'cve':       img(WAZUCAPT, 'wazuh_04_cve.png'),
    'mitre':     img(WAZUCAPT, 'wazuh_05_mitre.png'),
    'siem':      img(CAPMON,   'brave_screenshot_localhost.png'),
    'rgpd':      img(CAPMON,   'brave_screenshot_localhost (1).png'),
    'rbac':      img(CAPMON,   'brave_screenshot_localhost (2).png'),
    'waf':       img(CAPMON,   'brave_screenshot_localhost (3).png'),
    'score':     img(CAPMON,   'brave_screenshot_localhost (4).png'),
}

W = Inches(13.33); H = Inches(7.5)
prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]
TOTAL = 30

# ═══════════════════ HELPERS ═══════════════════════════════════════════════

def ns(): return prs.slides.add_slide(BLANK)

def bg(s, c=BG2):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def r(s, x, y, w, h, fc=NAVY_M, bc=None, lw=0.4):
    from pptx.util import Pt
    sh = s.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if bc: sh.line.color.rgb = bc; sh.line.width = Pt(lw)
    else:  sh.line.fill.background()
    return sh

def tx(s, text, x, y, w, h, sz=10, bold=False, color=W80,
       align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    ru = p.add_run(); ru.text = text
    ru.font.size = Pt(sz); ru.font.bold = bold
    ru.font.italic = italic; ru.font.color.rgb = color
    ru.font.name = 'Calibri'
    return tb

def hdr(s, n):
    tx(s, f'PFE Cybersécurité · YNOV Campus 2026',
       Inches(0.35), Inches(0.08), Inches(8), Inches(0.28),
       sz=7.5, color=GREY)
    tx(s, f'{n:02d} / {TOTAL}',
       W - Inches(1.6), Inches(0.08), Inches(1.5), Inches(0.28),
       sz=8.5, color=GREY, align=PP_ALIGN.RIGHT)

def lbl(s, t, x=Inches(0.45), y=Inches(0.45)):
    tx(s, t, x, y, Inches(10), Inches(0.22), sz=7.5, bold=True, color=GREY)

def ttl(s, t, x=Inches(0.45), y=Inches(0.7), sz=22):
    tx(s, t, x, y, Inches(12.6), Inches(0.85), sz=sz, bold=True, color=WHITE)

def div(s, x=Inches(0.45), y=Inches(1.58), w=Inches(1.2)):
    sh = s.shapes.add_shape(1, x, y, w, Pt(2.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.fill.background()

def notes(s, t):
    ny = H - Inches(1.52); nh = Inches(1.44)
    sh = s.shapes.add_shape(1, 0, ny, W, nh)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG1
    sh.line.color.rgb = NAVY_A; sh.line.width = Pt(0.5)
    tx(s, 'NOTE', Inches(0.4), ny + Pt(4), Inches(2), Inches(0.2),
       sz=6.5, bold=True, color=GREY)
    tx(s, t, Inches(0.4), ny + Inches(0.28), W - Inches(0.8), nh - Inches(0.35),
       sz=8.5, color=W80)

def kpi_row(s, items, y=Inches(3.8), h=Inches(0.8)):
    n = len(items); cw = (W - Inches(0.9)) / n; x0 = Inches(0.45)
    for i, (v, l) in enumerate(items):
        cx = x0 + i * cw
        r(s, cx, y, cw - Pt(4), h, NAVY_M, BORD)
        tx(s, v, cx, y + Pt(4), cw - Pt(4), Inches(0.42),
           sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(s, l, cx + Inches(0.05), y + Inches(0.45), cw - Inches(0.1), Inches(0.28),
           sz=7, color=GREY, align=PP_ALIGN.CENTER)

def card_grid(s, items, x=Inches(0.45), y=Inches(1.7), cols=2,
              cw=None, ch=Inches(0.95), gap=Inches(0.12)):
    n = len(items)
    rows = (n + cols - 1) // cols
    if cw is None: cw = (W - Inches(0.9) - gap * (cols - 1)) / cols
    for i, (title, body) in enumerate(items):
        col = i % cols; row = i // cols
        cx = x + col * (cw + gap); cy = y + row * (ch + gap)
        r(s, cx, cy, cw, ch, NAVY_M, BORD)
        tx(s, title, cx + Inches(0.15), cy + Pt(4), cw - Inches(0.2), Inches(0.24),
           sz=8.5, bold=True, color=WHITE)
        tx(s, body, cx + Inches(0.15), cy + Inches(0.3), cw - Inches(0.2), ch - Inches(0.36),
           sz=7.5, color=W80)

def tbl_hdr(s, cols, x, y, tw, rh=Inches(0.3)):
    sh = s.shapes.add_shape(1, x, y, tw, rh)
    sh.fill.solid(); sh.fill.fore_color.rgb = NAVY_L; sh.line.fill.background()
    cx = x
    for (t, fr) in cols:
        cw = Inches(fr * tw.inches)
        tx(s, t.upper(), cx + Pt(4), y + Pt(3), cw - Pt(6), rh - Pt(4), sz=7, bold=True, color=GREY)
        cx += cw

def tbl_row(s, vals, fracs, x, y, tw, rh=Inches(0.29), even=False):
    fc = NAVY_M if even else RGBColor(0x0f, 0x22, 0x3a)
    sh = s.shapes.add_shape(1, x, y, tw, rh)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc; sh.line.fill.background()
    cx = x
    for v, fr in zip(vals, fracs):
        cw = Inches(fr * tw.inches)
        tx(s, v, cx + Pt(4), y + Pt(2), cw - Pt(6), rh - Pt(2), sz=8, color=W80)
        cx += cw

def screenshot(s, path, x, y, w, h, cap=''):
    r(s, x - Pt(2), y - Pt(14), w + Pt(4), Pt(13), NAVY_L, None)
    if cap:
        tx(s, cap, x, y - Pt(12), w, Pt(11), sz=6, color=GREY)
    r(s, x, y, w, h, NAVY_M, BORD)
    if path and os.path.exists(path):
        try: s.shapes.add_picture(path, x, y, w, h)
        except: pass
    if cap:
        r(s, x, y + h, w, Pt(12), BG1, None)
        tx(s, cap, x + Inches(0.1), y + h + Pt(1), w - Inches(0.1), Pt(11),
           sz=6, color=GREY, italic=True)

def li_list(s, items, x, y, w, eh=Inches(0.38)):
    cy = y
    for (b, t) in items:
        r(s, x, cy, w, eh - Pt(2), NAVY_M, BORD, 0.3)
        tx(s, f'{b}  {t}' if b else t, x + Inches(0.15), cy + Pt(2),
           w - Inches(0.18), eh - Pt(5), sz=8.5, color=W80)
        cy += eh

def section_slide(n, title, sub):
    s = ns(); bg(s, BG1); hdr(s, n)
    r(s, 0, H * 0.35, W, H * 0.32, NAVY_M, None)
    sh = s.shapes.add_shape(1, Inches(0.45), H * 0.35, Pt(3), H * 0.32)
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE; sh.line.fill.background()
    tx(s, title, Inches(0.7), H * 0.37, W - Inches(0.9), H * 0.2,
       sz=30, bold=True, color=WHITE)
    tx(s, sub, Inches(0.7), H * 0.57, W - Inches(0.9), H * 0.1,
       sz=11, color=W80, italic=True)
    return s

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 01 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG1)
r(s, 0, 0, W, H * 0.06, NAVY_A, None)
r(s, 0, H - H * 0.06, W, H * 0.06, NAVY_A, None)
r(s, Inches(0.45), Inches(0.9), W - Inches(0.9), H * 0.72, NAVY, None)
tx(s, 'PFE CYBERSÉCURITÉ / CYBERDÉFENSE',
   Inches(0.45), Inches(0.48), W - Inches(0.9), Inches(0.3),
   sz=8, bold=True, color=GREY, align=PP_ALIGN.CENTER)
tx(s, 'SÉCURITÉ & MONITORING',
   Inches(0.45), Inches(1.0), W - Inches(0.9), Inches(1.0),
   sz=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, 'Application Web — Gestion Scolaire YNOV Campus',
   Inches(0.45), Inches(2.05), W - Inches(0.9), Inches(0.5),
   sz=17, color=W80, align=PP_ALIGN.CENTER)
tx(s, 'Phase 3 — Post-Production  ·  WAF  ·  SIEM  ·  DAST  ·  Wazuh 4.7.4  ·  CDC §3.3 à 100%',
   Inches(0.45), Inches(2.6), W - Inches(0.9), Inches(0.3),
   sz=10, color=GREY, align=PP_ALIGN.CENTER, italic=True)
badges = ['OWASP Top 10', 'RGPD UE 2016/679', 'CDC §3.3', 'Wazuh 4.7.4']
for i, b in enumerate(badges):
    bx = Inches(2.4 + i * 2.2)
    r(s, bx, Inches(3.1), Inches(2.0), Inches(0.4), NAVY_M, BORD)
    tx(s, b, bx, Inches(3.12), Inches(2.0), Inches(0.35), sz=8.5, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
r(s, Inches(0.45), Inches(3.7), W - Inches(0.9), Pt(1), NAVY_A, None)
kpi_row(s, [('10','Catégories OWASP'),('6','Articles RGPD'),('12','Tests DAST'),
             ('2 500+','FIM Events'),('17','CVE Détectées'),('1 600+','MITRE ATT&CK')],
        y=Inches(3.82), h=Inches(0.72))
tx(s, 'Amine BAHOU  /  Anass Akker  —  YNOV Campus 2026',
   Inches(0.45), H - Inches(0.48), W - Inches(0.9), Inches(0.3),
   sz=9, color=GREY, align=PP_ALIGN.CENTER)
print('01 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 02 — PLAN
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 2); lbl(s, 'NAVIGATION'); ttl(s, 'Plan de la Présentation — 30 Slides'); div(s)
toc = [
    ('01','Cover — Contexte & KPI'),('02','Plan de la présentation'),
    ('03','Contexte du projet — Stack technique'),('04','Audit OWASP Top 10 (2021)'),
    ('05','Vulnérabilités & Corrections appliquées'),('06','RBAC — Gestion des accès par rôle'),
    ('07','Conformité RGPD — 6 articles'),('08','Audit Logs Firestore immuables'),
    ('09','Security Headers · JWT · Chiffrement'),('10','Rate Limiting & Protection DoS'),
    ('11','Architecture sécurité en couches'),('12','Tests DAST — 12 tests OWASP — 92/100'),
    ('13','[ Monitoring Applicatif ]'),('14','Dashboard Accès Sécurisé — Score 100/100'),
    ('15','Dashboard WAF OWASP Top 10'),('16','Dashboard SIEM Logs Temps Réel'),
    ('17','Dashboard RGPD Conformité'),('18','Dashboard RBAC & Audit Logs'),
    ('19','[ Wazuh SIEM Infrastructure ]'),('20','Architecture Wazuh + Installation Docker'),
    ('21','Agents Wazuh — Coverage 100%'),('22','Agent 002 — Vue d\'ensemble'),
    ('23','FIM — 2 500+ Événements'),('24','CVE — 17 Vulnérabilités'),
    ('25','MITRE ATT&CK — 1 600+ Alertes'),('26','5 Modules Wazuh Actifs'),
    ('27','Monitoring à 2 niveaux — Defense in Depth'),('28','Synthèse CDC §3.3 à 100%'),
    ('29','Bilan global & Perspectives'),('30','Merci — Questions'),
]
cols2 = 2; rpp = 15
for i, (n, t) in enumerate(toc):
    col = i // rpp; row = i % rpp
    cx = Inches(0.45 + col * 6.5); cy = Inches(1.72 + row * 0.325)
    r(s, cx, cy, Inches(6.2), Inches(0.3), NAVY_M, BORD, 0.3)
    tx(s, n, cx + Inches(0.1), cy + Pt(2), Inches(0.5), Inches(0.25), sz=7.5, bold=True, color=GREY)
    tx(s, t, cx + Inches(0.65), cy + Pt(2), Inches(5.4), Inches(0.25), sz=8.5, color=W80)
notes(s, "Plan structuré en 3 grandes parties : (1) Sécurité Applicative — OWASP, RBAC, RGPD, Audit Logs, JWT, Chiffrement, DAST ; (2) Monitoring Applicatif — Dashboard /monitoring avec 5 captures ; (3) Infrastructure Wazuh SIEM — Architecture, agents, FIM, CVE, MITRE ATT&CK. La progression va de la conception sécurisée vers la surveillance en temps réel.")
print('02 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 03 — CONTEXTE DU PROJET
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 3); lbl(s, 'CONTEXTE'); ttl(s, 'Contexte du Projet — Application YNOV Campus'); div(s)
stack = [
    ('Frontend','React 18 + TypeScript + Tailwind CSS'),
    ('Backend','Node.js + Express (Firebase Cloud Functions)'),
    ('Base de données','Firestore NoSQL temps réel'),
    ('Authentification','JWT HS256 + Firebase Auth'),
    ('Hébergement','Firebase Hosting + Azure Container Apps'),
    ('SIEM','Wazuh 4.7.4 — Docker (3 conteneurs)'),
]
card_grid(s, stack, x=Inches(0.45), y=Inches(1.72), cols=3,
          cw=Inches(4.1), ch=Inches(0.85), gap=Inches(0.1))
r(s, Inches(0.45), Inches(3.68), W - Inches(0.9), Inches(0.5), NAVY_M, BORD)
tx(s, 'Rôles :   admin   ·   sous-admin   ·   comptable   ·   étudiant / parent',
   Inches(0.6), Inches(3.72), W - Inches(1.1), Inches(0.4), sz=11, bold=True, color=WHITE)
r(s, Inches(0.45), Inches(4.28), W - Inches(0.9), Inches(0.42), NAVY_M, BORD)
tx(s, 'Données sensibles : élèves · paiements · factures · parents → RGPD obligatoire',
   Inches(0.6), Inches(4.3), W - Inches(1.1), Inches(0.35), sz=9.5, color=W80, italic=True)
notes(s, "L'application GestionScolaire traite des données personnelles sensibles d'élèves : noms, paiements, factures, informations des parents. Ce contexte rend le RGPD (UE 2016/679) obligatoire et non optionnel. La stack React 18 + Node.js + Firestore + Firebase est moderne et performante, mais expose des surfaces d'attaque sur chaque couche (frontend XSS, API injections, BDD règles Firestore, auth JWT). Le SIEM Wazuh 4.7.4 déployé en Docker 3 conteneurs surveille l'infrastructure complète. Les 4 rôles RBAC garantissent le principe du moindre privilège.")
print('03 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 04 — OWASP TOP 10
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 4); lbl(s, 'AUDIT SÉCURITÉ'); ttl(s, 'Audit OWASP Top 10 (2021) — 10 Catégories Vérifiées'); div(s)
owasp = [
    ('A01 — Broken Access Control','RBAC · checkRole() · Firestore Rules','CORRIGÉ'),
    ('A02 — Cryptographic Failures','AES-256-CBC · HTTPS · bcrypt 12 rounds','CORRIGÉ'),
    ('A03 — Injection SQLi·XSS·Cmd','WAF waf.js — détection + blocage → 403','CORRIGÉ'),
    ('A04 — Insecure Design','Architecture JWT + RBAC dès la conception','CORRIGÉ'),
    ('A05 — Security Misconfiguration','Helmet · CORS strict · Headers HTTP sécurisés','CORRIGÉ'),
    ('A06 — Vulnerable Components','Wazuh Vulnerability Detection — 17 CVE','SURVEILLÉ'),
    ('A07 — Auth & Session Failures','Rate limiting · lockout 5 échecs · JWT 24h','CORRIGÉ'),
    ('A08 — Data Integrity Failures','auditLogs immuables · JWT signé HS256','CORRIGÉ'),
    ('A09 — Security Logging Failures','9 types audit logs Firestore · SIEM','CORRIGÉ'),
    ('A10 — SSRF','Non applicable — pas d\'appels serveur↔serveur','N/A'),
]
tw = Inches(12.4); tx0 = Inches(0.45); ty0 = Inches(1.75)
cols4 = [('Catégorie OWASP',3.5),('Implémentation',5.0),('Statut',1.3)]
fracs4 = [c[1]/9.8 for c in cols4]
tbl_hdr(s, cols4, tx0, ty0, tw)
for i, (cat, impl, st) in enumerate(owasp):
    tbl_row(s, [cat, impl, st], fracs4, tx0, ty0 + Inches(0.3)*(i+1), tw, even=(i%2==0))
notes(s, "Sur les 10 catégories OWASP Top 10 (2021), 8 sont intégralement corrigées, 1 est sous surveillance continue (A06 Vulnerable Components via Wazuh CVE detection — 17 CVE identifiées), et 1 est non applicable (A10 SSRF — l'architecture Firebase ne comporte pas d'appels serveur-à-serveur internes). Toutes les catégories CRITIQUES (A01 à A05) sont couvertes par des implémentations techniques vérifiables. A09 (Security Logging) est couvert par le dashboard SIEM avec 9 types d'auditLogs immuables.")
print('04 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 05 — VULNÉRABILITÉS & CORRECTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 5); lbl(s, 'SÉCURITÉ APPLICATIVE'); ttl(s, 'Vulnérabilités Identifiées & Corrections Appliquées'); div(s)
vulns = [
    ('Injection SQL','CRITIQUE','WAF waf.js — SELECT/UNION/DROP/OR 1=1 → HTTP 403 + WAF_BLOCK loggué'),
    ('XSS Cross-Site Scripting','HAUTE','WAF — <script> · eval() · onerror= → HTTP 403 · CSP default-src \'self\''),
    ('Brute force login','HAUTE','Rate limit 10 req/15min + lockout après 5 échecs → AUTH_LOCKOUT Firestore'),
    ('Escalade de privilèges','HAUTE','checkRole() re-lit Firestore à chaque requête — jamais depuis JWT seul'),
    ('Tokens JWT sans expiration','MOYENNE','exp: 24h + SESSION_EXPIRED loggué · secret en variable d\'environnement'),
    ('Données sensibles en clair','HAUTE','AES-256-CBC + IV aléatoire 128 bits · clé en variable d\'environnement'),
]
tw = Inches(12.4); tx0 = Inches(0.45)
cols5 = [('Vulnérabilité',2.5),('Sévérité',1.2),('Correction appliquée',8.7)]
fracs5 = [c[1]/12.4 for c in cols5]
tbl_hdr(s, cols5, tx0, Inches(1.75), tw)
for i, (v, sev, fix) in enumerate(vulns):
    tbl_row(s, [v, sev, fix], fracs5, tx0, Inches(2.05) + Inches(0.34)*i, tw, Inches(0.31), even=(i%2==0))
r(s, tx0, Inches(4.25), tw, Inches(0.35), NAVY_M, BORD)
tx(s, 'Toutes les vulnérabilités CRITIQUES et HAUTES ont été corrigées. Le WAF est la première ligne de défense — appliqué avant même l\'authentification.',
   tx0 + Inches(0.1), Inches(4.28), tw - Inches(0.2), Inches(0.28), sz=8.5, color=W80, italic=True)
notes(s, "L'audit AppSec a identifié 6 vulnérabilités principales. Le WAF waf.js (middleware Express.js) constitue la première couche de défense contre SQLi, XSS et les injections. La re-lecture du rôle Firestore à chaque requête est une décision architecturale clé : même si un token JWT est volé, le rôle peut être révoqué immédiatement sans attendre l'expiration du token. AES-256-CBC avec IV aléatoire garantit que le même texte chiffré deux fois donne des résultats différents, résistant à l'analyse de patterns.")
print('05 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 06 — RBAC
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 6); lbl(s, 'CDC §3.3 — RBAC'); ttl(s, 'RBAC — Gestion des Accès par Rôle — Moindre Privilège'); div(s)
roles = [
    ('admin','Accès complet','monitoring · users · toutes données · configuration · exports'),
    ('sous-admin','Gestion opérationnelle','étudiants · classes · paiements · factures · rapports'),
    ('comptable','Finance uniquement','lecture paiements · génération factures · exports uniquement'),
    ('étudiant / parent','Portail personnel','propres données uniquement · lecture seule · pas d\'admin'),
]
cw = Inches(3.0); cy0 = Inches(1.75)
for i, (role, subtitle, details) in enumerate(roles):
    cx = Inches(0.45 + i * 3.12)
    r(s, cx, cy0, cw, Inches(2.6), NAVY_M, BORD)
    r(s, cx, cy0, cw, Inches(0.38), NAVY_L, None)
    tx(s, role, cx, cy0 + Pt(4), cw, Inches(0.32), sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, subtitle, cx + Inches(0.1), cy0 + Inches(0.45), cw - Inches(0.2), Inches(0.28),
       sz=8.5, color=GREY, align=PP_ALIGN.CENTER)
    for j, d in enumerate(details.split(' · ')):
        ty = cy0 + Inches(0.8) + j * Inches(0.38)
        tx(s, f'▸  {d}', cx + Inches(0.15), ty, cw - Inches(0.2), Inches(0.35), sz=8.5, color=W80)
r(s, Inches(0.45), Inches(4.5), W - Inches(0.9), Inches(0.55), NAVY_M, BORD)
tx(s, 'Implémentation :  checkRole([\'admin\',\'sous-admin\']) sur chaque route  ·  Re-lecture rôle Firestore à CHAQUE requête (jamais depuis JWT uniquement)',
   Inches(0.6), Inches(4.55), W - Inches(1.1), Inches(0.42), sz=9, color=W80)
notes(s, "La décision clé du RBAC est que le rôle est re-lu depuis Firestore à CHAQUE requête HTTP, et non extrait uniquement du JWT. Cela signifie que si un administrateur révoque un rôle, l'effet est immédiat — même si le token JWT de l'utilisateur n'a pas encore expiré. Le principe du moindre privilège est strictement appliqué : un comptable ne peut que lire les paiements et générer des factures, il ne peut pas accéder aux données personnelles ou à la configuration système. Un étudiant/parent ne voit QUE ses propres données.")
print('06 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 07 — RGPD
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 7); lbl(s, 'CDC §3.3 — CONFORMITÉ'); ttl(s, 'Conformité RGPD (UE 2016/679) — 6 Articles Implémentés'); div(s)
rgpd = [
    ('Art. 15 — Droit d\'accès','Export données personnelles\n→ DATA_EXPORT loggué dans auditLogs\nGET /users/:id/export'),
    ('Art. 16 — Rectification','Modification données\n→ traçabilité complète dans Firestore\nCRUD loggué avec serverTimestamp()'),
    ('Art. 17 — Droit à l\'effacement','Anonymisation DATA_ANONYMIZE\n→ champs remplacés · log conservé\nDELETE /users/:id/data'),
    ('Art. 25 — Privacy by Design','AES-256-CBC · HTTPS · RBAC\n→ principe du moindre privilège\ndès la conception architecture'),
    ('Art. 32 — Sécurité du traitement','WAF + JWT HS256 + bcrypt 12 rounds\n→ rate limiting + lockout\nHTTPS/HSTS max-age=31536000'),
    ('Art. 33 — Notification violation','Alerte CRITIQUE auditLogs\n→ dashboard /monitoring\nProcédure 72h documentée'),
]
cw2 = Inches(4.0); ch2 = Inches(1.35); gap2 = Inches(0.1)
for i, (art, desc) in enumerate(rgpd):
    col = i % 3; row = i // 3
    cx = Inches(0.45 + col * (cw2 + gap2))
    cy = Inches(1.75 + row * (ch2 + gap2))
    r(s, cx, cy, cw2, ch2, NAVY_M, BORD)
    r(s, cx, cy, cw2, Inches(0.32), NAVY_L, None)
    tx(s, art, cx, cy + Pt(3), cw2, Inches(0.28), sz=8.5, bold=True, color=WHITE)
    tx(s, desc, cx + Inches(0.12), cy + Inches(0.36), cw2 - Inches(0.2), ch2 - Inches(0.42), sz=7.5, color=W80)
notes(s, "L'Article 33 (notification de violation) est particulièrement important : en cas d'incident de sécurité, l'organisation a 72 heures pour notifier l'autorité de contrôle (CNIL). Le dashboard /monitoring avec ses alertes automatiques (brute force, lockouts, WAF attacks) permet cette détection rapide. L'Article 25 (Privacy by Design) est respecté car l'architecture sécurisée a été conçue dès le départ, et non ajoutée après coup. Le module GDPR de Wazuh surveille la conformité au niveau infrastructure et complète les auditLogs Firestore.")
print('07 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 08 — AUDIT LOGS
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 8); lbl(s, 'CDC §3.3 — JOURNALISATION'); ttl(s, 'Audit Logs Firestore — 9 Types d\'Événements Immuables'); div(s)
events = [
    ('AUTH_SUCCESS','Connexion réussie · INFO','userId + IP + horodatage'),
    ('AUTH_FAILURE','Mauvais mot de passe · WARNING','userId + IP + tentative n°'),
    ('AUTH_LOCKOUT','Blocage après 5 échecs · CRITIQUE','userId + IP + durée blocage'),
    ('LOGOUT','Déconnexion explicite · INFO','userId + durée session'),
    ('SESSION_EXPIRED','Timeout inactivité 30 min · INFO','userId + last_activity'),
    ('ACCESS_DENIED','Rôle insuffisant · WARNING','userId + route + rôle tenté'),
    ('DATA_EXPORT','Export RGPD Art.15 · INFO','userId + données exportées'),
    ('DATA_ANONYMIZE','Anonymisation Art.17 · WARNING','userId + champs anonymisés'),
    ('WAF_BLOCK','Attaque bloquée · CRITIQUE','IP + type attaque + payload'),
]
tw = Inches(12.4); tx0 = Inches(0.45)
cols8 = [('Événement',2.2),('Niveau',2.5),('Données enregistrées',7.7)]
fracs8 = [c[1]/12.4 for c in cols8]
tbl_hdr(s, cols8, tx0, Inches(1.75), tw)
for i, (ev, lvl, data) in enumerate(events):
    tbl_row(s, [ev, lvl, data], fracs8, tx0, Inches(2.05) + Inches(0.31)*i, tw, Inches(0.28), even=(i%2==0))
r(s, tx0, Inches(5.0), tw, Inches(0.35), NAVY_M, BORD)
tx(s, 'Firestore Rules :  allow update: if false  ·  allow delete: if false  →  logs IMMUABLES  ·  serverTimestamp() non manipulable côté client',
   tx0 + Inches(0.1), Inches(5.04), tw - Inches(0.2), Inches(0.28), sz=8.5, color=W80, italic=True)
notes(s, "L'immuabilité des auditLogs est garantie à deux niveaux : (1) côté code backend, les logs ne sont créés qu'en append-only (jamais mis à jour ni supprimés) ; (2) côté Firestore Security Rules, les règles allow update: if false et allow delete: if false garantissent que même si le code backend est compromis, les logs ne peuvent PAS être modifiés. serverTimestamp() côté serveur empêche la manipulation des horodatages côté client. Ces 9 types d'événements couvrent 100% des actions sensibles pour la conformité RGPD et les investigations forensiques.")
print('08 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 09 — SECURITY HEADERS + JWT + CHIFFREMENT
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 9); lbl(s, 'SÉCURITÉ TECHNIQUE'); ttl(s, 'Security Headers · JWT HS256 · Chiffrement AES-256-CBC'); div(s)
# 3 blocs côte à côte
blocs = [
    ('Security Headers (Helmet.js)',
     'Strict-Transport-Security : max-age=31536000 (Force HTTPS 1 an)\nX-Frame-Options : DENY (Anti-clickjacking)\nX-Content-Type-Options : nosniff (Anti-MIME)\nContent-Security-Policy : default-src \'self\'\nX-XSS-Protection : 1; mode=block'),
    ('JWT HS256 — Authentification',
     'Secret : variable d\'environnement — jamais en dur\nExpiration : 24h — SESSION_EXPIRED loggué\nPayload : userId + role + email (jamais mdp)\nVérification : à chaque requête via checkRole()\nLockout : 5 échecs → AUTH_LOCKOUT + blocage 5 min'),
    ('AES-256-CBC · bcrypt · Secrets',
     'AES-256-CBC : clé 256 bits env · IV aléatoire 128 bits\nbcrypt : 12 rounds · comparaison temps constant\nAnti-timing attack · anti-rainbow table\nJWT_SECRET · ENCRYPTION_KEY : .env jamais en code\n.gitignore · injection sécurisée Firebase config'),
]
for i, (title, body) in enumerate(blocs):
    cx = Inches(0.45 + i * 4.3); bw = Inches(4.15); bh = Inches(2.65)
    r(s, cx, Inches(1.75), bw, bh, NAVY_M, BORD)
    r(s, cx, Inches(1.75), bw, Inches(0.38), NAVY_L, None)
    tx(s, title, cx, Inches(1.77), bw, Inches(0.32), sz=9, bold=True, color=WHITE)
    tx(s, body, cx + Inches(0.12), Inches(2.18), bw - Inches(0.2), bh - Inches(0.52), sz=8, color=W80)
notes(s, "La sécurité est renforcée à chaque couche : Helmet.js applique automatiquement 5 headers de sécurité HTTP. JWT HS256 avec expiration 24h et vérification Firestore garantit que les sessions ne peuvent pas être prolongées indéfiniment. AES-256-CBC avec IV aléatoire à chaque chiffrement garantit que deux chiffrements du même texte donnent des résultats différents, résistant à l'analyse de fréquence. bcrypt avec 12 rounds est résistant aux attaques GPU modernes. La règle absolue : aucun secret ne doit jamais apparaître dans le code source (ni hardcodé ni versionné dans Git).")
print('09 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RATE LIMITING & DoS
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 10); lbl(s, 'PROTECTION DoS'); ttl(s, 'Rate Limiting & Protection DoS — express-rate-limit'); div(s)
rl_items = [
    ('Global — 100 req / 15 min / IP','HTTP 429','Toutes les routes API · protection générale'),
    ('Auth login — 10 req / 15 min / IP','HTTP 429','Route /auth/login uniquement · anti-brute force'),
    ('Brute force — 5 échecs → lockout 5 min','AUTH_LOCKOUT','Tentatives connexion échouées · loggué Firestore'),
    ('Payload max — 10 KB maximum','HTTP 413','Protection DoS upload / injection massive'),
    ('Timeout — 28 secondes maximum','HTTP 503','Requêtes longues / Firestore timeout'),
]
tw = Inches(12.4); tx0 = Inches(0.45)
cols10 = [('Règle',3.6),('Code HTTP',1.4),('Périmètre & Action',7.4)]
fracs10 = [c[1]/12.4 for c in cols10]
tbl_hdr(s, cols10, tx0, Inches(1.75), tw)
for i, (rule, code, scope) in enumerate(rl_items):
    tbl_row(s, [rule, code, scope], fracs10, tx0, Inches(2.05) + Inches(0.42)*i, tw, Inches(0.38), even=(i%2==0))
r(s, tx0, Inches(4.3), tw, Inches(0.45), NAVY_M, BORD)
tx(s, 'Rate limiting appliqué par IP. En production, Redis partagera le compteur entre toutes les instances Cloud Functions pour une protection globale cohérente.',
   tx0 + Inches(0.1), Inches(4.33), tw - Inches(0.2), Inches(0.38), sz=9, color=W80, italic=True)
notes(s, "Le rate limiting multi-couches protège contre différents types d'attaques : le limite global (100 req/15min) protège contre les DoS volumétriques, la limite /auth/login (10 req/15min) protège contre le brute force d'authentification, et le lockout après 5 échecs (AUTH_LOCKOUT loggué) ajoute une couche supplémentaire. La limite payload 10KB empêche les injections massives via le corps des requêtes. En production avec plusieurs instances Cloud Functions, Redis sera nécessaire pour partager les compteurs entre instances.")
print('10 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ARCHITECTURE COUCHES
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 11); lbl(s, 'ARCHITECTURE SÉCURITÉ'); ttl(s, 'Pipeline de Sécurité — De la Requête à l\'Alerte SIEM'); div(s)
pipeline = [
    ('WAF (waf.js)','SQLi · XSS · PathTraversal · CmdInj · Agents suspects  →  HTTP 403 + WAF_BLOCK'),
    ('Rate Limiter','100 req/15min global · 5 échecs → AUTH_LOCKOUT  →  HTTP 429'),
    ('verifyJWT','JWT HS256 · expiration 24h · signature  →  HTTP 401 si invalide'),
    ('checkRole() Firestore','Re-lecture rôle BDD à chaque requête  →  HTTP 403 + ACCESS_DENIED si insuffisant'),
    ('Route Handler','Traitement métier sécurisé  →  Réponse applicative'),
    ('auditLogger','9 types events · Firestore immuable · serverTimestamp()'),
    ('Dashboard /monitoring','Score /100 · WAF stats · SIEM logs · 4 alertes auto · 60s refresh'),
    ('Wazuh SIEM','CIS Benchmark · MITRE ATT&CK · FIM · CVE · GDPR · Agent actif'),
]
for i, (name, desc) in enumerate(pipeline):
    cy = Inches(1.72 + i * 0.52)
    r(s, Inches(0.45), cy, W - Inches(0.9), Inches(0.48), NAVY_M, BORD, 0.3)
    r(s, Inches(0.45), cy, Inches(2.8), Inches(0.48), NAVY_L, None)
    tx(s, name, Inches(0.55), cy + Pt(3), Inches(2.65), Inches(0.4), sz=9, bold=True, color=WHITE)
    tx(s, desc, Inches(3.4), cy + Pt(3), W - Inches(3.75), Inches(0.4), sz=8.5, color=W80)
    if i < len(pipeline) - 1:
        r(s, W/2 - Pt(3), cy + Inches(0.48), Pt(6), Pt(5), WHITE, None)
notes(s, "Chaque requête HTTP traverse obligatoirement 4 couches de sécurité avant d'atteindre le handler métier : WAF (attaques OWASP) → Rate Limiter (DoS/brute force) → verifyJWT (authentification) → checkRole() Firestore (autorisation). Ce pipeline garantit qu'aucune requête malveillante ne peut atteindre les données. Après traitement, un auditLog est créé, visible dans le dashboard /monitoring (Score /100, auto-refresh 60s) et les événements critiques remontent à Wazuh SIEM pour corrélation MITRE ATT&CK.")
print('11 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — TESTS DAST
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 12); lbl(s, 'TESTS DE SÉCURITÉ AUTOMATISÉS'); ttl(s, 'Scanner DAST — 12 Tests OWASP — Score 92/100'); div(s)
r(s, Inches(0.45), Inches(1.78), W - Inches(0.9), Inches(0.4), NAVY_M, BORD)
tx(s, 'node back/functions/scripts/security_scan.js  http://localhost:5001/.../api   →   Score DAST final : 92/100 — 11/12 tests réussis · 1 warning (rate limit header non exposé)',
   Inches(0.6), Inches(1.82), W - Inches(1.1), Inches(0.32), sz=8.5, color=WHITE)
tests = [
    ('T01','API opérationnelle','HTTP 200 + status:true','✓'),
    ('T02','Headers sécurité (HSTS/CSP)','HSTS · X-Frame · CSP présents','✓'),
    ('T03','Auth requise','HTTP 401 sans token','✓'),
    ('T04','Rate limiting actif','HTTP 429 à la 11e requête','✓'),
    ('T05','WAF — SQL Injection','HTTP 403 bloqué immédiatement','✓'),
    ('T06','WAF — Cross-Site Scripting','HTTP 403 bloqué immédiatement','✓'),
    ('T07','WAF — Path Traversal','HTTP 403 bloqué immédiatement','✓'),
    ('T08','WAF — Agents suspects (sqlmap)','HTTP 403 bloqué','✓'),
    ('T09','Payload > 10 KB','HTTP 413 rejeté','✓'),
    ('T10','Escalade de privilèges (rôle forcé)','Rôle \'etudiant\' — accès bloqué','✓'),
    ('T11','/monitoring admin only','HTTP 401/403 non-admin','✓'),
    ('T12','CORS origines interdites','Header absent evil-attacker.com','~ warning'),
]
tw = Inches(12.4); tx0 = Inches(0.45)
cols12 = [('ID',0.6),('Test',3.5),('Résultat',5.5),('Statut',1.0)]
fracs12 = [c[1]/10.6 for c in cols12]
tbl_hdr(s, cols12, tx0, Inches(2.3), tw)
for i, (tid, name, res, st) in enumerate(tests):
    tbl_row(s, [tid, name, res, st], fracs12, tx0, Inches(2.58) + Inches(0.245)*i, tw, Inches(0.23), even=(i%2==0))
notes(s, "Le scanner DAST security_scan.js est un outil maison développé en Node.js qui exécute 12 tests automatisés ciblant l'OWASP Top 10. Il peut être intégré dans un pipeline CI/CD pour valider la sécurité à chaque déploiement. Le score 92/100 (11/12 tests réussis) reflète un excellent niveau de sécurité. Le seul warning concerne l'exposition du header rate-limit (Retry-After), non critique. Le rapport est généré en JSON + HTML, réutilisable pour les audits.")
print('12 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SECTION BREAK — MONITORING APPLICATIF
# ══════════════════════════════════════════════════════════════════════════════
s = section_slide(13, 'PARTIE 1 — MONITORING APPLICATIF',
                  '5 captures · WAF · auditLogs · SIEM · Score /100 · CDC §3.3')
kpi_row(s, [('10','Étapes tech'),('9','Audit logs'),('12','Tests DAST'),('/100','Score sécu')],
        y=Inches(4.4), h=Inches(0.75))
print('13 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DASHBOARD ACCÈS SÉCURISÉ
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 14); lbl(s, 'DASHBOARD /MONITORING — SECTION 1'); ttl(s, 'Dashboard Accès Sécurisé — Score 100 / 100', sz=20); div(s)
li_list(s, [
    ('Score 100/100 :','Calculé sur 24h — authFailures + lockouts + accessDenied + wafBlocks'),
    ('Protections actives :','JWT HS256 forcé · bcrypt saltRounds=10 · Rate limit 10/15min · Timeout 30min'),
    ('Alertes auto :','>20 auth failures → CRITIQUE · >5 lockouts → CRITIQUE · >10 WAF blocks → CRITIQUE'),
], x=Inches(0.45), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.58))
screenshot(s, IMGS['score'], Inches(5.35), Inches(1.75), Inches(7.53), Inches(3.45),
           'Section 1 — Score sécurité · JWT actif · Rate limiting · Lockout brute force')
notes(s, "Le score de sécurité /100 est calculé dynamiquement en temps réel sur les 24 dernières heures : score = 100 − authFailures(max-20) − lockouts(max-15) − accessDenied(max-15) − wafBlocks(max-20). Vert (80-100) = sécurité normale, Orange (60-79) = alerte, Rouge (0-59) = incident. Le score se rafraîchit automatiquement toutes les 60 secondes. En cas d'attaque démonstrative (5 connexions échouées), le score baisse en temps réel devant le jury — effet visuel immédiat.")
print('14 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DASHBOARD WAF
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 15); lbl(s, 'DASHBOARD /MONITORING — WAF'); ttl(s, 'Dashboard WAF — OWASP Top 10 Actif · Toutes Attaques Bloquées', sz=20); div(s)
li_list(s, [
    ('A03 SQLi :','SELECT/UNION/DROP → HTTP 403 + WAF_BLOCK'),
    ('A03 XSS :','<script>/onerror=/eval() → HTTP 403 + WAF_BLOCK'),
    ('A01 LFI :','../../ / %2e%2e → HTTP 403 + WAF_BLOCK'),
    ('Cmd Inj :','; && | $( ) → HTTP 403 + WAF_BLOCK'),
    ('Scanners :','sqlmap/nikto/nmap → HTTP 403 bloqué'),
], x=Inches(0.45), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.47))
screenshot(s, IMGS['waf'], Inches(5.35), Inches(1.75), Inches(7.53), Inches(3.45),
           'WAF OWASP — SQLi · XSS · Path Traversal · CMD Injection · Agents suspects bloqués')
notes(s, "Le dashboard WAF affiche en temps réel les compteurs de blocages par type d'attaque, les barres de progression, et un tableau des 10 dernières attaques avec IP source, chemin et méthode HTTP. Chaque blocage génère un WAF_BLOCK dans les auditLogs Firestore, visible dans l'onglet SIEM. Les compteurs à zéro indiquent un environnement sain au moment de la démo. En démonstration live : envoyer une requête SQLi depuis curl génère immédiatement un compteur > 0 et un log WAF_BLOCK visible dans le tableau.")
print('15 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DASHBOARD SIEM LOGS
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 16); lbl(s, 'DASHBOARD /MONITORING — SIEM LOGS'); ttl(s, 'SIEM Logs Temps Réel — Audit Trail Complet', sz=20); div(s)
li_list(s, [
    ('Colonnes :','IP · email · rôle · chemin · action · horodatage'),
    ('Source :','Collection Firestore auditLogs — 24 dernières heures'),
    ('0 events :','Surveillance normale — aucune anomalie détectée'),
    ('Admin only :','<ProtectedRoute roles={[\'admin\']}> · redirection /unauthorized'),
], x=Inches(0.45), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.52))
r(s, Inches(0.45), Inches(3.88), Inches(4.75), Inches(0.5), NAVY_M, BORD)
tx(s, 'C\'est un "boîtier noir" infalsifiable : personne ne peut effacer ou modifier ces preuves, même pas un administrateur.',
   Inches(0.6), Inches(3.92), Inches(4.45), Inches(0.4), sz=8, color=W80, italic=True)
screenshot(s, IMGS['siem'], Inches(5.35), Inches(1.75), Inches(7.53), Inches(3.45),
           'SIEM — Journal 20 derniers événements · IP · email · rôle · chemin · action · horodatage')
notes(s, "L'onglet SIEM Logs est le journal de bord complet de l'application. Il agrège les 20 derniers événements de sécurité en temps réel depuis Firestore, avec toutes les métadonnées : adresse IP de l'appelant, email de l'utilisateur, rôle, route appelée, type d'action, et horodatage serveur. Le dashboard est accessible uniquement aux admins (ProtectedRoute). Un non-admin qui tente d'accéder à /monitoring génère automatiquement un ACCESS_DENIED loggué. Cet audit trail complet est la preuve irréfutable pour tout contrôle de conformité RGPD.")
print('16 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DASHBOARD RGPD
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 17); lbl(s, 'DASHBOARD /MONITORING — RGPD'); ttl(s, 'Dashboard Conformité RGPD — Art. 15 · 16 · 17 · 33', sz=20); div(s)
li_list(s, [
    ('Art. 15 :','GET /users/:id/export → DATA_EXPORT loggué'),
    ('Art. 17 :','DELETE /users/:id/data → DATA_ANONYMIZE loggué'),
    ('Chiffrement :','AES-256-CBC actif + HTTPS/HSTS max-age=1an'),
    ('Art. 33 :','Procédure notification 72h documentée'),
], x=Inches(0.45), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.52))
screenshot(s, IMGS['rgpd'], Inches(5.35), Inches(1.75), Inches(7.53), Inches(3.45),
           'Section 2 — RGPD : Exports (Art.15) · Anonymisations (Art.17) · Conformité AES-256 · Art.33')
notes(s, "Ce dashboard RGPD permet à un DPO (Data Protection Officer) de vérifier la conformité sans accès technique aux systèmes backend. Les compteurs Exports (Art.15) et Anonymisations (Art.17) à zéro signifient qu'aucune demande RGPD n'a été traitée dans les 24h. Les coches vertes confirment que tous les mécanismes techniques sont opérationnels : AES-256-CBC, HTTPS/HSTS, Art.15/16/17/33 implémentés. Ce tableau de bord est directement dérivé des auditLogs Firestore, sans traitement intermédiaire.")
print('17 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — DASHBOARD RBAC & AUDIT LOGS
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 18); lbl(s, 'DASHBOARD /MONITORING — RBAC'); ttl(s, 'Dashboard RBAC & Audit Logs — Deny by Default', sz=20); div(s)
li_list(s, [
    ('Deny by default :','Firestore Rules + ProtectedRoute /unauthorized'),
    ('4 rôles :','Admin · Sous-admin · Comptable · Étudiant'),
    ('Intégrité :','allow update: if false · allow delete: if false'),
    ('9 events :','AUTH + SESSION + ACCESS + RGPD + WAF_BLOCK'),
], x=Inches(0.45), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.52))
screenshot(s, IMGS['rbac'], Inches(5.35), Inches(1.75), Inches(7.53), Inches(3.45),
           'Section 3/4 — Accès refusés RBAC · 9 types audit logs immuables · Intégrité logs')
notes(s, "L'architecture RBAC suit le principe 'deny by default' : tout accès est interdit sauf ce qui est explicitement autorisé. La section 3 (RBAC) affiche les tentatives d'accès non autorisées (ACCESS_DENIED) en temps réel. La section 4 (Audit Logs) montre les compteurs d'événements par catégorie : Auth, RBAC, RGPD. L'intégrité des logs est garantie par les règles Firestore (allow update/delete: if false) — même l'admin système ne peut pas effacer les preuves.")
print('18 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — SECTION BREAK — WAZUH SIEM
# ══════════════════════════════════════════════════════════════════════════════
s = section_slide(19, 'PARTIE 2 — WAZUH SIEM INFRASTRUCTURE',
                  'Monitoring Infrastructure · MITRE ATT&CK · GDPR · CIS Benchmark')
kpi_row(s, [('213','Événements < 24h'),('3','Conteneurs Docker'),
             ('5','Modules actifs'),('100%','Agents Coverage')],
        y=Inches(4.4), h=Inches(0.75))
print('19 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — ARCHITECTURE WAZUH + INSTALLATION
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 20); lbl(s, 'WAZUH SIEM'); ttl(s, 'Architecture Wazuh 4.7.4 — Déploiement Docker — 4 Commandes'); div(s)
# Architecture colonne gauche
arch_items = [
    ('Application Web (React)', 'Frontend · auditLogs Firestore · dashboard /monitoring'),
    ('API Backend (Node.js)','WAF · JWT · RBAC · Handler · auditLogger'),
    ('Wazuh Agent', 'OSSEC agent · Port 1514 · Chiffrement AES · Push events → Manager'),
    ('Wazuh Manager','Collecte · corrélation · moteur règles · Port 1514/55000'),
    ('Wazuh Indexer (OpenSearch)','Stockage alertes · Port 9200 · Full-text search'),
    ('Wazuh Dashboard','Interface Kibana fork · Port 443 HTTPS TLS'),
]
for i, (comp, desc) in enumerate(arch_items):
    cy = Inches(1.75 + i * 0.53)
    r(s, Inches(0.45), cy, Inches(6.1), Inches(0.48), NAVY_M, BORD, 0.3)
    tx(s, comp, Inches(0.6), cy + Pt(3), Inches(2.2), Inches(0.4), sz=8.5, bold=True, color=WHITE)
    tx(s, desc, Inches(2.85), cy + Pt(3), Inches(3.6), Inches(0.4), sz=8, color=W80)
    if i < len(arch_items) - 1 and i != 1:
        r(s, Inches(1.2), cy + Inches(0.48), Pt(6), Pt(4), WHITE, None)
# Installation colonne droite
r(s, Inches(6.8), Inches(1.75), Inches(6.1), Inches(3.25), NAVY_M, BORD)
r(s, Inches(6.8), Inches(1.75), Inches(6.1), Inches(0.3), NAVY_L, None)
tx(s, 'Installation — 4 Commandes', Inches(6.8), Inches(1.77), Inches(6.0), Inches(0.26), sz=9, bold=True, color=WHITE)
cmds = [
    ('# 1. Cloner Wazuh Docker v4.7.4','git clone -b v4.7.4 github.com/wazuh/wazuh-docker.git'),
    ('# 2. Single-node','cd wazuh-docker/single-node'),
    ('# 3. Certificats TLS','docker compose -f generate-indexer-certs.yml run --rm generator'),
    ('# 4. Démarrer','docker compose up -d'),
    ('# Accès dashboard','https://localhost  →  admin / SecretPassword'),
]
for i, (comment, cmd) in enumerate(cmds):
    cy = Inches(2.15) + i * Inches(0.46)
    tx(s, comment, Inches(6.95), cy, Inches(5.8), Inches(0.22), sz=7.5, color=GREY, italic=True)
    r(s, Inches(6.95), cy + Inches(0.22), Inches(5.8), Inches(0.22), BG1, None)
    tx(s, cmd, Inches(7.05), cy + Inches(0.22), Inches(5.7), Inches(0.2), sz=8, color=WHITE, bold=True)
notes(s, "Wazuh se déploie en 3 conteneurs Docker : wazuh-manager (collecte et corrélation), wazuh-indexer (OpenSearch pour le stockage et la recherche full-text), wazuh-dashboard (interface Kibana fork, HTTPS). L'installation complète prend ~10 minutes avec 4 commandes. Les certificats TLS sont générés automatiquement. Le manager écoute sur le port 1514 pour les agents et 55000 pour l'API REST. Toutes les communications agent ↔ manager sont chiffrées AES. Le dashboard est accessible sur https://localhost avec authentification.")
print('20 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — AGENTS WAZUH — COVERAGE 100%
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 21); lbl(s, 'WAZUH SIEM — AGENTS'); ttl(s, 'Agents Wazuh — 2 Agents Actifs — Coverage 100%'); div(s)
screenshot(s, IMGS['agents'], Inches(0.45), Inches(1.75), Inches(7.6), Inches(3.4),
           'Agents list — Agent 001: frais-gestion-scolaire (Debian 12) + Agent 002: main-machine (macOS 15.7.4) · Actifs')
li_list(s, [
    ('Agent 001 :','frais-gestion-scolaire · Debian 12 · 172.20.0.5 · v4.7.4 · Active'),
    ('Agent 002 :','main-machine · macOS 15.7.4 · 127.0.0.1 · v4.7.4 · Active'),
    ('Coverage :','100% — 2 actifs · 0 déconnecté · 0 pending'),
    ('Protocole :','OSSEC · Port 1514 · Chiffrement AES · Temps réel'),
], x=Inches(8.2), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.52))
notes(s, "Les 2 agents Wazuh couvrent l'intégralité de l'infrastructure : Agent 001 surveille le conteneur Docker hébergeant l'application GestionScolaire (Debian 12, réseau Docker 172.20.0.x), Agent 002 surveille la machine de développement macOS 15.7.4. Le coverage 100% signifie qu'aucune partie de l'infrastructure n'échappe à la surveillance SIEM. L'authentification agent ↔ manager utilise une clé pré-partagée générée à l'enrôlement, et toutes les communications sont chiffrées AES. Les événements remontent en temps réel pour corrélation.")
print('21 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — AGENT 002 VUE D'ENSEMBLE
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 22); lbl(s, 'WAZUH — AGENT 002'); ttl(s, 'Agent 002 — Dashboard macOS 15.7.4 — Vue d\'Ensemble Complète', sz=20); div(s)
screenshot(s, IMGS['overview'], Inches(0.45), Inches(1.75), Inches(7.6), Inches(3.4),
           'Agent 002 — Security events · Integrity monitoring · SCA · Vulnerabilities · MITRE ATT&CK · More')
li_list(s, [
    ('ID 002 :','macOS 15.7.4 · Active · Wazuh v4.7.4'),
    ('MITRE :','Impact ×1594 · Defense Evasion ×6'),
    ('FIM :','2 500+ events · Rule 550 Niveau 7'),
    ('SCA :','0 Passed / 10 · unix_audit — 0%'),
], x=Inches(8.2), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.52))
notes(s, "La vue d'ensemble de l'Agent 002 montre simultanément tous les modules actifs et leurs métriques : MITRE ATT&CK (1594 Impact + 6 Defense Evasion), FIM (2500+ events Rule 550), SCA Compliance (0/10), CVE Vulnerabilities. C'est cette vue qui permet à un opérateur SOC de constater d'un coup d'œil la santé globale de l'agent. Le SCA à 0% est un indicateur important : l'audit de configuration système unix_audit n'a passé aucun des 10 checks, révélant une surface d'attaque maximale au niveau configuration système.")
print('22 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — FIM
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 23); lbl(s, 'WAZUH — FILE INTEGRITY MONITORING'); ttl(s, 'FIM — 2 500+ Événements · Rule 550 Level 7 · root 99.44%', sz=20); div(s)
screenshot(s, IMGS['fim'], Inches(0.45), Inches(1.75), Inches(7.6), Inches(3.4),
           'FIM Dashboard — Most active users: root 99.44% · Actions: modified 100% · Events: 2500+ · /bin/bash /bin/cp /bin/df')
li_list(s, [
    ('Règle :','Rule 550 Niveau 7 — Integrity checksum changed'),
    ('root 99.44% :','Mises à jour système macOS — binaires /bin'),
    ('Fichiers :','/bin/bash · /bin/cp · /bin/df · /bin/echo · /usr/bin/*'),
    ('Risque :','Rootkit/backdoor si modifications non légitimes'),
], x=Inches(8.2), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.52))
notes(s, "Le FIM (File Integrity Monitoring) surveille les modifications des fichiers système critiques via hash SHA-256. Les 2500+ événements correspondent aux modifications des binaires macOS /bin/* lors des mises à jour système Sequoia 15.7.4 — c'est un comportement normal macOS. root à 99.44% est cohérent car les mises à jour OS s'effectuent avec les droits root. RGPD Art.5 : l'intégrité des données est maintenue. Ce module est équivalent à un IDS (Intrusion Detection System) au niveau filesystem — toute modification suspecte (rootkit, backdoor) déclencherait immédiatement une alerte.")
print('23 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — CVE
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 24); lbl(s, 'WAZUH — VULNERABILITY DETECTION'); ttl(s, 'CVE — 17 Vulnérabilités Identifiées — 8 High + 9 Medium', sz=20); div(s)
screenshot(s, IMGS['cve'], Inches(0.45), Inches(1.75), Inches(7.6), Inches(3.4),
           'CVE Dashboard — Severity: 0 Critical · 8 High · 9 Medium · Docker (8) · Python (4) · Excel (1) · lz4 (1)')
li_list(s, [
    ('CVE-2019-5736 :','CVSS3=8.6 · Docker container escape — runc'),
    ('8 High :','Docker 4.43.2 · Excel 16.107.3'),
    ('9 Medium :','Docker CVE-2020/2021 · lz4 · Python'),
    ('Action :','Mise à jour Docker Engine → priorité 1'),
], x=Inches(8.2), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.52))
notes(s, "Le Vulnerability Detector compare les packages installés sur chaque agent avec la base NVD NIST pour identifier les CVE connues. CVE-2019-5736 (CVSS3=8.6) est la plus critique : elle permet à un attaquant depuis l'intérieur d'un conteneur Docker d'écrire dans le binaire runc de l'hôte, permettant une évasion de conteneur et élévation de privilèges vers root. Patch prioritaire : mise à jour Docker Engine ≥ 26.1.5. Les 9 CVE Medium (Docker 2020-2021, lz4, Python) sont moins urgentes mais doivent être patchées dans le plan de maintenance. npm audit fix réglera les CVE Node.js/Python.")
print('24 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — MITRE ATT&CK
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 25); lbl(s, 'WAZUH — MITRE ATT&CK'); ttl(s, 'MITRE ATT&CK — 1 600+ Alertes — Impact T1499 × 1594', sz=20); div(s)
screenshot(s, IMGS['mitre'], Inches(0.45), Inches(1.75), Inches(7.6), Inches(3.4),
           'MITRE ATT&CK Dashboard — Top tactics: Impact ×1594 · Defense Evasion ×6 · PCI DSS 11.5/10.6.1/10.2.6')
li_list(s, [
    ('Impact T1499 :','×1594 (96%) — Endpoint DoS / modifications binaires /bin'),
    ('Defense Evasion :','×6 — T1562 Disable or Modify Tools'),
    ('PCI DSS :','11.5 (1594) · 10.6.1 (34) · 10.2.6 (12) · 10.2.7 (2)'),
    ('RGPD Art.25 :','Privacy by Design ✓ — Surveillance continue'),
], x=Inches(8.2), y=Inches(1.75), w=Inches(4.75), eh=Inches(0.52))
notes(s, "MITRE ATT&CK est le framework universel de classification des comportements malveillants utilisé par tous les SOC professionnels mondiaux. Wazuh mappe automatiquement ses règles aux techniques ATT&CK. Les 1594 événements T1499 (Endpoint Denial of Service / Resource Exhaustion) correspondent aux modifications FIM des binaires /bin — Wazuh corrèle ce pattern à une technique d'Impact potentielle. Les 6 événements Defense Evasion (T1562 Disable or Modify Tools) signalent des tentatives de désactivation d'outils de sécurité. La couverture PCI DSS (11.5, 10.6.1, 10.2.6) valide aussi la conformité aux normes de sécurité des paiements.")
print('25 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — 5 MODULES WAZUH
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 26); lbl(s, 'WAZUH SIEM — MODULES'); ttl(s, 'Wazuh — 5 Modules de Sécurité Actifs — CDC §3.3'); div(s)
mods = [
    ('Security Events','Journalisation temps réel',
     'Tableau central : événements classés par niveau (3 à 15)\nIP source · agent · règle déclenchée · horodatage\nBase de toute la corrélation SIEM'),
    ('MITRE ATT&CK','Classification internationale',
     'T1499 Impact ×1594 · T1562 Defense Evasion ×6\nTaxonomie universelle SOC\nCorrélation automatique techniques d\'attaque'),
    ('GDPR — Module RGPD','Conformité Art. 15/17/33',
     'Art.15 (accès) · Art.17 (effacement) · Art.33 (violation)\nComplète nos auditLogs applicatifs\nSurveillance conformité infrastructure'),
    ('Integrity Monitoring','FIM — Hash SHA-256',
     '2 500+ events · /bin/bash /bin/cp modifiés\nAlerte si rootkit ou backdoor détecté\nÉquivalent IDS filesystem'),
    ('Vulnerabilities','CVE — CVSS',
     '17 CVE · 8 High + 9 Medium · Docker dominant\nCVE-2019-5736 CVSS3=8.6 → patch prioritaire\nScan packages vs NVD NIST automatique'),
]
cw3 = Inches(2.45); ch3 = Inches(2.2)
for i, (mod, sub, desc) in enumerate(mods):
    cx = Inches(0.45 + i * 2.58)
    r(s, cx, Inches(1.75), cw3, ch3, NAVY_M, BORD)
    r(s, cx, Inches(1.75), cw3, Inches(0.5), NAVY_L, None)
    tx(s, mod, cx, Inches(1.77), cw3, Inches(0.28), sz=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, sub, cx, Inches(2.08), cw3, Inches(0.24), sz=7.5, color=GREY, align=PP_ALIGN.CENTER, italic=True)
    tx(s, desc, cx + Inches(0.12), Inches(2.38), cw3 - Inches(0.2), ch3 - Inches(0.72), sz=7.5, color=W80)
notes(s, "Les 5 modules couvrent l'ensemble du spectre de sécurité SIEM : Security Events (logs temps réel), MITRE ATT&CK (classification internationale des menaces), GDPR (conformité réglementaire infrastructure), Integrity Monitoring (IDS filesystem), et Vulnerabilities (CVE/CVSS). Le module GDPR de Wazuh est directement lié aux exigences RGPD — il surveille la conformité au niveau infrastructure et complète les auditLogs applicatifs Firestore. Ensemble, ces 5 modules constituent une surveillance complète couvrant la sécurité, la conformité et la détection de menaces.")
print('26 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — MONITORING 2 NIVEAUX — DEFENSE IN DEPTH
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 27); lbl(s, 'DÉFENSE EN PROFONDEUR'); ttl(s, 'Monitoring à 2 Niveaux — Defense in Depth — CDC §3.3 à 100%'); div(s)
# Colonne 1 — Applicatif
r(s, Inches(0.45), Inches(1.75), Inches(5.9), Inches(3.5), NAVY_M, BORD)
r(s, Inches(0.45), Inches(1.75), Inches(5.9), Inches(0.38), NAVY_L, None)
tx(s, 'COUCHE APPLICATIVE — auditLogs Firestore + /monitoring',
   Inches(0.45), Inches(1.77), Inches(5.85), Inches(0.32), sz=8.5, bold=True, color=WHITE)
app_items = [
    'WAF — SQLi · XSS · PathTraversal · CmdInj · Agents suspects',
    '9 types audit logs immuables Firestore (append-only)',
    'Score /100 temps réel + 4 alertes automatiques',
    'Scanner DAST 12 tests OWASP — Score 92/100',
    'Dashboard SIEM 3 onglets (admin only · 60s refresh)',
    'RGPD Art.15/16/17/33 · AES-256-CBC · HTTPS/HSTS',
]
for i, item in enumerate(app_items):
    ty = Inches(2.22) + i * Inches(0.37)
    tx(s, f'✓  {item}', Inches(0.6), ty, Inches(5.6), Inches(0.34), sz=8, color=W80)

# Flèche centrale
tx(s, '⟺', Inches(6.4), H/2 - Inches(1.2), Inches(0.8), Inches(0.5), sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Colonne 2 — Infrastructure
r(s, Inches(7.0), Inches(1.75), Inches(5.9), Inches(3.5), NAVY_M, BORD)
r(s, Inches(7.0), Inches(1.75), Inches(5.9), Inches(0.38), NAVY_L, None)
tx(s, 'COUCHE INFRASTRUCTURE — Wazuh SIEM Agent actif',
   Inches(7.0), Inches(1.77), Inches(5.85), Inches(0.32), sz=8.5, bold=True, color=WHITE)
inf_items = [
    'Agent 001 actif — coverage 100% (Debian 12)',
    '213 événements détectés en < 24h sans config extra',
    'CIS Benchmark Debian 12 — audit configuration',
    'Module MITRE ATT&CK — T1499/T1562 classifiés',
    'Module GDPR — Art.15/17/33 infrastructure',
    'Vulnerability Detection — 17 CVE/CVSS détectées',
]
for i, item in enumerate(inf_items):
    ty = Inches(2.22) + i * Inches(0.37)
    tx(s, f'✓  {item}', Inches(7.15), ty, Inches(5.6), Inches(0.34), sz=8, color=W80)

r(s, Inches(0.45), Inches(5.38), W - Inches(0.9), Inches(0.38), NAVY_A, BORD)
tx(s, '→  CDC §3.3 couvert à 100% : Accès sécurisé (RBAC) · Conformité RGPD · Journalisation des actions · HTTPS · Monitoring',
   Inches(0.6), Inches(5.42), W - Inches(1.1), Inches(0.3), sz=9, bold=True, color=WHITE)
notes(s, "La stratégie Defense in Depth (défense en profondeur) consiste à avoir plusieurs couches de sécurité indépendantes : si une couche est contournée, les autres continuent de protéger. La couche applicative (WAF + RBAC + JWT + AuditLogs + Dashboard) protège contre les attaques web. La couche infrastructure (Wazuh SIEM) protège contre les menaces système (rootkits, CVE, misconfigurations, brute force SSH). Ensemble, elles couvrent intégralement le CDC §3.3 sur 5 axes : accès sécurisé, HTTPS, RGPD, RBAC, journalisation.")
print('27 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — SYNTHÈSE CDC §3.3
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 28); lbl(s, 'SYNTHÈSE FINALE'); ttl(s, 'Synthèse — Couverture CDC §3.3 à 100%'); div(s)
tw = Inches(12.4); tx0 = Inches(0.45)
cols28 = [('Exigence CDC §3.3',2.8),('Étapes',0.9),('Implémentation technique',5.5),('Wazuh',1.5),('Statut',1.0)]
fracs28 = [c[1]/11.7 for c in cols28]
tbl_hdr(s, cols28, tx0, Inches(1.75), tw)
cdc = [
    ('Accès sécurisé id/mdp','1·2·3·5','JWT HS256 · bcrypt 12r · WAF · rate limit · lockout 5 échecs','Security Events','✓ FAIT'),
    ('HTTPS obligatoire','3','Firebase HTTPS auto · Helmet HSTS max-age=31536000','—','✓ FAIT'),
    ('RGPD conforme','5·6·8','DATA_EXPORT · DATA_ANONYMIZE · AES-256 · Art.15/16/17/33','Module GDPR','✓ FAIT'),
    ('RBAC — Gestion rôles','5·6·8·9','checkRole() Firestore · 4 rôles · Firestore Rules deny','Security Events','✓ FAIT'),
    ('Journalisation actions','6·7·8','9 events immuables · SIEM dashboard · serverTimestamp()','213+ events','✓ FAIT'),
    ('Détection vulnérabilités','Wazuh','17 CVE identifiées · Plan correctif prioritisé','Vuln Detector','✓ FAIT'),
    ('Intégrité fichiers système','Wazuh','FIM 2500+ · Rule 550 · Hash SHA-256','Integrity Mon.','✓ FAIT'),
    ('Détection menaces','Wazuh','MITRE ATT&CK T1499/T1562 · 1600+ alertes','MITRE ATT&CK','✓ FAIT'),
]
for i, row in enumerate(cdc):
    tbl_row(s, row, fracs28, tx0, Inches(2.04) + Inches(0.31)*i, tw, Inches(0.29), even=(i%2==0))
r(s, tx0, Inches(4.54), tw, Inches(0.35), NAVY_A, BORD)
tx(s, 'CDC §3.3 couvert à 100% : chaque exigence a une implémentation technique vérifiable et démontrables en live devant le jury.',
   tx0 + Inches(0.1), Inches(4.58), tw - Inches(0.2), Inches(0.28), sz=9, bold=True, color=WHITE)
notes(s, "Ce tableau de synthèse est la preuve définitive de la couverture CDC §3.3 à 100%. Chaque ligne correspond à une exigence du cahier des charges avec son implémentation technique précise et la corrélation avec Wazuh SIEM. Toutes les exigences sont marquées FAIT avec des preuves vérifiables : captures d'écran dashboard, logs auditLogs, rapports DAST, métriques Wazuh. La colonne Wazuh montre comment le SIEM renforce et complète la sécurité applicative au niveau infrastructure.")
print('28 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — BILAN GLOBAL & PERSPECTIVES
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s); hdr(s, 29); lbl(s, 'BILAN GLOBAL'); ttl(s, 'Bilan Global — Points Forts & Perspectives Post-Soutenance'); div(s)
# Gauche — Points forts
r(s, Inches(0.45), Inches(1.75), Inches(6.1), Inches(3.55), NAVY_M, BORD)
r(s, Inches(0.45), Inches(1.75), Inches(6.1), Inches(0.3), NAVY_L, None)
tx(s, 'Points Forts Réalisés', Inches(0.45), Inches(1.77), Inches(6.0), Inches(0.26), sz=9, bold=True, color=WHITE)
points = [
    'WAF maison — 5 types OWASP bloqués (waf.js)',
    '9 audit logs immuables Firestore (append-only)',
    'Dashboard SIEM /monitoring — Score /100 temps réel',
    'Scanner DAST 12 tests OWASP automatisés (92/100)',
    'Wazuh SIEM — 213 événements < 24h détectés',
    'Agent Wazuh actif — coverage 100%',
    'CDC §3.3 couvert à 100% — toutes exigences',
]
for i, p in enumerate(points):
    ty = Inches(2.12) + i * Inches(0.43)
    tx(s, f'✓  {p}', Inches(0.6), ty, Inches(5.8), Inches(0.38), sz=9, color=W80)
# Droite — Perspectives
r(s, Inches(6.72), Inches(1.75), Inches(6.2), Inches(3.55), NAVY_M, BORD)
r(s, Inches(6.72), Inches(1.75), Inches(6.2), Inches(0.3), NAVY_L, None)
tx(s, 'Perspectives Post-Soutenance', Inches(6.72), Inches(1.77), Inches(6.1), Inches(0.26), sz=9, bold=True, color=WHITE)
persp = [
    'Déploiement Azure Container Apps (nginx SPA)',
    'Alertes email/SMS automatiques (Twilio / SendGrid)',
    'Intégration CI/CD pipeline sécurité (GitHub Actions)',
    'Wazuh agent sur serveur production (hors Docker)',
    'Rotation automatique clés JWT (Redis blacklist)',
    'Dashboard Grafana (métriques Firebase + Wazuh)',
    'Pentest professionnel post-déploiement (OWASP ZAP)',
]
for i, p in enumerate(persp):
    ty = Inches(2.12) + i * Inches(0.43)
    tx(s, f'→  {p}', Inches(6.87), ty, Inches(5.9), Inches(0.38), sz=9, color=W80)
notes(s, "Le projet démontre une maîtrise complète de la sécurité applicative et infrastructure. Les points forts couvrent l'ensemble du cycle de vie sécuritaire : conception (Security by Design), implémentation (WAF, JWT, RBAC), surveillance (Wazuh SIEM, Dashboard /monitoring), et tests (DAST 92/100). Les perspectives montrent la maturité de l'approche : la sécurité ne s'arrête pas à la soutenance — le pentest professionnel, les alertes automatiques, et l'intégration CI/CD sont les prochaines étapes naturelles pour un déploiement en production réelle.")
print('29 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — MERCI / CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG1)
r(s, 0, 0, W, H * 0.07, NAVY_A, None)
r(s, 0, H - H * 0.07, W, H * 0.07, NAVY_A, None)
r(s, Inches(0.45), Inches(0.8), W - Inches(0.9), H * 0.75, NAVY, None)
hdr(s, 30)
tx(s, 'Merci', Inches(0.45), Inches(1.0), W - Inches(0.9), Inches(1.0),
   sz=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, 'Questions & Démonstration Live', Inches(0.45), Inches(2.1), W - Inches(0.9), Inches(0.45),
   sz=16, bold=True, color=W80, align=PP_ALIGN.CENTER)
r(s, Inches(2.0), Inches(2.65), W - Inches(4.0), Pt(1.5), WHITE, None)
quote = ('"Le monitoring couvre intégralement le CDC §3.3 sur deux niveaux :\n'
         'monitoring applicatif (WAF · 9 audit logs Firestore · dashboard SIEM /monitoring · DAST 12 tests)\n'
         'et monitoring infrastructure Wazuh (213 événements · MITRE ATT&CK · GDPR · CIS Benchmark).\n'
         'Ensemble, ils assurent une défense en profondeur conforme RGPD."')
tx(s, quote, Inches(0.9), Inches(2.85), W - Inches(1.8), Inches(1.1),
   sz=10, color=W80, italic=True, align=PP_ALIGN.CENTER)
r(s, Inches(2.0), Inches(4.1), W - Inches(4.0), Pt(1.5), WHITE, None)
tx(s, 'Amine BAHOU  /  Anass Akker', Inches(0.45), Inches(4.25), W - Inches(0.9), Inches(0.35),
   sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, 'PFE Bachelor Cybersécurité / Cyberdéfense — YNOV Campus 2026',
   Inches(0.45), Inches(4.65), W - Inches(0.9), Inches(0.3),
   sz=10, color=GREY, align=PP_ALIGN.CENTER)
print('30 ✓')

# ══════════════════════════════════════════════════════════════════════════════
# SAUVEGARDE
# ══════════════════════════════════════════════════════════════════════════════
out = os.path.join(BASE, 'PRESENTATION_PFE_SECURITE_30SLIDES.pptx')
prs.save(out)
sz = os.path.getsize(out)
print(f'\n✓ {out}')
print(f'  Taille : {sz/1024/1024:.1f} Mo  ·  {TOTAL} slides')
