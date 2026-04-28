#!/usr/bin/env python3
"""
Ajoute à PRESENTATION_PFE_SECURITE_UPDATED.pptx :
  - Slide 32 : Toutes les commandes de sécurité (WAF, JWT, DAST, Wazuh...)
  - Slide 33 : Guide Démonstration Live devant le jury (étape par étape)
"""

import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
import os

BASE = "/Users/anass/Downloads/frais-gestionScolaire 4"
SRC  = os.path.join(BASE, "PRESENTATION_PFE_SECURITE_UPDATED.pptx")
DST  = os.path.join(BASE, "PRESENTATION_PFE_SECURITE_UPDATED.pptx")

# ─── Couleurs thème ───────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0F, 0x17, 0x2A)   # fond sombre
C_ACCENT   = RGBColor(0x38, 0xBD, 0xF8)   # cyan
C_GREEN    = RGBColor(0x10, 0xB9, 0x81)   # vert teal
C_YELLOW   = RGBColor(0xFB, 0xBF, 0x24)   # jaune
C_RED      = RGBColor(0xEF, 0x44, 0x44)   # rouge
C_HDR      = RGBColor(0x25, 0x63, 0xEB)   # bleu header
C_CARD     = RGBColor(0x1E, 0x29, 0x3B)   # card bg
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DIM      = RGBColor(0x94, 0xA3, 0xB8)
C_CODE_BG  = RGBColor(0x0D, 0x1A, 0x2E)   # fond code

# ─── Dimensions slide (widescreen 13.33" x 7.50") ────────────────────────────
SW = Inches(13.33)
SH = Inches(7.50)


def rgb_to_hex(r):
    return "{:02X}{:02X}{:02X}".format(r[0], r[1], r[2]) if isinstance(r, tuple) else \
           "{:02X}{:02X}{:02X}".format(r.r, r.g, r.b)


def add_rect(slide, l, t, w, h, fill_color, border_color=None, border_w=0):
    """Ajoute un rectangle coloré."""
    sp = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if border_color:
        sp.line.color.rgb = border_color
        sp.line.width = Pt(border_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_text(slide, text, l, t, w, h, font_size=14, font_color=C_WHITE,
             bold=False, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Ajoute un TextBox."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Consolas" if font_color == C_ACCENT or italic else "Calibri"
    return txb


def add_multiline_text(slide, lines, l, t, w, h, font_size=12):
    """
    lines = list of (text, color, bold, is_code)
    """
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, color, bold, is_code) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Consolas" if is_code else "Calibri"
    return txb


# ─────────────────────────────────────────────────────────────────────────────
# Copie d'un slide layout vierge (dernier layout = blanc)
# ─────────────────────────────────────────────────────────────────────────────
def add_blank_slide(prs):
    """Ajoute un slide vierge basé sur le layout 'Blank'."""
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", "vierge", "titre vierge", "blank slide"):
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6]  # index 6 = Blank dans la plupart des thèmes
    return prs.slides.add_slide(blank_layout)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 32 — COMMANDES DE SÉCURITÉ
# ─────────────────────────────────────────────────────────────────────────────
def build_commands_slide(prs):
    slide = add_blank_slide(prs)
    slide_num = len(prs.slides)

    # Fond sombre
    add_rect(slide, 0, 0, SW, SH, C_BG)

    # ── Header bar ────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, SW, Inches(0.42), RGBColor(0x14, 0x1E, 0x37))
    add_rect(slide, 0, Inches(0.42), SW, Pt(2), C_ACCENT)

    add_text(slide, "PFE Cybersécurité · YNOV Campus 2026",
             Inches(0.35), Inches(0.10), Inches(8), Inches(0.25),
             font_size=11, font_color=C_DIM)
    add_text(slide, f"{slide_num:02d} / 33",
             Inches(12.50), Inches(0.10), Inches(0.75), Inches(0.25),
             font_size=11, font_color=C_ACCENT, align=PP_ALIGN.RIGHT)

    # ── Titre ─────────────────────────────────────────────────────────────────
    add_rect(slide, Inches(0.45), Inches(0.55), Inches(3.8), Pt(3), C_ACCENT)
    add_text(slide, "ANNEXE — COMMANDES DE SÉCURITÉ",
             Inches(0.45), Inches(0.45), Inches(9), Inches(0.35),
             font_size=20, font_color=C_ACCENT, bold=True)
    add_text(slide, "Toutes les commandes de test/démonstration par catégorie",
             Inches(0.45), Inches(0.80), Inches(10), Inches(0.28),
             font_size=13, font_color=C_WHITE)

    # ── Layout 2 colonnes x 3 sections ───────────────────────────────────────
    sections = [
        {
            "title": "🛡️  WAF — Test Blocage",
            "color": C_HDR,
            "lines": [
                ("# SQLi — doit retourner 403", C_DIM, False, False),
                ("curl -s -o /dev/null -w \"Réponse HTTP: %{http_code}\" \\", C_ACCENT, False, True),
                ("  'localhost:5001/gestionadminastration/us-central1/api", C_ACCENT, False, True),
                ("   /etudiants?id=1%20OR%201=1'", C_ACCENT, False, True),
                ("", C_WHITE, False, False),
                ("# XSS — doit retourner 403", C_DIM, False, False),
                ("curl -s -o /dev/null -w \"Réponse HTTP: %{http_code}\" \\", C_ACCENT, False, True),
                ("  'localhost:5001/gestionadminastration/us-central1/api", C_ACCENT, False, True),
                ("   /etudiants?nom=<script>alert(1)</script>'", C_ACCENT, False, True),
            ]
        },
        {
            "title": "🔐  Auth JWT — Brute Force",
            "color": RGBColor(0x16, 0x50, 0x80),
            "lines": [
                ("# Brute force → lockout après 5 essais", C_DIM, False, False),
                ("for i in {1..6}; do", C_ACCENT, False, True),
                ("  curl -s -o /dev/null -w \"Tentative $i → %{http_code}\\n\" \\", C_ACCENT, False, True),
                ("  -X POST localhost:5001/gestionadminastration/", C_ACCENT, False, True),
                ("          us-central1/api/auth/login \\", C_ACCENT, False, True),
                ("  -H 'Content-Type: application/json' \\", C_ACCENT, False, True),
                ("  -d '{\"email\":\"admin@school.fr\",\"password\":\"WRONG\"}'", C_ACCENT, False, True),
                ("done", C_ACCENT, False, True),
                ("# → Tentative 6 = 429 Too Many Requests", C_GREEN, False, False),
            ]
        },
        {
            "title": "🔎  DAST Scanner Automatique",
            "color": RGBColor(0x15, 0x70, 0x3E),
            "lines": [
                ("# Lancer le scan de vulnérabilités", C_DIM, False, False),
                ("cd \"/Users/anass/Downloads/frais-gestionScolaire 4/back/functions\"", C_ACCENT, False, True),
                ("node scripts/security_scan.js", C_ACCENT, False, True),
                ("", C_WHITE, False, False),
                ("# Résultat attendu :", C_DIM, False, False),
                ("[PASS] WAF actif — SQLi bloqué", C_GREEN, True, False),
                ("[PASS] Rate limiting — 429 après 5 req", C_GREEN, True, False),
                ("[PASS] JWT validation — 401 token invalide", C_GREEN, True, False),
                ("[PASS] Headers sécurité — HSTS présent", C_GREEN, True, False),
                ("[WARN] 17 CVE détectées (audit npm)", C_YELLOW, True, False),
            ]
        },
        {
            "title": "🚀  Démarrage Application",
            "color": RGBColor(0x05, 0x54, 0x3A),
            "lines": [
                ("# Terminal 1 — Backend", C_DIM, False, False),
                ("cd \"/Users/anass/Downloads/frais-gestionScolaire 4/back/functions\"", C_ACCENT, False, True),
                ("node server_local.js", C_ACCENT, False, True),
                ("# → http://localhost:5001", C_GREEN, False, False),
                ("", C_WHITE, False, False),
                ("# Terminal 2 — Frontend", C_DIM, False, False),
                ("cd \"/Users/anass/Downloads/frais-gestionScolaire 4/front\"", C_ACCENT, False, True),
                ("npm run dev", C_ACCENT, False, True),
                ("# → http://localhost:5173", C_GREEN, False, False),
            ]
        },
        {
            "title": "🐳  Wazuh — Démarrage Docker",
            "color": RGBColor(0x1D, 0x4E, 0x89),
            "lines": [
                ("# 1. Symlink sans espace (1 seule fois)", C_DIM, False, False),
                ("ln -s \".../front/wazuh-docker/single-node\" /tmp/wazuh-node", C_ACCENT, False, True),
                ("cd /tmp/wazuh-node", C_ACCENT, False, True),
                ("", C_WHITE, False, False),
                ("# 2. Démarrer les 3 conteneurs", C_DIM, False, False),
                ("docker compose up -d", C_ACCENT, False, True),
                ("# → manager · indexer · dashboard", C_GREEN, False, False),
                ("", C_WHITE, False, False),
                ("# 3. Vérifier (attendre ~2 min)", C_DIM, False, False),
                ("docker compose ps   # 3 conteneurs Up", C_ACCENT, False, True),
            ]
        },
        {
            "title": "🔗  Wazuh — Agent main-machine",
            "color": RGBColor(0x5B, 0x21, 0x86),
            "lines": [
                ("# 1. Arrêter l'agent", C_DIM, False, False),
                ("sudo /Library/Ossec/bin/wazuh-control stop", C_ACCENT, False, True),
                ("# 2. Pointer vers manager localhost", C_DIM, False, False),
                ("sudo sed -i \"\" 's|<address>.*</address>|", C_ACCENT, False, True),
                ("<address>127.0.0.1</address>|' /Library/Ossec/etc/ossec.conf", C_ACCENT, False, True),
                ("# 3. Vérifier l'adresse", C_DIM, False, False),
                ("sudo grep \"<address>\" /Library/Ossec/etc/ossec.conf", C_ACCENT, False, True),
                ("# 4. Re-enregistrer main-machine", C_DIM, False, False),
                ("sudo /Library/Ossec/bin/agent-auth -m 127.0.0.1 -A main-machine", C_ACCENT, False, True),
                ("# 5. Redémarrer + statut", C_DIM, False, False),
                ("sudo /Library/Ossec/bin/wazuh-control start", C_ACCENT, False, True),
                ("sudo /Library/Ossec/bin/wazuh-control status", C_ACCENT, False, True),
                ("# → https://localhost → main-machine Active", C_GREEN, False, False),
            ]
        },
        {
            "title": "🔒  Chiffrement AES-256",
            "color": RGBColor(0x92, 0x40, 0x0E),
            "lines": [
                ("# Test chiffrement (Node.js REPL)", C_DIM, False, False),
                ("node -e \"", C_ACCENT, False, True),
                ("const enc = require('./src/utils/encryption');", C_ACCENT, False, True),
                ("const c = enc.encrypt('donnee-sensible');", C_ACCENT, False, True),
                ("console.log('Chiffré:', c);", C_ACCENT, False, True),
                ("console.log('Déchiffré:', enc.decrypt(c));\"", C_ACCENT, False, True),
                ("", C_WHITE, False, False),
                ("# Vérifier headers sécurité HTTPS", C_DIM, False, False),
                ("curl -I http://localhost:5001/api/auth/login", C_ACCENT, False, True),
                ("# → Strict-Transport-Security: max-age=31536000", C_GREEN, False, False),
            ]
        },
        {
            "title": "📝  RGPD — Routes de Conformité",
            "color": RGBColor(0x1F, 0x4E, 0x79),
            "lines": [
                ("# Art. 15 — Export données (avec token admin)", C_DIM, False, False),
                ("curl -H 'Authorization: Bearer $TOKEN' \\", C_ACCENT, False, True),
                ("  http://localhost:5001/api/users/USER_ID/export", C_ACCENT, False, True),
                ("", C_WHITE, False, False),
                ("# Art. 17 — Anonymisation", C_DIM, False, False),
                ("curl -X DELETE \\", C_ACCENT, False, True),
                ("  -H 'Authorization: Bearer $TOKEN' \\", C_ACCENT, False, True),
                ("  http://localhost:5001/api/users/USER_ID/data", C_ACCENT, False, True),
                ("", C_WHITE, False, False),
                ("# → Vérifier auditLogs Firestore", C_GREEN, False, False),
            ]
        },
    ]

    # ── Séparateur vertical central ────────────────────────────────────────────
    add_rect(slide, Inches(6.60), Inches(1.10), Pt(1.5), Inches(6.05), RGBColor(0x32, 0x46, 0x64))

    # 2 colonnes : gauche (index pairs), droite (index impairs)
    col_x = [Inches(0.28), Inches(6.78)]
    section_h = Inches(1.44)
    section_w = Inches(6.22)
    start_y = Inches(1.12)
    gap_y = Inches(0.08)

    for i, sec in enumerate(sections):
        col = i % 2          # 0=gauche, 1=droite
        row = i // 2
        x = col_x[col]
        y = start_y + row * (section_h + gap_y)

        # Card bg + border
        add_rect(slide, x, y, section_w, section_h, C_CARD,
                 border_color=RGBColor(0x2A, 0x3D, 0x58), border_w=0.8)

        # Left accent stripe
        add_rect(slide, x, y, Inches(0.06), section_h, sec["color"])

        # Header band
        add_rect(slide, x + Inches(0.06), y, section_w - Inches(0.06), Inches(0.33), sec["color"])

        # Title
        add_text(slide, sec["title"],
                 x + Inches(0.14), y + Inches(0.05), section_w - Inches(0.20), Inches(0.26),
                 font_size=12.5, font_color=C_WHITE, bold=True)

        # Code background
        add_rect(slide, x + Inches(0.10), y + Inches(0.37),
                 section_w - Inches(0.18), section_h - Inches(0.43),
                 C_CODE_BG)

        # Code lines
        add_multiline_text(slide, sec["lines"],
                           x + Inches(0.15), y + Inches(0.40),
                           section_w - Inches(0.25), section_h - Inches(0.48),
                           font_size=9.0)

    # Footer note
    add_rect(slide, 0, Inches(7.18), SW, Inches(0.32), RGBColor(0x0D, 0x14, 0x28))
    add_rect(slide, 0, Inches(7.18), SW, Pt(2), C_ACCENT)
    add_text(slide,
             "⚡  Pré-requis : back (node server_local.js :5001) + frontend (npm run dev :8081) + Wazuh Docker · TOKEN = réponse /auth/login",
             Inches(0.35), Inches(7.20), Inches(12.5), Inches(0.26),
             font_size=10.5, font_color=C_DIM)

    print(f"✅ Slide {slide_num} — Commandes de sécurité créé")
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 33 — GUIDE DÉMONSTRATION LIVE
# ─────────────────────────────────────────────────────────────────────────────
def build_demo_guide_slide(prs):
    slide = add_blank_slide(prs)
    slide_num = len(prs.slides)

    # Fond sombre
    add_rect(slide, 0, 0, SW, SH, C_BG)

    # Header bar
    add_rect(slide, 0, 0, SW, Inches(0.42), RGBColor(0x14, 0x1E, 0x37))
    add_rect(slide, 0, Inches(0.42), SW, Pt(2), C_ACCENT)
    add_text(slide, "PFE Cybersécurité · YNOV Campus 2026",
             Inches(0.35), Inches(0.10), Inches(8), Inches(0.25),
             font_size=11, font_color=C_DIM)
    add_text(slide, f"{slide_num:02d} / 33",
             Inches(12.50), Inches(0.10), Inches(0.75), Inches(0.25),
             font_size=11, font_color=C_ACCENT, align=PP_ALIGN.RIGHT)

    # Titre
    add_rect(slide, Inches(0.45), Inches(0.55), Inches(4.2), Pt(3), C_GREEN)
    add_text(slide, "DÉMONSTRATION LIVE — Guide Jury",
             Inches(0.45), Inches(0.45), Inches(10), Inches(0.35),
             font_size=20, font_color=C_GREEN, bold=True)
    add_text(slide, "Script pas-à-pas · 8 min · 7 démonstrations · 3 terminaux · OWASP Top 10 + Wazuh SIEM + RGPD",
             Inches(0.45), Inches(0.80), Inches(12.5), Inches(0.28),
             font_size=12, font_color=C_WHITE)

    # ── Étapes ───────────────────────────────────────────────────────────────
    steps = [
        {
            "num": "01",
            "time": "1 min",
            "title": "WAF — Injection SQL bloquée",
            "color": C_RED,
            "cmd": "curl -s -o /dev/null -w \"Réponse HTTP: %{http_code}\\n\" \\\n  'http://localhost:5001/gestionadminastration/us-central1/api\\\n   /etudiants?id=1%20OR%201=1'",
            "result": "Réponse HTTP: 403",
            "say": "J'essaie une injection SQL dans l'URL. Le WAF middleware intercepte la requête avant qu'elle touche la base de données → 403 Forbidden.",
        },
        {
            "num": "02",
            "time": "30s",
            "title": "WAF — XSS bloqué",
            "color": RGBColor(0xDC, 0x26, 0x26),
            "cmd": "curl -s -o /dev/null -w \"Réponse HTTP: %{http_code}\\n\" \\\n  'http://localhost:5001/gestionadminastration/us-central1/api\\\n   /etudiants?nom=<script>alert(1)</script>'",
            "result": "Réponse HTTP: 403",
            "say": "Même chose pour le XSS — la balise script est bloquée immédiatement par le WAF. OWASP A03.",
        },
        {
            "num": "03",
            "time": "1 min",
            "title": "Brute Force — Rate Limiting (OWASP A07)",
            "color": C_YELLOW,
            "cmd": "for i in {1..6}; do\n  curl -s -o /dev/null -w \"Tentative $i → %{http_code}\\n\" \\\n  -X POST http://localhost:5001/gestionadminastration/us-central1/api/auth/login \\\n  -H 'Content-Type: application/json' \\\n  -d '{\"email\":\"admin@school.fr\",\"password\":\"WRONG\"}'\ndone",
            "result": "Tentative 1→401  2→401  3→401  4→401  5→401  6→429",
            "say": "5 tentatives échouées = compte verrouillé 15 minutes. Protection brute force — OWASP A07.",
        },
        {
            "num": "04",
            "time": "1 min",
            "title": "Scan DAST automatique — OWASP Top 10",
            "color": C_GREEN,
            "cmd": "node \"/Users/anass/Downloads/frais-gestionScolaire 4/back/functions/scripts/security_scan.js\"",
            "result": "[PASS] WAF · JWT · Headers · Rate limiting · CORS  |  [WARN] 17 CVE npm",
            "say": "Ce script automatise tous les tests OWASP Top 10. Tout passe sauf 17 CVE détectées dans les dépendances npm — détectées mais non critiques.",
        },
        {
            "num": "05",
            "time": "2 min",
            "title": "Wazuh SIEM — Dashboard en direct",
            "color": RGBColor(0x5B, 0x21, 0x86),
            "cmd": "Ouvrir https://localhost  →  admin / SecretPassword\n→ Security Events · File Integrity · MITRE ATT&CK · Policy Monitoring",
            "result": "436 420 events · FIM 89% · T1565 Stored Data · Policy anomalies",
            "say": "436 420 événements collectés. FIM surveille les fichiers système. Corrélation MITRE ATT&CK automatique. Vérification continue de la config OS.",
        },
        {
            "num": "06",
            "time": "1 min",
            "title": "Dashboard Monitoring intégré",
            "color": RGBColor(0x06, 0x7C, 0x4A),
            "cmd": "Ouvrir http://localhost:8081/monitoring  (frontend React)",
            "result": "Alertes temps réel · Graphiques d'activité · Événements sécurité",
            "say": "Ce dashboard est intégré dans l'application — l'admin voit les alertes de sécurité directement depuis l'interface, sans aller sur Wazuh.",
        },
        {
            "num": "07",
            "time": "1 min",
            "title": "RGPD — Traçabilité Art.5 & Art.32",
            "color": RGBColor(0x1F, 0x4E, 0x79),
            "cmd": "1. Onglet privé → http://localhost:8081/login\n   → Taper un MAUVAIS mot de passe → connexion échouée\n2. http://localhost:8081/monitoring → SIEM Logs → Rafraîchir\n   → Voir auth_failure avec timestamp + IP",
            "result": "auth_failure loggué · timestamp serveur · IP · email tracés dans Firestore",
            "say": "Je simule une tentative de connexion échouée — l'événement auth_failure est immédiatement tracé dans Firestore avec timestamp, IP, identité. Conformité RGPD Art.5 et Art.32.",
        },
    ]

    # ── Séparateur vertical central ────────────────────────────────────────────
    add_rect(slide, Inches(6.60), Inches(1.05), Pt(1.5), Inches(6.00), RGBColor(0x32, 0x46, 0x64))

    # Layout : 4 étapes colonne gauche, 3 étapes colonne droite
    col_x = [Inches(0.25), Inches(6.72)]
    step_h = Inches(1.42)
    step_w = Inches(6.28)
    start_y = Inches(1.08)
    gap = Inches(0.06)

    for i, step in enumerate(steps):
        col = i % 2
        row = i // 2
        x = col_x[col]
        y = start_y + row * (step_h + gap)

        # ── Card background ───────────────────────────────────────────────────
        add_rect(slide, x, y, step_w, step_h, C_CARD,
                 border_color=RGBColor(0x2A, 0x3D, 0x58), border_w=0.8)

        # Colored left stripe (4px)
        add_rect(slide, x, y, Inches(0.06), step_h, step["color"])

        # ── Header band ───────────────────────────────────────────────────────
        add_rect(slide, x + Inches(0.06), y, step_w - Inches(0.06), Inches(0.30), step["color"])

        # Numéro étape (badge style)
        add_text(slide, f"ÉTAPE {step['num']}",
                 x + Inches(0.12), y + Inches(0.04), Inches(1.10), Inches(0.23),
                 font_size=10, font_color=C_WHITE, bold=True)

        # Titre étape
        add_text(slide, step["title"],
                 x + Inches(1.25), y + Inches(0.04), Inches(3.60), Inches(0.23),
                 font_size=10.5, font_color=C_WHITE, bold=True)

        # Durée (alignée à droite)
        add_text(slide, f"⏱ {step['time']}",
                 x + Inches(4.70), y + Inches(0.04), Inches(1.40), Inches(0.23),
                 font_size=10, font_color=C_WHITE, bold=False, align=PP_ALIGN.RIGHT)

        # ── Terminal / Commande ───────────────────────────────────────────────
        add_rect(slide, x + Inches(0.10), y + Inches(0.33),
                 step_w - Inches(0.18), Inches(0.42), C_CODE_BG,
                 border_color=RGBColor(0x1E, 0x4A, 0x6E), border_w=0.5)
        add_text(slide, f"$ {step['cmd']}",
                 x + Inches(0.14), y + Inches(0.34), step_w - Inches(0.24), Inches(0.40),
                 font_size=8.0, font_color=C_ACCENT)

        # ── Résultat attendu ──────────────────────────────────────────────────
        add_rect(slide, x + Inches(0.10), y + Inches(0.77),
                 step_w - Inches(0.18), Inches(0.23),
                 RGBColor(0x07, 0x20, 0x10), border_color=C_GREEN, border_w=0.6)
        add_text(slide, f"✔  {step['result']}",
                 x + Inches(0.14), y + Inches(0.78), step_w - Inches(0.28), Inches(0.20),
                 font_size=8.5, font_color=C_GREEN, bold=True)

        # ── Séparateur + quote jury ───────────────────────────────────────────
        add_rect(slide, x + Inches(0.10), y + Inches(1.03),
                 step_w - Inches(0.18), Pt(0.8), RGBColor(0x32, 0x46, 0x64))
        add_text(slide, f'🎤  {step["say"]}',
                 x + Inches(0.12), y + Inches(1.05), step_w - Inches(0.22), Inches(0.34),
                 font_size=8.0, font_color=C_DIM)

    # ── Footer conseils ───────────────────────────────────────────────────────
    add_rect(slide, 0, Inches(7.10), SW, Inches(0.40), RGBColor(0x07, 0x10, 0x20))
    add_rect(slide, 0, Inches(7.10), SW, Pt(2), C_GREEN)
    add_text(slide,
             "💡  AVANT DE COMMENCER :  3 terminaux ouverts  ·  T1 = Backend :5001  ·  T2 = Frontend :8081  ·  T3 = Wazuh Docker (/tmp/wazuh-node)  ·  Token JWT dans le presse-papier  ·  Total démo ≈ 8 min",
             Inches(0.35), Inches(7.12), Inches(12.5), Inches(0.34),
             font_size=9.5, font_color=C_YELLOW, bold=True)

    print(f"✅ Slide {slide_num} — Guide Démonstration Live créé")
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation(SRC)
print(f"📂 Présentation chargée : {len(prs.slides)} slides")

# Supprimer les slides 32+ si déjà présents (re-génération propre)
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
while len(prs.slides) > 31:
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    rId = sldId.get(f"{{{NS_R}}}id")
    if rId:
        prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)
    print(f"  🗑  Slide supprimé (re-génération) — reste {len(prs.slides)}")

build_commands_slide(prs)
build_demo_guide_slide(prs)

prs.save(DST)
print(f"\n✅ Sauvegardé : {DST}")
print(f"   Nombre total de slides : {len(prs.slides)}")
