#!/usr/bin/env python3
"""
Ajoute 13 slides Q&A (12 blocs + 1 intro) à la présentation existante
SECURITE_MONITORING_YNOV_2026.pptx.pptx
→ Produit : SECURITE_MONITORING_YNOV_2026_AVEC_QA.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import copy

# ── Constantes couleurs ──────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x1B, 0x2A)   # navy foncé
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CYAN    = RGBColor(0x00, 0xB4, 0xD8)
YELLOW  = RGBColor(0xFF, 0xD6, 0x00)
GREEN   = RGBColor(0x2D, 0xC6, 0x53)
ORANGE  = RGBColor(0xFF, 0x6B, 0x35)
RED     = RGBColor(0xE6, 0x3A, 0x3A)
PURPLE  = RGBColor(0x7B, 0x2F, 0xBE)
LBLUE   = RGBColor(0x1A, 0x78, 0xC2)
TEAL    = RGBColor(0x0D, 0x9E, 0x8A)
PINK    = RGBColor(0xD6, 0x24, 0x8D)
LGRAY   = RGBColor(0xCC, 0xCC, 0xCC)

# ── Helpers ──────────────────────────────────────────────────────────────────
def cm(v): return Emu(int(v * 360000))

def fill_bg(slide, color=BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=16, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             font="Calibri", wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox

def add_para(tf, text, size=13, bold=False, color=WHITE, indent=False,
             font="Calibri", align=PP_ALIGN.LEFT):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    if indent:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return p

def add_multiline_textbox(slide, lines, l, t, w, h,
                           default_size=13, default_color=WHITE,
                           font="Calibri"):
    """lines = list of (text, size, bold, color) or just str"""
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if isinstance(line, str):
            text, size, bold, color = line, default_size, False, default_color
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            bold = line[2] if len(line) > 2 else False
            color = line[3] if len(line) > 3 else default_color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return txBox

def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}", cm(85), cm(14.5), cm(8), cm(1),
             size=9, color=LGRAY, align=PP_ALIGN.RIGHT)

def header_bar(slide, title, color=LBLUE):
    add_rect(slide, cm(0.5), cm(0.4), cm(32), cm(1.7), color)
    add_text(slide, title, cm(0.8), cm(0.45), cm(31), cm(1.6),
             size=20, bold=True, color=WHITE)

def section_tag(slide, label, color=CYAN):
    add_text(slide, label, cm(0.5), cm(0.05), cm(20), cm(0.5),
             size=7, color=color)

def footer_note(slide, text):
    add_rect(slide, cm(0), cm(14.1), cm(34), cm(1.0), RGBColor(0x08, 0x10, 0x1A))
    add_text(slide, text, cm(0.5), cm(14.15), cm(33), cm(0.9),
             size=8, italic=True, color=LGRAY)

# ── Q&A data ─────────────────────────────────────────────────────────────────
QA_BLOCS = [
    {
        "section": "Q&A — QUESTIONS JURY",
        "title": "Questions Jury — Présentation des 12 Blocs Thématiques",
        "color": PURPLE,
        "intro": True,
        "blocs": [
            "BLOC 1 — WAF & Sécurité Applicative",
            "BLOC 2 — Authentification & JWT",
            "BLOC 3 — Chiffrement AES-256-CBC",
            "BLOC 4 — Audit Logs & Traçabilité",
            "BLOC 5 — RBAC & Contrôle d'Accès",
            "BLOC 6 — Wazuh SIEM",
            "BLOC 7 — Conformité RGPD",
            "BLOC 8 — Architecture & Defense in Depth",
            "BLOC 9 — Tests & Validation",
            "BLOC 10 — Incidents & Réponse",
            "BLOC 11 — Monitoring & Performances",
            "BLOC 12 — Questions Pièges / Approfondissement",
        ],
    },
    {
        "section": "Q&A — BLOC 1",
        "title": "WAF & Sécurité Applicative",
        "color": RED,
        "qas": [
            {
                "q": "Q1 : C'est quoi ton WAF et comment il fonctionne ?",
                "pts": [
                    "waf.js = middleware Express appliqué AVANT l'authentification",
                    "Détecte SQLi / XSS / Path Traversal / CMD Injection / Scanners via regex",
                    "Si attaque → HTTP 403 + WAF_BLOCK enregistré dans Firestore auditLogs",
                ],
            },
            {
                "q": "Q2 : Pourquoi un WAF applicatif et pas un WAF réseau (nginx, Cloudflare) ?",
                "pts": [
                    "Coût zéro, déployé directement dans le code Node.js/Express",
                    "Adapté au contexte Firebase Functions (pas d'infra réseau contrôlable)",
                    "Logging métier précis : on sait qui a attaqué, quand, quelle route",
                ],
            },
            {
                "q": "Q3 : Comment tu valides que ton WAF bloque vraiment les attaques ?",
                "pts": [
                    "security_scan.js = scanner DAST maison, 12 tests automatisés",
                    "Score obtenu : 92/100 sur toutes les catégories (SQLi, XSS, Headers...)",
                    "Résultats reproductibles : node security_scan.js depuis le terminal",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 2",
        "title": "Authentification & JWT",
        "color": ORANGE,
        "qas": [
            {
                "q": "Q4 : Comment fonctionne l'authentification dans ton application ?",
                "pts": [
                    "Firebase Auth génère un idToken JWT signé côté client",
                    "Backend vérifie le token avec admin.auth().verifyIdToken() à chaque requête",
                    "Si invalide ou expiré → HTTP 401, aucune donnée exposée",
                ],
            },
            {
                "q": "Q5 : Quelle est la durée de vie des tokens JWT ?",
                "pts": [
                    "Access token : 24 heures (usage quotidien normal)",
                    "Refresh token : 7 jours (renouvellement silencieux)",
                    "Secret stocké dans .env, jamais dans le code source",
                ],
            },
            {
                "q": "Q6 : Comment tu protèges contre le brute-force sur le login ?",
                "pts": [
                    "Rate limiter : 10 requêtes / 15 min par IP sur /login",
                    "Après 5 échecs → AUTH_LOCKOUT loggué + HTTP 429",
                    "bcrypt 12 rounds : hash volontairement lent pour ralentir les attaques",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 3",
        "title": "Chiffrement AES-256-CBC",
        "color": TEAL,
        "qas": [
            {
                "q": "Q7 : Pourquoi chiffrer les données personnelles dans Firestore ?",
                "pts": [
                    "Conformité RGPD : téléphone et adresse = données sensibles",
                    "Protection defense in depth : si Firestore compromis, données illisibles",
                    "AES-256-CBC = standard industriel approuvé NIST",
                ],
            },
            {
                "q": "Q8 : Comment fonctionne concrètement ton chiffrement AES ?",
                "pts": [
                    "IV aléatoire 128 bits généré à chaque chiffrement avec crypto.randomBytes(16)",
                    "Format stocké : iv_hex:ciphertext_hex dans le champ Firestore",
                    "Déchiffrement : on extrait l'IV du préfixe, puis createDecipheriv()",
                ],
            },
            {
                "q": "Q9 : Où est stockée la clé de chiffrement ENCRYPTION_KEY ?",
                "pts": [
                    "Variable d'environnement obligatoire — jamais dans le code",
                    "Sur Firebase : firebase functions:secrets:set ENCRYPTION_KEY",
                    "En local : fichier .env exclu par .gitignore",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 4",
        "title": "Audit Logs & Traçabilité",
        "color": LBLUE,
        "qas": [
            {
                "q": "Q10 : Comment garantis-tu l'intégrité des logs d'audit ?",
                "pts": [
                    "Règle Firestore : allow update, delete: if false → logs immuables",
                    "9 types d'événements : LOGIN, LOGOUT, CREATE, UPDATE, DELETE...",
                    "Chaque log contient : userId, action, timestamp, IP, userAgent",
                ],
            },
            {
                "q": "Q11 : À quoi servent tes audit logs en pratique ?",
                "pts": [
                    "Forensique : reconstruire la chronologie d'un incident",
                    "Conformité : preuve d'accès et modifications pour audit RGPD",
                    "Monitoring : détecter comportements anormaux (ex. 50 DELETE en 1 min)",
                ],
            },
            {
                "q": "Q12 : Qui peut lire les audit logs ?",
                "pts": [
                    "Admin uniquement via interface dédiée — règles Firestore RBAC",
                    "Wazuh les ingère aussi pour corrélation SIEM côté infrastructure",
                    "Export possible en CSV/Excel pour rapports de conformité",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 5",
        "title": "RBAC & Contrôle d'Accès",
        "color": PURPLE,
        "qas": [
            {
                "q": "Q13 : Combien de rôles dans ton RBAC et qui fait quoi ?",
                "pts": [
                    "admin : accès total, gestion utilisateurs, backup, audit",
                    "sous-admin / comptable : gestion frais, paiements, étudiants",
                    "étudiant / parent : lecture seule de leur propre dossier uniquement",
                ],
            },
            {
                "q": "Q14 : Comment le RBAC est-il appliqué techniquement ?",
                "pts": [
                    "Middleware verifyRole() sur chaque route Express sensible",
                    "Double vérification : Firebase Auth (token) + Firestore (rôle utilisateur)",
                    "Firestore Rules comme dernier rempart si le backend est contourné",
                ],
            },
            {
                "q": "Q15 : Que se passe-t-il si quelqu'un essaie d'accéder hors de son rôle ?",
                "pts": [
                    "HTTP 403 retourné immédiatement par le middleware",
                    "Tentative loggée dans auditLogs avec type UNAUTHORIZED_ACCESS",
                    "Wazuh reçoit l'alerte et peut déclencher une règle de blocage IP",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 6",
        "title": "Wazuh SIEM",
        "color": GREEN,
        "qas": [
            {
                "q": "Q16 : Pourquoi avoir intégré Wazuh dans ton projet ?",
                "pts": [
                    "Surveillance infrastructure en temps réel : logs + FIM + CVE",
                    "436 420 événements analysés, alertes automatiques",
                    "Réponse active possible : blocage IP, isolation hôte compromis",
                ],
            },
            {
                "q": "Q17 : Qu'est-ce que le FIM (File Integrity Monitoring) dans Wazuh ?",
                "pts": [
                    "Surveille les modifications de fichiers critiques (hash MD5/SHA256)",
                    "Si un fichier système change → alerte immédiate niveau 7+",
                    "MITRE T1565.001 : Stored Data Manipulation — représente ~95% des alertes",
                ],
            },
            {
                "q": "Q18 : Comment Wazuh détecte les vulnérabilités CVE ?",
                "pts": [
                    "Agent Wazuh scanne les packages installés et compare avec NVD",
                    "CVE-2019-5736 détecté : CVSS3=8.6, Docker container escape, priorité P1",
                    "Rapport généré avec recommandations de remediation",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 7",
        "title": "Conformité RGPD",
        "color": TEAL,
        "qas": [
            {
                "q": "Q19 : Quelles mesures RGPD as-tu implémentées ?",
                "pts": [
                    "Chiffrement AES-256 des données sensibles (téléphone, adresse)",
                    "Droit à l'effacement : endpoint /api/users/rgpd/delete disponible",
                    "Audit logs pour traçabilité des accès aux données personnelles",
                ],
            },
            {
                "q": "Q20 : Comment tu gères le consentement et les droits des utilisateurs ?",
                "pts": [
                    "Portail étudiant/parent : accès lecture seule à son propre dossier",
                    "Suppression compte : rgpd.js anonymise ou supprime les données",
                    "Export données : format JSON/CSV disponible sur demande",
                ],
            },
            {
                "q": "Q21 : Est-ce que ton application couvre le cahier des charges sécurité ?",
                "pts": [
                    "CDC §3.3 : 100% couvert, 8 exigences validées",
                    "Chaque exigence tracée vers une implémentation technique concrète",
                    "Validé par le scanner DAST (score 92/100) et les tests unitaires",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 8",
        "title": "Architecture & Defense in Depth",
        "color": LBLUE,
        "qas": [
            {
                "q": "Q22 : C'est quoi ta stratégie de défense en profondeur ?",
                "pts": [
                    "Couche 1 applicative : WAF + JWT + RBAC + Rate Limiter + AES",
                    "Couche 2 infrastructure : Wazuh SIEM + FIM + CVE + Rootcheck",
                    "Si couche 1 compromise → couche 2 détecte et alerte en temps réel",
                ],
            },
            {
                "q": "Q23 : Quelle est l'architecture globale de l'application ?",
                "pts": [
                    "Frontend React/TypeScript → Firebase Hosting",
                    "Backend Node.js/Express → Firebase Functions (serverless)",
                    "Base de données Firestore + Storage Firebase",
                ],
            },
            {
                "q": "Q24 : Pourquoi Firebase Functions pour le backend ?",
                "pts": [
                    "Scalabilité automatique, pas de serveur à gérer",
                    "Isolation par fonction : surface d'attaque réduite",
                    "Intégration native Firebase Auth + Firestore Rules",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 9",
        "title": "Tests & Validation",
        "color": ORANGE,
        "qas": [
            {
                "q": "Q25 : Comment tu as testé la sécurité de ton application ?",
                "pts": [
                    "Scanner DAST maison security_scan.js : 12 tests, score 92/100",
                    "Tests unitaires encryption : encryption.test.js (Jest)",
                    "Tests manuels : injection SQL, XSS, path traversal sur toutes les routes",
                ],
            },
            {
                "q": "Q26 : Qu'est-ce que ton scanner DAST teste exactement ?",
                "pts": [
                    "Injections : SQLi, XSS, Path Traversal, Command Injection",
                    "Headers de sécurité : HSTS, X-Frame-Options, CSP, X-Content-Type",
                    "Authentification : tokens expirés, rôles croisés, endpoints sans auth",
                ],
            },
            {
                "q": "Q27 : Quel est le score 92/100 — qu'est-ce qui manque pour 100% ?",
                "pts": [
                    "8 points manquants : CSP strict non encore activé en production",
                    "Amélioration prévue : nonce-based CSP pour les scripts inline",
                    "L'application reste à un niveau de sécurité élevé selon OWASP Top 10",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 10",
        "title": "Incidents & Réponse",
        "color": RED,
        "qas": [
            {
                "q": "Q28 : Que se passe-t-il si une intrusion est détectée ?",
                "pts": [
                    "Wazuh génère une alerte avec niveau de sévérité (1-15)",
                    "Réponse active possible : script de blocage IP automatique",
                    "Audit log enregistre l'événement pour forensique post-incident",
                ],
            },
            {
                "q": "Q29 : As-tu simulé des attaques sur ton application ?",
                "pts": [
                    "Oui : tests SQLi et XSS bloqués par WAF → HTTP 403 confirmé",
                    "CVE simulée sur Docker : détectée par Wazuh en < 30 secondes",
                    "FIM testé : modification de fichier critique → alerte niveau 7 générée",
                ],
            },
            {
                "q": "Q30 : Comment tu procèdes pour un rollback en cas d'incident ?",
                "pts": [
                    "Backup automatique Firestore avant chaque déploiement",
                    "ROLLBACK.md documente la procédure étape par étape",
                    "Firebase Hosting : rollback en 1 commande firebase hosting:clone",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 11",
        "title": "Monitoring & Performances",
        "color": GREEN,
        "qas": [
            {
                "q": "Q31 : Comment tu surveilles les performances de l'application ?",
                "pts": [
                    "Middleware performance.js : mesure le temps de réponse de chaque route",
                    "Seuil d'alerte : > 2000ms → log SLOW_REQUEST dans auditLogs",
                    "Dashboard Wazuh : métriques système CPU/RAM en temps réel",
                ],
            },
            {
                "q": "Q32 : Combien d'événements Wazuh as-tu analysés ?",
                "pts": [
                    "436 420 événements collectés sur la période de test",
                    "~95% = T1565.001 (Stored Data Manipulation) — FIM très actif",
                    "Rootcheck (Policy Monitoring) : vérification des bonnes pratiques système",
                ],
            },
            {
                "q": "Q33 : Ton application est-elle prête pour la production ?",
                "pts": [
                    "Frontend déployable sur Firebase Hosting (Dockerfile + nginx.conf prêts)",
                    "Backend Functions déjà configuré pour prod (.env + secrets Firebase)",
                    "Pipeline CI/CD GitLab (.gitlab-ci.yml) avec tests automatiques",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 12A",
        "title": "Questions Pièges — Approfondissement (1/2)",
        "color": PURPLE,
        "qas": [
            {
                "q": "Q34 : Quelle différence entre authentification et autorisation ?",
                "pts": [
                    "Authentification : 'qui es-tu ?' → Firebase Auth vérifie le token JWT",
                    "Autorisation : 'que peux-tu faire ?' → RBAC vérifie le rôle Firestore",
                    "Les deux sont indépendants : un token valide ≠ accès autorisé",
                ],
            },
            {
                "q": "Q35 : Pourquoi bcrypt et pas argon2 pour les mots de passe ?",
                "pts": [
                    "Firebase Auth gère le hashing côté serveur (bcrypt par défaut)",
                    "Le choix n'est pas le nôtre : Firebase impose son implémentation",
                    "bcrypt 12 rounds reste sécurisé : ~250ms par hash, anti-brute force efficace",
                ],
            },
            {
                "q": "Q36 : Si ENCRYPTION_KEY est dans .env, comment tu la protèges en prod ?",
                "pts": [
                    "Firebase Secrets Manager : firebase functions:secrets:set ENCRYPTION_KEY",
                    "Jamais dans le code source, jamais dans git (.gitignore inclut .env)",
                    "Rotation possible : nouvelle clé → re-chiffrement des données au déploiement",
                ],
            },
        ],
    },
    {
        "section": "Q&A — BLOC 12B",
        "title": "Questions Pièges — Approfondissement (2/2)",
        "color": PINK,
        "qas": [
            {
                "q": "Q37 : Peux-tu déchiffrer les données si tu perds la clé AES ?",
                "pts": [
                    "Non — aucune récupération possible sans la clé",
                    "C'est une feature, pas un bug : garantit la confidentialité absolue",
                    "Backup sécurisé de la clé requis (coffre-fort, HSM en entreprise)",
                ],
            },
            {
                "q": "Q38 : Quelle est la différence entre Wazuh et un simple pare-feu ?",
                "pts": [
                    "Pare-feu : filtre réseau statique (ports/IP) — couche 3/4",
                    "Wazuh : analyse comportementale dynamique — couche 7 + OS",
                    "Wazuh corrèle les événements dans le temps → détecte les APT",
                ],
            },
            {
                "q": "Q39-40 : Conformité ISO 27001 + Prochain axe d'amélioration ?",
                "pts": [
                    "ISO 27001 partielle : A.9 accès, A.10 crypto, A.12 ops — manque SMSI formel",
                    "Court terme : CSP strict nonce-based → score 100/100 DAST",
                    "Long terme : MFA admin + pentest pro + certification ISO 27001",
                ],
            },
        ],
    },
]

# ── Génération des slides ────────────────────────────────────────────────────

def build_intro_slide(prs, bloc_data, slide_num, total):
    """Slide d'introduction listant les 12 blocs Q&A"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    fill_bg(slide)

    section_tag(slide, "Q&A — QUESTIONS JURY", CYAN)
    header_bar(slide, "Questions Jury — 12 Blocs Thématiques — 40 Questions", PURPLE)

    # Deux colonnes de 6 blocs
    blocs = bloc_data["blocs"]
    cols = [blocs[:6], blocs[6:]]
    col_colors = [CYAN, YELLOW]
    x_positions = [cm(0.5), cm(17)]

    for col_idx, (col, x) in enumerate(zip(cols, x_positions)):
        for i, bloc in enumerate(col):
            y = cm(2.4) + i * cm(1.8)
            add_rect(slide, x, y, cm(15.5), cm(1.5), RGBColor(0x16, 0x2B, 0x40))
            add_text(slide, bloc, x + cm(0.3), y + cm(0.15), cm(15), cm(1.2),
                     size=13, bold=True, color=col_colors[col_idx])

    footer_note(slide, "Format : 3 points par réponse · Mémorisable en ~30s · Toutes réponses vérifiables en live demo")
    slide_number(slide, slide_num, total)


def build_qa_slide(prs, bloc_data, slide_num, total):
    """Slide Q&A avec 3 questions × 3 points chacune"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    fill_bg(slide)

    section_tag(slide, bloc_data["section"], CYAN)
    header_bar(slide, bloc_data["title"], bloc_data["color"])

    qas = bloc_data["qas"]
    # Layout: 3 blocs verticaux
    block_h = cm(3.6)
    y_start = cm(2.3)
    gap = cm(0.3)

    q_colors = [YELLOW, CYAN, GREEN]

    for i, qa in enumerate(qas):
        y = y_start + i * (block_h + gap)
        # Fond du bloc question
        add_rect(slide, cm(0.5), y, cm(32.5), cm(0.8), RGBColor(0x16, 0x2B, 0x40))
        # Question
        add_text(slide, qa["q"], cm(0.7), y + cm(0.05), cm(32), cm(0.75),
                 size=12, bold=True, color=q_colors[i % 3])
        # Points
        pts_y = y + cm(0.85)
        for j, pt in enumerate(qa["pts"]):
            pt_y = pts_y + j * cm(0.75)
            add_text(slide, f"- {pt}", cm(1.0), pt_y, cm(32), cm(0.7),
                     size=11, color=WHITE)

    footer_note(slide, f"Score DAST 92/100 · CDC §3.3 100% · Wazuh 436 420 events · AES-256-CBC · JWT HS256 · bcrypt 12r")
    slide_number(slide, slide_num, total)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    src = "SECURITE_MONITORING_YNOV_2026.pptx.pptx"
    dst = "SECURITE_MONITORING_YNOV_2026_AVEC_QA.pptx"

    prs = Presentation(src)
    print(f"Présentation chargée : {len(prs.slides)} slides")

    # Nombre total de slides final
    n_qa_slides = len(QA_BLOCS)  # 1 intro + 13 blocs = 14
    total = len(prs.slides) + n_qa_slides
    slide_offset = len(prs.slides)

    for idx, bloc in enumerate(QA_BLOCS):
        slide_num = slide_offset + idx + 1
        if bloc.get("intro"):
            build_intro_slide(prs, bloc, slide_num, total)
            print(f"  ✓ Slide {slide_num}/{total} — Intro Q&A")
        else:
            build_qa_slide(prs, bloc, slide_num, total)
            print(f"  ✓ Slide {slide_num}/{total} — {bloc['section']}")

    prs.save(dst)
    import os
    size_ko = os.path.getsize(dst) // 1024
    print(f"\nFichier généré : {dst} ({size_ko} Ko)")
    print(f"Total slides   : {total} (20 originaux + {n_qa_slides} Q&A)")


if __name__ == "__main__":
    main()
