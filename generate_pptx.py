#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génération PPTX — Sécurité & Monitoring — YNOV Campus 2026
Amine BAHOU / Anass Akker
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

BASE = "/Users/anass/Downloads/frais-gestionScolaire 4"

# ── Couleurs ──────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x4B)   # fond principal
BLUE       = RGBColor(0x14, 0x5D, 0xA7)   # bleu moyen
BLUE_LIGHT = RGBColor(0x1E, 0x88, 0xE5)   # bleu clair accent
BLUE_CARD  = RGBColor(0x0D, 0x47, 0xA1)   # carte bleu foncé
CYAN       = RGBColor(0x42, 0xA5, 0xF5)   # titres accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)
GREEN      = RGBColor(0x1E, 0x88, 0xE5)   # bleu clair (statut positif)
ORANGE     = RGBColor(0x15, 0x65, 0xC0)   # bleu moyen (avertissement)
RED        = RGBColor(0x0D, 0x47, 0xA1)   # bleu foncé (critique)

W  = Inches(13.33)   # largeur slide 16:9
H  = Inches(7.5)     # hauteur slide 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # layout vide

# ─────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=NAVY, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    fill_obj = shape.fill
    fill_obj.solid()
    fill_obj.fore_color.rgb = fill
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False,
             wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb

def slide_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title, subtitle=None):
    """Barre de titre bleue en haut du slide."""
    add_rect(slide, 0, 0, W, Inches(1.3), fill=BLUE_CARD)
    add_text(slide, title,
             Inches(0.4), Inches(0.12), Inches(12), Inches(0.75),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.82), Inches(12), Inches(0.4),
                 size=13, bold=False, color=CYAN, align=PP_ALIGN.LEFT)

def badge(slide, text, x, y, w, h, bg=BLUE_CARD, color=WHITE, size=11, bold=True):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, text, x + Inches(0.05), y + Inches(0.04),
             w - Inches(0.1), h - Inches(0.08),
             size=size, bold=bold, color=color, align=PP_ALIGN.CENTER)

def try_image(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
        return True
    else:
        add_rect(slide, x, y, w, h, fill=BLUE_CARD)
        add_text(slide, f"[Image: {os.path.basename(path)}]",
                 x + Inches(0.1), y + h//2 - Inches(0.2),
                 w - Inches(0.2), Inches(0.4),
                 size=10, color=CYAN, align=PP_ALIGN.CENTER)
        return False

def footer(slide, text):
    add_text(slide, text, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35),
             size=9, italic=True, color=CYAN, align=PP_ALIGN.LEFT)

def separator(slide, y, color=BLUE_LIGHT):
    add_rect(slide, Inches(0.3), y, Inches(12.73), Inches(0.03), fill=color)

# ─────────────────────────────────────────────────────────
# SLIDE 1 — TITRE
# ─────────────────────────────────────────────────────────
def slide_01():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)

    # Bande du haut
    add_rect(sl, 0, 0, W, Inches(0.4), fill=BLUE)
    add_text(sl, "PFE CYBERSÉCURITÉ / CYBERDÉFENSE",
             Inches(0.3), Inches(0.05), Inches(12), Inches(0.3),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Titre principal
    add_text(sl, "SÉCURITÉ & MONITORING",
             Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.5),
             size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Sous-titre
    add_text(sl, "Application Web — Gestion Scolaire YNOV Campus",
             Inches(1), Inches(2.15), Inches(11.3), Inches(0.6),
             size=22, bold=False, color=CYAN, align=PP_ALIGN.CENTER)

    add_text(sl, "Phase 3 — Post-Production  ·  WAF  ·  SIEM  ·  DAST  ·  Wazuh 4.7.4  ·  CDC §3.3 à 100%",
             Inches(1), Inches(2.7), Inches(11.3), Inches(0.4),
             size=13, italic=True, color=CYAN, align=PP_ALIGN.CENTER)

    separator(sl, Inches(3.2))

    # Tags
    tags = ["OWASP Top 10", "RGPD UE 2016/679", "CDC §3.3", "Wazuh 4.7.4"]
    tw = Inches(2.8)
    gap = Inches(0.2)
    start = (W - 4*tw - 3*gap) / 2
    for i, t in enumerate(tags):
        badge(sl, t, start + i*(tw+gap), Inches(3.4), tw, Inches(0.5),
              bg=BLUE, size=12)

    # KPI cards
    kpis = [
        ("10", "Catégories OWASP"),
        ("6", "Articles RGPD"),
        ("12", "Tests DAST"),
        ("2 500+", "FIM Events"),
        ("17", "CVE Détectées"),
        ("1 600+", "MITRE ATT&CK"),
    ]
    cw = Inches(1.95)
    ch = Inches(1.1)
    cx_start = Inches(0.4)
    cy = Inches(4.1)
    for i, (num, lbl) in enumerate(kpis):
        cx = cx_start + i * (cw + Inches(0.15))
        add_rect(sl, cx, cy, cw, ch, fill=BLUE_LIGHT)
        add_text(sl, num, cx, cy + Inches(0.05), cw, Inches(0.55),
                 size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, cx, cy + Inches(0.6), cw, Inches(0.45),
                 size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Auteurs
    add_text(sl, "Amine BAHOU  /  Anass Akker  —  YNOV Campus 2026",
             Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45),
             size=13, bold=False, color=CYAN, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
# SLIDE 2 — CONTEXTE & OBJECTIFS
# ─────────────────────────────────────────────────────────
def slide_02():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "CONTEXTE & OBJECTIFS DU PROJET", "02 / 20  —  Périmètre PFE")

    # Colonne gauche
    add_rect(sl, Inches(0.3), Inches(1.5), Inches(5.8), Inches(2.1), fill=BLUE_CARD)
    add_text(sl, "Application Web", Inches(0.4), Inches(1.55), Inches(5.6), Inches(0.45),
             size=15, bold=True, color=CYAN)
    items_app = [
        "Stack : React (frontend) + Node.js (backend)",
        "Base de données : Firebase / Firestore",
        "Déploiement : Docker + Firebase Functions",
        "Utilisateurs : admin, sous-admin, comptable, étudiant/parent",
    ]
    for j, it in enumerate(items_app):
        add_text(sl, f"▸  {it}", Inches(0.45), Inches(2.0) + j*Inches(0.38),
                 Inches(5.6), Inches(0.36), size=11, color=WHITE)

    add_rect(sl, Inches(0.3), Inches(3.75), Inches(5.8), Inches(2.1), fill=BLUE_CARD)
    add_text(sl, "Objectifs PFE — Phase 3", Inches(0.4), Inches(3.8), Inches(5.6), Inches(0.45),
             size=15, bold=True, color=CYAN)
    items_obj = [
        "Sécuriser l'application web (OWASP Top 10)",
        "Implémenter WAF, RBAC, audit logs immuables",
        "Assurer la conformité RGPD (6 articles)",
        "Déployer Wazuh SIEM — surveillance infrastructure",
        "Atteindre CDC §3.3 à 100% (8/8 exigences)",
    ]
    for j, it in enumerate(items_obj):
        add_text(sl, f"✓  {it}", Inches(0.45), Inches(4.25) + j*Inches(0.32),
                 Inches(5.6), Inches(0.3), size=11, color=WHITE)

    # Colonne droite — frameworks
    add_rect(sl, Inches(6.4), Inches(1.5), Inches(6.6), Inches(0.55), fill=BLUE)
    add_text(sl, "Frameworks & Standards Appliqués", Inches(6.45), Inches(1.53),
             Inches(6.5), Inches(0.5), size=14, bold=True, color=WHITE)

    standards = [
        ("OWASP Top 10 (2021)", "10 catégories vérifiées — 8 corrigées, 1 surveillée, 1 N/A"),
        ("RGPD UE 2016/679", "6 articles implémentés — Art.15/16/17/25/32/33"),
        ("CDC §3.3", "8/8 exigences couvertes — toutes vérifiables en live"),
        ("Wazuh 4.7.4", "SIEM + FIM + CVE + MITRE ATT&CK + GDPR module"),
        ("MITRE ATT&CK", "1 600+ alertes classifiées — T1499/T1562/T1565"),
        ("CIS Benchmark", "Debian 12 — audit configuration automatique"),
    ]
    for j, (std, desc) in enumerate(standards):
        add_rect(sl, Inches(6.4), Inches(2.2) + j*Inches(0.83),
                 Inches(6.6), Inches(0.75), fill=BLUE_CARD)
        add_text(sl, std, Inches(6.5), Inches(2.23) + j*Inches(0.83),
                 Inches(6.4), Inches(0.33), size=12, bold=True, color=CYAN)
        add_text(sl, desc, Inches(6.5), Inches(2.54) + j*Inches(0.83),
                 Inches(6.4), Inches(0.35), size=10, color=WHITE)

    footer(sl, "Projet réel déployé — Stack moderne — Conformité légale et technique prouvée.")

# ─────────────────────────────────────────────────────────
# SLIDE 3 — PIPELINE DE SÉCURITÉ (Architecture)
# ─────────────────────────────────────────────────────────
def slide_03():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "PIPELINE DE SÉCURITÉ — 8 COUCHES", "03 / 20  —  De la Requête à l'Alerte SIEM")

    layers = [
        ("WAF (waf.js)",          "SQLi · XSS · PathTraversal · CmdInj · Agents suspects  →  HTTP 403 + WAF_BLOCK",  BLUE_LIGHT),
        ("Rate Limiter",          "100 req/15min global · 5 échecs  →  AUTH_LOCKOUT  →  HTTP 429",                    BLUE),
        ("verifyJWT",             "JWT HS256 · expiration 24h · signature  →  HTTP 401 si invalide",                   BLUE_CARD),
        ("checkRole() Firestore", "Re-lecture rôle BDD à chaque requête  →  HTTP 403 + ACCESS_DENIED si insuffisant", BLUE_LIGHT),
        ("Route Handler",         "Traitement métier sécurisé  →  Réponse applicative",                                BLUE),
        ("auditLogger",           "9 types events · Firestore immuable · serverTimestamp()",                           BLUE_CARD),
        ("Dashboard /monitoring", "Score /100 · WAF stats · SIEM logs · 4 alertes auto · 60s refresh",                BLUE_LIGHT),
        ("Wazuh SIEM",            "CIS Benchmark · MITRE ATT&CK · FIM · CVE · GDPR · Agent actif",                   BLUE),
    ]

    # Flèche verticale gauche
    add_rect(sl, Inches(0.35), Inches(1.45), Inches(0.08), Inches(5.7), fill=CYAN)

    for i, (name, desc, col) in enumerate(layers):
        y = Inches(1.45) + i * Inches(0.72)
        # Bloc nom
        add_rect(sl, Inches(0.5), y + Inches(0.05), Inches(2.8), Inches(0.58), fill=col)
        add_text(sl, name, Inches(0.55), y + Inches(0.09), Inches(2.7), Inches(0.46),
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        # Description
        add_text(sl, desc, Inches(3.45), y + Inches(0.12), Inches(9.5), Inches(0.46),
                 size=11, color=WHITE)
        # Flèche
        if i < 7:
            add_text(sl, "▼", Inches(0.3), y + Inches(0.62), Inches(0.2), Inches(0.25),
                     size=10, color=CYAN, align=PP_ALIGN.CENTER)

    footer(sl, "Chaque requête filtrée par 4 couches avant d'accéder aux données. Toutes les actions loggées et analysées via Wazuh.")

# ─────────────────────────────────────────────────────────
# SLIDE 4 — OWASP TOP 10
# ─────────────────────────────────────────────────────────
def slide_04():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "AUDIT OWASP TOP 10 (2021) — 10 Catégories Vérifiées", "04 / 20")

    # En-tête table
    cols_w = [Inches(3.2), Inches(6.8), Inches(1.9)]
    cols_x = [Inches(0.3), Inches(3.55), Inches(10.4)]
    headers = ["CATÉGORIE OWASP", "IMPLÉMENTATION", "STATUT"]
    for j, (h, cx, cw) in enumerate(zip(headers, cols_x, cols_w)):
        add_rect(sl, cx, Inches(1.45), cw - Inches(0.05), Inches(0.45), fill=BLUE)
        add_text(sl, h, cx + Inches(0.05), Inches(1.47), cw - Inches(0.1), Inches(0.4),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("A01 — Broken Access Control",     "RBAC · checkRole() · Firestore Rules",                                "CORRIGÉ",    GREEN),
        ("A02 — Cryptographic Failures",     "AES-256-CBC · HTTPS · bcrypt 12 rounds",                             "CORRIGÉ",    GREEN),
        ("A03 — Injection SQLi·XSS·Cmd",    "WAF waf.js  →  HTTP 403 + WAF_BLOCK",                                "CORRIGÉ",    GREEN),
        ("A04 — Insecure Design",            "Architecture JWT + RBAC dès la conception",                          "CORRIGÉ",    GREEN),
        ("A05 — Security Misconfiguration",  "Helmet · CORS strict · Headers HTTP",                                "CORRIGÉ",    GREEN),
        ("A06 — Vulnerable Components",      "Wazuh Vulnerability Detection — 17 CVE",                             "SURVEILLÉ",  ORANGE),
        ("A07 — Auth & Session Failures",    "Rate limiting · lockout 5 échecs · JWT 24h",                         "CORRIGÉ",    GREEN),
        ("A08 — Data Integrity Failures",    "auditLogs immuables · JWT signé HS256",                              "CORRIGÉ",    GREEN),
        ("A09 — Security Logging Failures",  "9 types audit logs Firestore · SIEM",                                "CORRIGÉ",    GREEN),
        ("A10 — SSRF",                       "Non applicable — pas d'appels serveur→serveur",                      "N/A",        BLUE_LIGHT),
    ]

    for i, (cat, impl, status, scol) in enumerate(rows):
        row_bg = BLUE_CARD if i % 2 == 0 else RGBColor(0x12, 0x30, 0x6B)
        y = Inches(1.95) + i * Inches(0.5)
        for j, (cx, cw) in enumerate(zip(cols_x, cols_w)):
            add_rect(sl, cx, y, cw - Inches(0.05), Inches(0.46), fill=row_bg)
        add_text(sl, cat,    cols_x[0]+Inches(0.05), y+Inches(0.05), cols_w[0]-Inches(0.1), Inches(0.38), size=10.5, color=WHITE)
        add_text(sl, impl,   cols_x[1]+Inches(0.05), y+Inches(0.05), cols_w[1]-Inches(0.1), Inches(0.38), size=10.5, color=WHITE)
        # Badge statut
        badge(sl, status, cols_x[2]+Inches(0.1), y+Inches(0.06), cols_w[2]-Inches(0.2), Inches(0.33),
              bg=scol, size=10)

    footer(sl, "8 CORRIGÉES · 1 SURVEILLÉE par Wazuh · 1 N/A — Toutes failles critiques couvertes.")

# ─────────────────────────────────────────────────────────
# SLIDE 5 — WAF
# ─────────────────────────────────────────────────────────
def slide_05():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "WAF waf.js — 5 Vecteurs OWASP Bloqués · Première Ligne de Défense", "05 / 20  —  Sécurité Applicative")

    rows_waf = [
        ("Injection SQL",         "CRITIQUE", "SELECT/UNION/DROP/OR 1=1  →  HTTP 403 + WAF_BLOCK loggué",           RED),
        ("XSS Cross-Site Scripting","HAUTE",  "<script> · eval() · onerror=  →  HTTP 403 · CSP default-src 'self'", ORANGE),
        ("Path Traversal",        "HAUTE",    "../../ · %2e%2e · /etc/passwd  →  HTTP 403 + WAF_BLOCK",             ORANGE),
        ("Command Injection",     "HAUTE",    "; && | $()  →  HTTP 403 + WAF_BLOCK loggué",                         ORANGE),
        ("Scanners suspects",     "HAUTE",    "sqlmap/nikto/nmap User-Agent  →  HTTP 403 bloqué",                   ORANGE),
    ]

    # En-têtes
    hdrs = [("TYPE D'ATTAQUE", Inches(2.8)), ("SÉVÉRITÉ", Inches(1.3)), ("ACTION WAF", Inches(8.7))]
    hx = [Inches(0.3), Inches(3.15), Inches(4.5)]
    for txt, w, x in zip([h[0] for h in hdrs], [h[1] for h in hdrs], hx):
        add_rect(sl, x, Inches(1.5), w-Inches(0.05), Inches(0.42), fill=BLUE)
        add_text(sl, txt, x+Inches(0.05), Inches(1.52), w-Inches(0.1), Inches(0.36),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (attack, sev, action, scol) in enumerate(rows_waf):
        bg = BLUE_CARD if i % 2 == 0 else RGBColor(0x12, 0x30, 0x6B)
        y = Inches(1.97) + i * Inches(0.52)
        for x, w in zip(hx, [h[1] for h in hdrs]):
            add_rect(sl, x, y, w-Inches(0.05), Inches(0.48), fill=bg)
        add_text(sl, attack, hx[0]+Inches(0.05), y+Inches(0.07), Inches(2.7), Inches(0.36),
                 size=11, bold=True, color=WHITE)
        badge(sl, sev, hx[1]+Inches(0.05), y+Inches(0.07), Inches(1.2), Inches(0.33), bg=scol, size=10)
        add_text(sl, action, hx[2]+Inches(0.05), y+Inches(0.07), Inches(8.5), Inches(0.36), size=11, color=WHITE)

    # Code exemple
    add_rect(sl, Inches(0.3), Inches(4.65), Inches(12.73), Inches(0.58), fill=RGBColor(0x05, 0x0F, 0x2E))
    add_text(sl, "# SQLi bloqué  →  curl '...?id=1 OR 1=1'  →  HTTP 403   |   XSS  →  curl '...?nom=<script>alert(1)</script>'  →  HTTP 403",
             Inches(0.45), Inches(4.7), Inches(12.5), Inches(0.45),
             size=10.5, color=GREEN, font="Courier New")

    # Brute Force box
    add_rect(sl, Inches(0.3), Inches(5.35), Inches(12.73), Inches(1.3), fill=BLUE_CARD)
    add_text(sl, "BRUTE FORCE PROTECTION", Inches(0.45), Inches(5.4), Inches(12), Inches(0.38),
             size=13, bold=True, color=WHITE)
    add_text(sl,
             "Rate limit 10 req/15min  ·  5 échecs  →  AUTH_LOCKOUT Firestore  →  HTTP 429  ·  Payload max 10 KB  →  HTTP 413  ·  Timeout 28s  →  HTTP 503",
             Inches(0.45), Inches(5.8), Inches(12.4), Inches(0.75),
             size=12, color=WHITE)

    footer(sl, "WAF appliqué avant même l'authentification — première ligne de défense. Toutes vulnérabilités CRITIQUES et HAUTES corrigées.")

# ─────────────────────────────────────────────────────────
# SLIDE 6 — RBAC
# ─────────────────────────────────────────────────────────
def slide_06():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "RBAC — Gestion des Accès par Rôle — Moindre Privilège", "06 / 20  —  CDC §3.3")

    roles = [
        ("admin",           "Accès complet",          ["monitoring", "users", "toutes données", "configuration", "exports"],               BLUE_LIGHT),
        ("sous-admin",      "Gestion opérationnelle",  ["étudiants", "classes", "paiements", "factures", "rapports"],                       BLUE),
        ("comptable",       "Finance uniquement",      ["lecture paiements", "génération factures", "exports uniquement"],                  BLUE_CARD),
        ("étudiant/parent", "Portail personnel",       ["propres données uniquement", "lecture seule", "pas d'admin"],                      RGBColor(0x1A, 0x4D, 0x8A)),
    ]

    rw = Inches(3.0)
    rh = Inches(4.2)
    gap = Inches(0.22)
    rx_start = Inches(0.3)

    for i, (role, subtitle, items, col) in enumerate(roles):
        rx = rx_start + i * (rw + gap)
        ry = Inches(1.5)
        add_rect(sl, rx, ry, rw, rh, fill=col)
        # Titre rôle
        add_rect(sl, rx, ry, rw, Inches(0.55), fill=RGBColor(0x0A, 0x20, 0x60))
        add_text(sl, role, rx, ry + Inches(0.08), rw, Inches(0.45),
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, subtitle, rx, ry + Inches(0.6), rw, Inches(0.35),
                 size=10, italic=True, color=GREY_LIGHT, align=PP_ALIGN.CENTER)
        for k, item in enumerate(items):
            add_text(sl, f"▸  {item}", rx + Inches(0.1), ry + Inches(1.05) + k*Inches(0.47),
                     rw - Inches(0.15), Inches(0.42), size=11, color=WHITE)

    # Note implémentation
    add_rect(sl, Inches(0.3), Inches(5.85), Inches(12.73), Inches(0.7), fill=BLUE_CARD)
    add_text(sl,
             "Implémentation : checkRole(['admin','sous-admin']) sur chaque route  ·  Re-lecture rôle Firestore à CHAQUE requête (jamais depuis JWT uniquement)",
             Inches(0.45), Inches(5.9), Inches(12.4), Inches(0.6),
             size=11, color=WHITE)

    footer(sl, "RBAC vérifie les rôles en temps réel. Si rôle révoqué → bloqué immédiatement à la prochaine requête.")

# ─────────────────────────────────────────────────────────
# SLIDE 7 — RGPD
# ─────────────────────────────────────────────────────────
def slide_07():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "CONFORMITÉ RGPD (UE 2016/679) — 6 Articles Implémentés", "07 / 20  —  CDC §3.3")

    articles = [
        ("Art. 15 — Droit d'accès",         "Export données personnelles\nDATA_EXPORT loggué · GET /users/:id/export"),
        ("Art. 16 — Rectification",          "Modification données\nCRUD loggué avec serverTimestamp()"),
        ("Art. 17 — Droit à l'effacement",   "Anonymisation DATA_ANONYMIZE\nchamps remplacés · log conservé"),
        ("Art. 25 — Privacy by Design",      "AES-256-CBC · HTTPS · RBAC\nprincipe du moindre privilège dès conception"),
        ("Art. 32 — Sécurité du traitement", "WAF + JWT HS256 + bcrypt 12 rounds\nrate limiting + HTTPS/HSTS max-age=31536000"),
        ("Art. 33 — Notification violation", "Alerte CRITIQUE auditLogs\nProcédure 72h documentée"),
    ]

    aw = Inches(4.1)
    ah = Inches(2.1)
    gap = Inches(0.22)
    ax_starts = [Inches(0.3), Inches(4.62), Inches(8.94)]
    ay_starts = [Inches(1.5), Inches(3.75)]

    for i, (art, desc) in enumerate(articles):
        row = i // 3
        col = i % 3
        ax = ax_starts[col]
        ay = ay_starts[row]
        add_rect(sl, ax, ay, aw, ah, fill=BLUE_CARD)
        add_text(sl, art, ax + Inches(0.1), ay + Inches(0.1), aw - Inches(0.2), Inches(0.45),
                 size=12, bold=True, color=CYAN)
        add_text(sl, desc, ax + Inches(0.1), ay + Inches(0.6), aw - Inches(0.2), Inches(1.4),
                 size=11, color=WHITE)

    footer(sl, "Privacy by Design (Art.25) dès la conception. Procédure 72h (Art.33) automatisée via le dashboard /monitoring.")

# ─────────────────────────────────────────────────────────
# SLIDE 8 — JOURNALISATION & DAST
# ─────────────────────────────────────────────────────────
def slide_08():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "9 Audit Logs Immuables + Security Headers + DAST 92/100", "08 / 20  —  CDC §3.3 — Journalisation")

    logs = [
        ("AUTH_SUCCESS",    "INFO",     "userId + IP + horodatage",               GREEN),
        ("AUTH_FAILURE",    "WARNING",  "userId + IP + tentative n°",             ORANGE),
        ("AUTH_LOCKOUT",    "CRITIQUE", "userId + IP + durée blocage",            RED),
        ("LOGOUT",          "INFO",     "userId + durée session",                 GREEN),
        ("SESSION_EXPIRED", "INFO",     "userId + last_activity",                 GREEN),
        ("ACCESS_DENIED",   "WARNING",  "userId + route + rôle tenté",            ORANGE),
        ("DATA_EXPORT",     "INFO",     "userId + données exportées — Art.15",    GREEN),
        ("DATA_ANONYMIZE",  "WARNING",  "userId + champs anonymisés — Art.17",    ORANGE),
        ("WAF_BLOCK",       "CRITIQUE", "IP + type attaque + payload",            RED),
    ]

    # Table logs (gauche)
    hx = [Inches(0.3), Inches(2.95), Inches(4.6)]
    hw = [Inches(2.6), Inches(1.6), Inches(5.5)]
    hdrs = ["EVENT", "NIVEAU", "DONNÉES LOGGÉES"]
    for j, (h, x, w) in enumerate(zip(hdrs, hx, hw)):
        add_rect(sl, x, Inches(1.5), w-Inches(0.04), Inches(0.38), fill=BLUE)
        add_text(sl, h, x+Inches(0.05), Inches(1.52), w-Inches(0.1), Inches(0.32),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (ev, lvl, data, col) in enumerate(logs):
        bg = BLUE_CARD if i % 2 == 0 else RGBColor(0x12, 0x30, 0x6B)
        y = Inches(1.93) + i * Inches(0.5)
        for x, w in zip(hx, hw):
            add_rect(sl, x, y, w-Inches(0.04), Inches(0.46), fill=bg)
        add_text(sl, ev,   hx[0]+Inches(0.05), y+Inches(0.06), hw[0]-Inches(0.1), Inches(0.34), size=9.5, bold=True, color=CYAN, font="Courier New")
        badge(sl, lvl, hx[1]+Inches(0.05), y+Inches(0.06), hw[1]-Inches(0.15), Inches(0.33), bg=col, size=8.5)
        add_text(sl, data, hx[2]+Inches(0.05), y+Inches(0.06), hw[2]-Inches(0.1), Inches(0.34), size=9.5, color=WHITE)

    # Firestore rules note
    add_rect(sl, Inches(0.3), Inches(6.52), Inches(9.85), Inches(0.5), fill=RGBColor(0x05, 0x0F, 0x2E))
    add_text(sl, "Firestore Rules : allow update: if false · allow delete: if false  →  LOGS IMMUABLES · serverTimestamp() non manipulable côté client",
             Inches(0.4), Inches(6.55), Inches(9.6), Inches(0.42), size=9.5, color=GREEN)

    # Droite — Headers & DAST
    add_rect(sl, Inches(10.35), Inches(1.5), Inches(2.65), Inches(2.5), fill=BLUE_CARD)
    add_text(sl, "Security Headers", Inches(10.45), Inches(1.55), Inches(2.5), Inches(0.4),
             size=12, bold=True, color=CYAN)
    headers = ["HSTS max-age=31536000", "X-Frame-Options DENY", "CSP default-src 'self'",
               "JWT exp:24h · HS256", "bcrypt 12 rounds", "AES-256-CBC données"]
    for k, h in enumerate(headers):
        add_text(sl, f"✓  {h}", Inches(10.45), Inches(2.0) + k*Inches(0.33),
                 Inches(2.5), Inches(0.3), size=9.5, color=WHITE)

    add_rect(sl, Inches(10.35), Inches(4.1), Inches(2.65), Inches(2.3), fill=BLUE_CARD)
    add_text(sl, "Scanner DAST", Inches(10.45), Inches(4.15), Inches(2.5), Inches(0.4),
             size=12, bold=True, color=CYAN)
    add_text(sl, "12 Tests OWASP", Inches(10.45), Inches(4.6), Inches(2.5), Inches(0.35),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    badge(sl, "SCORE 92/100", Inches(10.5), Inches(5.0), Inches(2.4), Inches(0.5), bg=GREEN, size=14)
    add_text(sl, "11/12 PASS  ·  1 warning CORS", Inches(10.45), Inches(5.6),
             Inches(2.5), Inches(0.35), size=9, color=WHITE, align=PP_ALIGN.CENTER)

    footer(sl, "Logs infalsifiables — serverTimestamp() côté serveur. DAST automatisé intégrable CI/CD.")

# ─────────────────────────────────────────────────────────
# SLIDE 9 — SCREENSHOT Dashboard Principal
# ─────────────────────────────────────────────────────────
def slide_09():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "DÉMONSTRATION — Dashboard Sécurité Principal", "09 / 20  —  Score · Auth · RGPD · RBAC · Audit")

    img_path = os.path.join(BASE, "DASH .png")
    try_image(sl, img_path, Inches(0.3), Inches(1.5), Inches(7.8), Inches(5.5))

    # Commentaires droite
    add_rect(sl, Inches(8.3), Inches(1.5), Inches(4.7), Inches(5.5), fill=BLUE_CARD)
    add_text(sl, "Dashboard Sécurité Complet", Inches(8.4), Inches(1.6), Inches(4.5), Inches(0.45),
             size=14, bold=True, color=WHITE)
    add_text(sl, "Auth · RGPD · RBAC · Audit · Score 50/100",
             Inches(8.4), Inches(2.05), Inches(4.5), Inches(0.35), size=11, color=CYAN)

    points = [
        ("Accès Sécurisé",    "3 connexions réussies · 5 échecs · 5 comptes bloqués · 3 sessions expirées"),
        ("RGPD",              "12 exports Art.15 · 8 anonymisations Art.17 — conformité totale"),
        ("RBAC",              "10 accès refusés en 24h · 4 dans la dernière heure"),
        ("Audit Logs",        "92 événements tracés — immuables, conservés 1 an"),
        ("Protections actives","JWT HS256 · Brute force · Rate limiting · bcrypt 12r · Anti-énumération"),
    ]
    for k, (label, text) in enumerate(points):
        y = Inches(2.55) + k * Inches(0.55)
        add_text(sl, f"✓  {label}", Inches(8.4), y, Inches(4.5), Inches(0.3),
                 size=11, bold=True, color=GREEN)
        add_text(sl, text, Inches(8.4), y + Inches(0.28), Inches(4.5), Inches(0.28),
                 size=9.5, color=WHITE)

    badge(sl, "TOUT DÉTECTÉ · BLOQUÉ · TRACÉ EN TEMPS RÉEL",
          Inches(8.3), Inches(6.65), Inches(4.7), Inches(0.4), bg=GREEN, size=10)

    footer(sl, "Score 50/100 — alerte active WAF en cours lors de cette capture (44 attaques bloquées).")

# ─────────────────────────────────────────────────────────
# SLIDE 10 — SCREENSHOT WAF Monitoring
# ─────────────────────────────────────────────────────────
def slide_10():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "DÉMONSTRATION — WAF Monitoring · 44 Attaques Bloquées", "10 / 20  —  Pare-feu Applicatif")

    img_path = os.path.join(BASE, "WAFCAPTSIEM", "WAF.png")
    try_image(sl, img_path, Inches(0.3), Inches(1.5), Inches(7.8), Inches(5.5))

    add_rect(sl, Inches(8.3), Inches(1.5), Inches(4.7), Inches(5.5), fill=BLUE_CARD)
    add_text(sl, "WAF — Résultats Live", Inches(8.4), Inches(1.6), Inches(4.5), Inches(0.45),
             size=14, bold=True, color=WHITE)

    stats = [
        ("44",  "Attaques bloquées (24h)", BLUE_LIGHT),
        ("6",   "Bloquées (1h)",           BLUE),
        ("44",  "Total WAF_BLOCK logs",    BLUE_CARD),
    ]
    for k, (num, lbl, col) in enumerate(stats):
        cx = Inches(8.35) + k * Inches(1.55)
        add_rect(sl, cx, Inches(2.15), Inches(1.45), Inches(0.9), fill=col)
        add_text(sl, num, cx, Inches(2.18), Inches(1.45), Inches(0.5),
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, cx, Inches(2.67), Inches(1.45), Inches(0.35),
                 size=8, color=WHITE, align=PP_ALIGN.CENTER)

    types_att = [
        ("Injection SQL",    "12 (27%)", RED),
        ("XSS",             "12 (27%)", ORANGE),
        ("Path Traversal",  "8 (18%)",  ORANGE),
        ("Scanner",         "6 (14%)",  BLUE_LIGHT),
        ("Autre",           "6 (14%)",  BLUE),
    ]
    add_text(sl, "Répartition par type :", Inches(8.4), Inches(3.2), Inches(4.5), Inches(0.35),
             size=11, bold=True, color=CYAN)
    for k, (t, cnt, col) in enumerate(types_att):
        y = Inches(3.6) + k * Inches(0.42)
        add_rect(sl, Inches(8.4), y, Inches(2.8), Inches(0.35), fill=col)
        add_text(sl, t,   Inches(8.45), y+Inches(0.04), Inches(2.0), Inches(0.28), size=10, color=WHITE)
        add_text(sl, cnt, Inches(10.3), y+Inches(0.04), Inches(0.9), Inches(0.28), size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    footer(sl, "WAF Firestore auditLogs — Détection SQL · XSS · Path Traversal · CMD Injection · Agents suspects. Alerte brute force active.")

# ─────────────────────────────────────────────────────────
# SLIDE 11 — SCREENSHOT Dashboard SIEM
# ─────────────────────────────────────────────────────────
def slide_11():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "DÉMONSTRATION — Dashboard SIEM · Journal des Événements", "11 / 20  —  Security Information & Event Management")

    img_path = os.path.join(BASE, "WAFCAPTSIEM", "SIEM.png")
    try_image(sl, img_path, Inches(0.3), Inches(1.5), Inches(7.8), Inches(5.5))

    add_rect(sl, Inches(8.3), Inches(1.5), Inches(4.7), Inches(5.5), fill=BLUE_CARD)
    add_text(sl, "SIEM — Agrégation Temps Réel", Inches(8.4), Inches(1.6), Inches(4.5), Inches(0.45),
             size=14, bold=True, color=WHITE)

    counters = [
        ("18",  "Événements Auth",   GREEN),
        ("10",  "Événements RBAC",   BLUE_LIGHT),
        ("20",  "Événements RGPD",   CYAN),
        ("44",  "Blocages WAF",      RED),
    ]
    for k, (n, l, c) in enumerate(counters):
        row, col = divmod(k, 2)
        cx = Inches(8.35) + col * Inches(2.3)
        cy = Inches(2.2) + row * Inches(0.85)
        add_rect(sl, cx, cy, Inches(2.15), Inches(0.75), fill=c)
        add_text(sl, n, cx, cy+Inches(0.04), Inches(2.15), Inches(0.42),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, l, cx, cy+Inches(0.44), Inches(2.15), Inches(0.28),
                 size=9, color=WHITE, align=PP_ALIGN.CENTER)

    integrity = ["allow update: if false", "allow delete: if false",
                 "Lecture réservée aux admins", "Conservation : 1 an",
                 "serverTimestamp() serveur", "9 types d'événements"]
    add_text(sl, "Intégrité du système :", Inches(8.4), Inches(4.1), Inches(4.5), Inches(0.35),
             size=11, bold=True, color=CYAN)
    for k, item in enumerate(integrity):
        add_text(sl, f"✓  {item}", Inches(8.4), Inches(4.48) + k*Inches(0.33),
                 Inches(4.5), Inches(0.3), size=10, color=WHITE)

    footer(sl, "Dashboard SIEM — 20 derniers événements · Rafraîchissement automatique 60s · Audit Firestore immuable.")

# ─────────────────────────────────────────────────────────
# SLIDE 12 — WAZUH ARCHITECTURE
# ─────────────────────────────────────────────────────────
def slide_12():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "WAZUH 4.7.4 — Architecture Docker · 2 Agents · Coverage 100%", "12 / 20  —  Infrastructure SIEM")

    layers_w = [
        ("Application Web (React)", "Frontend · auditLogs Firestore · dashboard /monitoring"),
        ("API Backend (Node.js)",   "WAF · JWT · RBAC · Handler · auditLogger"),
        ("Wazuh Agent",             "OSSEC · Port 1514 · Chiffrement AES · Push events → Manager"),
        ("Wazuh Manager",           "Collecte · corrélation · moteur règles · Port 1514/55000"),
        ("Wazuh Indexer (OpenSearch)","Stockage alertes · Port 9200 · Full-text search"),
        ("Wazuh Dashboard",         "Interface Kibana fork · Port 443 HTTPS TLS"),
    ]

    for i, (name, desc) in enumerate(layers_w):
        y = Inches(1.5) + i * Inches(0.8)
        col = BLUE_LIGHT if i % 2 == 0 else BLUE
        add_rect(sl, Inches(0.3), y, Inches(3.0), Inches(0.68), fill=col)
        add_text(sl, name, Inches(0.35), y+Inches(0.1), Inches(2.9), Inches(0.5),
                 size=11, bold=True, color=WHITE)
        add_text(sl, desc, Inches(3.45), y+Inches(0.15), Inches(5.4), Inches(0.44),
                 size=10.5, color=WHITE)
        if i < 5:
            add_text(sl, "▼", Inches(0.85), y + Inches(0.68), Inches(0.3), Inches(0.2),
                     size=9, color=CYAN, align=PP_ALIGN.CENTER)

    # Colonne droite — Install + Agents
    add_rect(sl, Inches(9.1), Inches(1.5), Inches(3.9), Inches(3.2), fill=RGBColor(0x05, 0x0F, 0x2E))
    add_text(sl, "Installation — 4 Commandes", Inches(9.2), Inches(1.55), Inches(3.7), Inches(0.38),
             size=12, bold=True, color=CYAN)
    cmds = [
        "git clone -b v4.7.4 github.com/wazuh/wazuh-docker.git",
        "cd wazuh-docker/single-node",
        "docker compose -f generate-indexer-certs.yml run --rm generator",
        "docker compose up -d",
        "→ https://localhost : admin / SecretPassword",
    ]
    for k, cmd in enumerate(cmds):
        c = CYAN if cmd.startswith("→") else GREEN
        add_text(sl, cmd, Inches(9.15), Inches(2.0) + k*Inches(0.47),
                 Inches(3.7), Inches(0.42), size=8.5, color=c, font="Courier New")

    add_rect(sl, Inches(9.1), Inches(4.85), Inches(3.9), Inches(2.0), fill=BLUE_CARD)
    add_text(sl, "Agents Actifs", Inches(9.2), Inches(4.9), Inches(3.7), Inches(0.38),
             size=12, bold=True, color=CYAN)

    agents = [
        ("Agent 001", "frais-gestion-scolaire · Debian 12 · Active"),
        ("Agent 002", "main-machine · macOS 15.7.4 · Active"),
    ]
    for k, (ag, desc) in enumerate(agents):
        add_text(sl, ag, Inches(9.2), Inches(5.35) + k*Inches(0.6), Inches(1.2), Inches(0.4),
                 size=11, bold=True, color=WHITE)
        add_text(sl, desc, Inches(10.45), Inches(5.35) + k*Inches(0.6), Inches(2.5), Inches(0.4),
                 size=9.5, color=GREY_LIGHT)

    badge(sl, "COVERAGE 100% — 2 actifs · 0 déconnecté",
          Inches(9.1), Inches(6.6), Inches(3.9), Inches(0.42), bg=GREEN, size=10)

    # Screenshot agents — interface Wazuh Agents (capture fournie)
    img = os.path.join(BASE, "WAZUCAPT", "wazuh_agents.png")
    try_image(sl, img, Inches(0.3), Inches(6.35), Inches(8.7), Inches(1.1))

    footer(sl, "Wazuh s'installe en 4 commandes Docker. Agent macOS envoie les événements en temps réel. Coverage 100%.")

# ─────────────────────────────────────────────────────────
# SLIDE 13 — WAZUH Security Events + FIM
# ─────────────────────────────────────────────────────────
def slide_13():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "WAZUH — Security Events · 436 420 Events · FIM SHA-256", "13 / 20  —  Surveillance Temps Réel")

    # Left: Security Events screenshot
    img_se = os.path.join(BASE, "CVAPTWAZUH", "SECURITYEVENTS.png")
    try_image(sl, img_se, Inches(0.3), Inches(1.5), Inches(6.3), Inches(2.5))

    add_rect(sl, Inches(0.3), Inches(4.1), Inches(6.3), Inches(0.7), fill=BLUE_CARD)
    add_text(sl, "Total Events : 436 420 — journalisation complète depuis l'installation",
             Inches(0.4), Inches(4.15), Inches(6.1), Inches(0.55), size=11, bold=True, color=WHITE)

    se_info = [
        "Top Rule : Rule 550 Level 7 — Integrity checksum changed",
        "MITRE : T1565.001 Stored Data Manipulation (~95%)",
        "Agent : main-machine · macOS 15.7.4 · Active",
        "Base de toute la corrélation SIEM — IP source · agent · règle · horodatage",
    ]
    for k, info in enumerate(se_info):
        add_text(sl, f"▸  {info}", Inches(0.35), Inches(4.9) + k*Inches(0.38),
                 Inches(6.3), Inches(0.34), size=10, color=WHITE)

    # Right: FIM
    add_rect(sl, Inches(6.8), Inches(1.5), Inches(6.2), Inches(0.42), fill=BLUE)
    add_text(sl, "FIM — File Integrity Monitoring", Inches(6.85), Inches(1.52),
             Inches(6.0), Inches(0.36), size=13, bold=True, color=WHITE)

    img_fim = os.path.join(BASE, "WAZUCAPT", "wazuh_03_fim.png")
    try_image(sl, img_fim, Inches(6.8), Inches(2.0), Inches(6.2), Inches(3.0))

    fim_info = [
        ("Règle", "Rule 550 Niveau 7 — Integrity checksum changed"),
        ("root 89.44%", "Modifications binaires système macOS — /var/bin/*"),
        ("Fichiers", "/var/bin/afsa · /var/bin/apfs · /private/var/db/locationd"),
        ("Risque", "Rootkit/backdoor si modifications non légitimes"),
        ("Hash", "SHA-256 en temps réel — équivalent IDS filesystem"),
    ]
    for k, (key, val) in enumerate(fim_info):
        y = Inches(5.1) + k*Inches(0.38)
        add_text(sl, key + " :", Inches(6.8), y, Inches(1.2), Inches(0.34),
                 size=10, bold=True, color=CYAN)
        add_text(sl, val, Inches(8.05), y, Inches(4.9), Inches(0.34),
                 size=10, color=WHITE)

    footer(sl, "FIM détecterait tout rootkit ou backdoor en temps réel. 436 420 events — journalisation complète depuis l'installation.")

# ─────────────────────────────────────────────────────────
# SLIDE 14 — WAZUH CVE
# ─────────────────────────────────────────────────────────
def slide_14():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "WAZUH — Vulnerability Detection · 17 CVE · 8 High + 9 Medium", "14 / 20  —  Docker P1 Action Requise")

    img_cve = os.path.join(BASE, "WAZUCAPT", "wazuh_04_cve.png")
    try_image(sl, img_cve, Inches(0.3), Inches(1.5), Inches(6.5), Inches(3.2))

    # Table CVE droite
    cve_rows = [
        ("CVE-2019-5736",  "Docker 4.43.2", "High",   "8.6", "Container escape → root sur l'hôte via runc", RED),
        ("CVE-2019-13509", "Docker 4.43.2", "High",   "7.5", "Information disclosure",                       ORANGE),
        ("CVE-2019-16884", "Docker 4.43.2", "High",   "7.5", "AppArmor bypass",                              ORANGE),
        ("CVE-2021-21284", "Docker 4.43.2", "Medium", "6.8", "Privilege escalation",                         BLUE_LIGHT),
        ("CVE-2001-0718",  "Excel 16.107.3","High",   "7.5", "Input validation failure",                     ORANGE),
        ("CVE-2014-4715",  "lz4 1.10.0",   "Medium", "5.0", "Heap overflow",                                BLUE_LIGHT),
    ]
    hdrs_cve = ["CVE", "Package", "Sév.", "CVSS3", "Description"]
    hx_c = [Inches(7.0), Inches(8.65), Inches(10.1), Inches(10.8), Inches(11.3)]
    hw_c = [Inches(1.6), Inches(1.4), Inches(0.65), Inches(0.45), Inches(1.7)]

    for j, (h, x, w) in enumerate(zip(hdrs_cve, hx_c, hw_c)):
        add_rect(sl, x, Inches(1.5), w-Inches(0.03), Inches(0.38), fill=BLUE)
        add_text(sl, h, x+Inches(0.03), Inches(1.52), w-Inches(0.06), Inches(0.32),
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (cve, pkg, sev, cvss, desc, scol) in enumerate(cve_rows):
        bg = BLUE_CARD if i % 2 == 0 else RGBColor(0x12, 0x30, 0x6B)
        y = Inches(1.93) + i * Inches(0.44)
        for x, w in zip(hx_c, hw_c):
            add_rect(sl, x, y, w-Inches(0.03), Inches(0.4), fill=bg)
        add_text(sl, cve,  hx_c[0]+Inches(0.03), y+Inches(0.04), hw_c[0]-Inches(0.06), Inches(0.32), size=8.5, bold=True, color=CYAN, font="Courier New")
        add_text(sl, pkg,  hx_c[1]+Inches(0.03), y+Inches(0.04), hw_c[1]-Inches(0.06), Inches(0.32), size=8.5, color=WHITE)
        badge(sl, sev, hx_c[2]+Inches(0.03), y+Inches(0.04), hw_c[2]-Inches(0.07), Inches(0.3), bg=scol, size=7.5)
        add_text(sl, cvss, hx_c[3]+Inches(0.03), y+Inches(0.04), hw_c[3]-Inches(0.06), Inches(0.32), size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, hx_c[4]+Inches(0.03), y+Inches(0.04), hw_c[4]-Inches(0.06), Inches(0.32), size=8, color=WHITE)

    # KPIs sévérité
    sev_boxes = [("0", "Critical", BLUE_CARD), ("8", "High", RED), ("9", "Medium", ORANGE), ("0", "Low", GREEN)]
    bw = Inches(1.55)
    for k, (n, l, c) in enumerate(sev_boxes):
        cx = Inches(7.0) + k*(bw + Inches(0.05))
        add_rect(sl, cx, Inches(4.75), bw, Inches(0.9), fill=c)
        add_text(sl, n, cx, Inches(4.78), bw, Inches(0.5), size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, l, cx, Inches(5.26), bw, Inches(0.32), size=10, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.3), Inches(4.85), Inches(6.55), Inches(0.85), fill=RGBColor(0x6A, 0x00, 0x00))
    add_text(sl,
             "ACTION PRIORITAIRE P1 : Mise à jour Docker Engine\ndocker pull docker:latest  →  élimine 8 CVE High dont CVE-2019-5736 (CVSS3=8.6 container escape)",
             Inches(0.4), Inches(4.9), Inches(6.3), Inches(0.75), size=10, bold=True, color=WHITE)

    footer(sl, "Wazuh Vulnerability Detector compare packages vs NVD NIST automatiquement. CVE-2019-5736 = priorité absolue post-soutenance.")

# ─────────────────────────────────────────────────────────
# SLIDE 15 — MITRE ATT&CK + Rootcheck
# ─────────────────────────────────────────────────────────
def slide_15():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "WAZUH — MITRE ATT&CK · T1565.001 · Rootcheck · 4 Anomalies", "15 / 20  —  Classification SOC internationale")

    img_mitre = os.path.join(BASE, "WAZUCAPT", "wazuh_05_mitre.png")
    try_image(sl, img_mitre, Inches(0.3), Inches(1.5), Inches(6.3), Inches(2.8))

    mitre_info = [
        ("T1565.001", "Stored Data Manipulation — tactic Impact (dominant ~95%)", RED),
        ("T1562",     "Disable or Modify Tools — tactic Defense Evasion (×6)",    ORANGE),
    ]
    add_text(sl, "MITRE ATT&CK :", Inches(0.35), Inches(4.4), Inches(6.2), Inches(0.35),
             size=12, bold=True, color=CYAN)
    for k, (tactic, desc, col) in enumerate(mitre_info):
        y = Inches(4.8) + k * Inches(0.5)
        badge(sl, tactic, Inches(0.35), y, Inches(1.3), Inches(0.38), bg=col, size=9.5)
        add_text(sl, desc, Inches(1.7), y + Inches(0.04), Inches(4.9), Inches(0.34), size=10, color=WHITE)

    add_text(sl, "Corrélation : PCI DSS 11.5 · RGPD Art.25 · Privacy by Design ✓",
             Inches(0.35), Inches(5.85), Inches(6.2), Inches(0.35), size=10.5, italic=True, color=CYAN)

    # Droite Rootcheck
    add_rect(sl, Inches(6.8), Inches(1.5), Inches(6.2), Inches(0.42), fill=BLUE)
    add_text(sl, "Rootcheck — Policy Monitoring — Host-based Anomaly Detection",
             Inches(6.85), Inches(1.52), Inches(6.0), Inches(0.36), size=12, bold=True, color=WHITE)

    img_pol = os.path.join(BASE, "CVAPTWAZUH", "POLICY MONITORING.png")
    try_image(sl, img_pol, Inches(6.8), Inches(2.0), Inches(6.2), Inches(2.5))

    rootcheck = [
        ("Trojaned files",  "Détection de binaires système compromis",                       RED),
        ("Hidden processes","Processus 26061 caché — tentative de dissimulation",             ORANGE),
        ("Promiscuity",     "Interface en3 en mode promiscuous — capture réseau détectée",    ORANGE),
        ("Module",          "Host-based Anomaly Detection · Rootcheck actif · temps réel",    GREEN),
    ]
    add_text(sl, "Anomalies détectées :", Inches(6.8), Inches(4.6), Inches(6.0), Inches(0.35),
             size=12, bold=True, color=CYAN)
    for k, (key, val, col) in enumerate(rootcheck):
        y = Inches(5.0) + k*Inches(0.47)
        badge(sl, key, Inches(6.8), y, Inches(1.6), Inches(0.36), bg=col, size=9)
        add_text(sl, val, Inches(8.45), y + Inches(0.03), Inches(4.5), Inches(0.32), size=10, color=WHITE)

    footer(sl, "MITRE ATT&CK = framework universel SOC. Rootcheck a détecté des anomalies réelles sur l'agent macOS — sniffing réseau possible.")

# ─────────────────────────────────────────────────────────
# SLIDE 16 — WAZUH 5 Modules
# ─────────────────────────────────────────────────────────
def slide_16():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "WAZUH — 5 Modules de Sécurité Actifs — CDC §3.3 Couvert", "16 / 20")

    modules = [
        ("Security Events",    BLUE_LIGHT, [
            "Tableau central — niveaux 3→15",
            "IP source · agent · règle · horodatage",
            "Base de toute la corrélation SIEM",
            "436 420 events journalisés",
        ]),
        ("MITRE ATT&CK",       BLUE, [
            "T1499 Impact ×1594 · T1562 Evasion ×6",
            "Taxonomie universelle SOC",
            "Corrélation automatique techniques",
            "PCI DSS 11.5 · RGPD Art.25 ✓",
        ]),
        ("GDPR Module",        BLUE_CARD, [
            "Art.15 accès · Art.17 effacement",
            "Art.33 violation",
            "Complète nos auditLogs applicatifs",
            "Conformité infrastructure",
        ]),
        ("Integrity Monitoring",RGBColor(0x1A, 0x4D, 0x8A), [
            "2 500+ events · /bin/bash modifié",
            "Alerte rootkit ou backdoor",
            "Équivalent IDS filesystem",
            "Hash SHA-256 en temps réel",
        ]),
        ("Vulnerabilities",    BLUE_LIGHT, [
            "17 CVE · 8 High + 9 Medium",
            "CVE-2019-5736 CVSS3=8.6 → P1",
            "Docker dominant",
            "Scan packages vs NVD NIST auto",
        ]),
    ]

    mw = Inches(2.4)
    mh = Inches(4.4)
    mx_start = Inches(0.3)
    gap = Inches(0.15)

    for i, (name, col, items) in enumerate(modules):
        mx = mx_start + i * (mw + gap)
        add_rect(sl, mx, Inches(1.5), mw, mh, fill=col)
        add_rect(sl, mx, Inches(1.5), mw, Inches(0.52), fill=RGBColor(0x0A, 0x20, 0x60))
        add_text(sl, name, mx, Inches(1.55), mw, Inches(0.44),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for k, item in enumerate(items):
            add_text(sl, f"▸  {item}", mx+Inches(0.08), Inches(2.15)+k*Inches(0.5),
                     mw-Inches(0.12), Inches(0.45), size=9.5, color=WHITE)

    badge(sl, "5 modules combinés : Menaces + RGPD + CVE + FIM + Journalisation — CDC §3.3 Couvert à 100%",
          Inches(0.3), Inches(6.1), Inches(12.73), Inches(0.48), bg=BLUE_CARD, size=11)

    footer(sl, "Coverage CDC §3.3 complète — Wazuh assure 3 des 8 exigences infrastructure en complément de la couche applicative.")

# ─────────────────────────────────────────────────────────
# SLIDE 17 — DÉFENSE EN PROFONDEUR
# ─────────────────────────────────────────────────────────
def slide_17():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "DÉFENSE EN PROFONDEUR — 2 Couches — CDC §3.3 à 100%", "17 / 20  —  Defense in Depth")

    # Couche 1
    add_rect(sl, Inches(0.3), Inches(1.5), Inches(6.15), Inches(4.8), fill=BLUE_CARD)
    add_rect(sl, Inches(0.3), Inches(1.5), Inches(6.15), Inches(0.52), fill=BLUE_LIGHT)
    add_text(sl, "COUCHE 1 — APPLICATIVE", Inches(0.4), Inches(1.54), Inches(5.9), Inches(0.44),
             size=14, bold=True, color=WHITE)
    add_text(sl, "auditLogs Firestore + /monitoring", Inches(0.4), Inches(1.98), Inches(5.9), Inches(0.3),
             size=10, italic=True, color=CYAN)

    c1_items = [
        "WAF — SQLi · XSS · PathTraversal · CmdInj · Agents suspects",
        "9 types audit logs immuables Firestore (append-only)",
        "Score /100 temps réel + 4 alertes automatiques",
        "Scanner DAST 12 tests OWASP — Score 92/100",
        "Dashboard SIEM 3 onglets (admin only · 60s refresh)",
        "RGPD Art.15/16/17/33 · AES-256-CBC · HTTPS/HSTS",
    ]
    for k, item in enumerate(c1_items):
        add_text(sl, f"✓  {item}", Inches(0.4), Inches(2.35) + k*Inches(0.5),
                 Inches(5.9), Inches(0.45), size=11, color=WHITE)

    # Couche 2
    add_rect(sl, Inches(6.65), Inches(1.5), Inches(6.35), Inches(4.8), fill=BLUE_CARD)
    add_rect(sl, Inches(6.65), Inches(1.5), Inches(6.35), Inches(0.52), fill=BLUE)
    add_text(sl, "COUCHE 2 — INFRASTRUCTURE", Inches(6.75), Inches(1.54), Inches(6.1), Inches(0.44),
             size=14, bold=True, color=WHITE)
    add_text(sl, "Wazuh SIEM Agent actif", Inches(6.75), Inches(1.98), Inches(6.1), Inches(0.3),
             size=10, italic=True, color=CYAN)

    c2_items = [
        "Agent 001 actif — coverage 100% (Debian 12)",
        "213 événements détectés en < 24h sans config extra",
        "CIS Benchmark Debian 12 — audit configuration",
        "Module MITRE ATT&CK — T1499/T1562 classifiés",
        "Module GDPR — Art.15/17/33 infrastructure",
        "Vulnerability Detection — 17 CVE/CVSS détectées",
    ]
    for k, item in enumerate(c2_items):
        add_text(sl, f"✓  {item}", Inches(6.75), Inches(2.35) + k*Inches(0.5),
                 Inches(6.1), Inches(0.45), size=11, color=WHITE)

    # Flèche ↔
    add_text(sl, "⟺", Inches(5.9), Inches(3.6), Inches(0.85), Inches(0.6),
             size=28, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # Résumé
    add_rect(sl, Inches(0.3), Inches(6.45), Inches(12.7), Inches(0.58), fill=BLUE)
    add_text(sl,
             "→  CDC §3.3 couvert à 100% : Accès sécurisé (RBAC) · Conformité RGPD · Journalisation · HTTPS · Monitoring",
             Inches(0.45), Inches(6.5), Inches(12.4), Inches(0.48), size=12, bold=True, color=WHITE)

    footer(sl, "Defense in Depth = 2 couches complémentaires. Applicatif sécurise les accès web. Wazuh surveille l'infrastructure système.")

# ─────────────────────────────────────────────────────────
# SLIDE 18 — SYNTHÈSE CDC §3.3
# ─────────────────────────────────────────────────────────
def slide_18():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "SYNTHÈSE — Couverture CDC §3.3 à 100% — 8/8 Exigences FAITES", "18 / 20")

    reqs = [
        ("Accès sécurisé id/mdp",   "1·2·3·5", "JWT HS256 · bcrypt 12r · WAF · rate limit · lockout 5 échecs", "Security Events", GREEN),
        ("HTTPS obligatoire",        "3",        "Firebase HTTPS auto · Helmet HSTS max-age=31536000",           "—",               GREEN),
        ("RGPD conforme",            "5·6·8",    "DATA_EXPORT · DATA_ANONYMIZE · AES-256 · Art.15/16/17/33",    "Module GDPR",     GREEN),
        ("RBAC — Gestion rôles",     "5·6·8·9",  "checkRole() Firestore · 4 rôles · Firestore Rules deny",       "Security Events", GREEN),
        ("Journalisation actions",   "6·7·8",    "9 events immuables · SIEM dashboard · serverTimestamp()",     "213+ events",     GREEN),
        ("Détection vulnérabilités", "Wazuh",    "17 CVE identifiées · Plan correctif prioritisé",               "Vuln Detector",   GREEN),
        ("Intégrité fichiers",       "Wazuh",    "FIM 2500+ · Rule 550 · Hash SHA-256",                          "Integrity Mon.",  GREEN),
        ("Détection menaces",        "Wazuh",    "MITRE ATT&CK T1499/T1562 · 1600+ alertes",                    "MITRE ATT&CK",    GREEN),
    ]

    hdrs_s = ["EXIGENCE CDC §3.3", "Réf.", "IMPLÉMENTATION", "Module Wazuh", "STATUT"]
    hx_s  = [Inches(0.3),  Inches(3.6),  Inches(4.25), Inches(10.4), Inches(11.75)]
    hw_s  = [Inches(3.25), Inches(0.6),  Inches(6.1),  Inches(1.3),  Inches(1.2)]

    for j, (h, x, w) in enumerate(zip(hdrs_s, hx_s, hw_s)):
        add_rect(sl, x, Inches(1.5), w-Inches(0.04), Inches(0.42), fill=BLUE)
        add_text(sl, h, x+Inches(0.04), Inches(1.52), w-Inches(0.08), Inches(0.36),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (req, ref, impl, wazuh, col) in enumerate(reqs):
        bg = BLUE_CARD if i % 2 == 0 else RGBColor(0x12, 0x30, 0x6B)
        y = Inches(1.97) + i * Inches(0.61)
        for x, w in zip(hx_s, hw_s):
            add_rect(sl, x, y, w-Inches(0.04), Inches(0.57), fill=bg)
        add_text(sl, req,   hx_s[0]+Inches(0.04), y+Inches(0.07), hw_s[0]-Inches(0.08), Inches(0.44), size=10, bold=True, color=WHITE)
        add_text(sl, ref,   hx_s[1]+Inches(0.04), y+Inches(0.07), hw_s[1]-Inches(0.08), Inches(0.44), size=9,  color=CYAN, align=PP_ALIGN.CENTER)
        add_text(sl, impl,  hx_s[2]+Inches(0.04), y+Inches(0.07), hw_s[2]-Inches(0.08), Inches(0.44), size=9.5, color=WHITE)
        add_text(sl, wazuh, hx_s[3]+Inches(0.04), y+Inches(0.07), hw_s[3]-Inches(0.08), Inches(0.44), size=9.5, color=CYAN)
        badge(sl, "✓ FAIT", hx_s[4]+Inches(0.1), y+Inches(0.1), hw_s[4]-Inches(0.2), Inches(0.36), bg=col, size=10)

    footer(sl, "CDC §3.3 couvert à 100% — 8 exigences · 8 implémentées · Toutes vérifiables en live devant le jury.")

# ─────────────────────────────────────────────────────────
# SLIDE 19 — KPI
# ─────────────────────────────────────────────────────────
def slide_19():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)
    add_title_bar(sl, "RÉCAPITULATIF KPI — Chiffres Clés du Projet", "19 / 20")

    kpis = [
        ("436 420",  "Security Events Wazuh",       BLUE_LIGHT),
        ("100/100",  "Score Dashboard App",          GREEN),
        ("92/100",   "Score DAST (12 tests)",        GREEN),
        ("17 CVE",   "8 High + 9 Medium",            RED),
        ("8/8",      "Exigences CDC §3.3",           GREEN),
        ("9",        "Types AuditLogs immuables",    BLUE),
        ("10",       "Catégories OWASP vérifiées",   BLUE_LIGHT),
        ("6",        "Articles RGPD couverts",       BLUE),
        ("4",        "Rôles RBAC moindre privilège", BLUE_CARD),
        ("12r",      "bcrypt rounds anti-brute",     GREEN),
        ("AES-256",  "Chiffrement données sensibles",BLUE_LIGHT),
        ("100%",     "Agents Wazuh coverage",        GREEN),
    ]

    kw = Inches(2.9)
    kh = Inches(1.3)
    gap = Inches(0.18)
    kx_start = Inches(0.3)

    for i, (num, lbl, col) in enumerate(kpis):
        row = i // 4
        col_idx = i % 4
        kx = kx_start + col_idx * (kw + gap)
        ky = Inches(1.55) + row * (kh + Inches(0.15))
        add_rect(sl, kx, ky, kw, kh, fill=col)
        add_text(sl, num, kx, ky + Inches(0.1), kw, Inches(0.7),
                 size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, kx, ky + Inches(0.8), kw, Inches(0.4),
                 size=10, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.3), Inches(6.3), Inches(12.73), Inches(0.72), fill=BLUE_CARD)
    add_text(sl,
             "CVE-2019-5736 CVSS3=8.6  ·  T1565.001 Stored Data Manipulation ~95%  ·  Rule 550 Level 7 FIM dominant  ·  Defense in Depth 2 couches",
             Inches(0.45), Inches(6.36), Inches(12.4), Inches(0.58),
             size=11, italic=True, color=CYAN, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
# SLIDE 20 — MERCI
# ─────────────────────────────────────────────────────────
def slide_20():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl, NAVY)

    # Header fin
    add_rect(sl, 0, 0, W, Inches(0.38), fill=BLUE)
    add_text(sl, "PFE CYBERSÉCURITÉ / CYBERDÉFENSE — YNOV Campus 2026",
             Inches(0.3), Inches(0.04), Inches(12.7), Inches(0.3),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "Merci",
             Inches(0.5), Inches(0.85), Inches(12.3), Inches(1.6),
             size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(2), Inches(2.4), Inches(9.33), Inches(0.06), fill=BLUE_LIGHT)

    add_text(sl, "Questions & Démonstration Live",
             Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.7),
             size=28, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # Citation
    add_rect(sl, Inches(1.2), Inches(3.45), Inches(10.93), Inches(1.65), fill=BLUE_CARD)
    add_text(sl,
             '"Le monitoring couvre intégralement le CDC §3.3 sur deux niveaux :\n'
             'monitoring applicatif (WAF · 9 audit logs Firestore · dashboard SIEM · DAST 12 tests)\n'
             'et monitoring infrastructure Wazuh (213 événements · MITRE ATT&CK · GDPR · CIS Benchmark).\n'
             'Ensemble, ils assurent une défense en profondeur conforme RGPD."',
             Inches(1.4), Inches(3.52), Inches(10.53), Inches(1.5),
             size=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "Amine BAHOU  /  Anass Akker",
             Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.6),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "PFE Bachelor Cybersécurité / Cyberdéfense — YNOV Campus 2026",
             Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.4),
             size=14, color=CYAN, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.3), Inches(6.5), Inches(12.73), Inches(0.58), fill=BLUE)
    tags = "OWASP Top 10  ·  RGPD UE 2016/679  ·  CDC §3.3 100%  ·  Wazuh 4.7.4  ·  JWT HS256  ·  AES-256-CBC  ·  Defense in Depth"
    add_text(sl, tags, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.44),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
# GÉNÉRATION
# ─────────────────────────────────────────────────────────
print("Génération des slides...")
slide_01(); print("  ✓ Slide 01 — Titre")
slide_02(); print("  ✓ Slide 02 — Contexte & Objectifs")
slide_03(); print("  ✓ Slide 03 — Pipeline Sécurité")
slide_04(); print("  ✓ Slide 04 — OWASP Top 10")
slide_05(); print("  ✓ Slide 05 — WAF")
slide_06(); print("  ✓ Slide 06 — RBAC")
slide_07(); print("  ✓ Slide 07 — RGPD")
slide_08(); print("  ✓ Slide 08 — Journalisation & DAST")
slide_09(); print("  ✓ Slide 09 — [Screenshot] Dashboard Principal")
slide_10(); print("  ✓ Slide 10 — [Screenshot] WAF Monitoring")
slide_11(); print("  ✓ Slide 11 — [Screenshot] Dashboard SIEM")
slide_12(); print("  ✓ Slide 12 — Wazuh Architecture")
slide_13(); print("  ✓ Slide 13 — [Screenshot] Security Events + FIM")
slide_14(); print("  ✓ Slide 14 — [Screenshot] CVE Detection")
slide_15(); print("  ✓ Slide 15 — [Screenshot] MITRE ATT&CK + Rootcheck")
slide_16(); print("  ✓ Slide 16 — Wazuh 5 Modules")
slide_17(); print("  ✓ Slide 17 — Défense en Profondeur")
slide_18(); print("  ✓ Slide 18 — Synthèse CDC §3.3")
slide_19(); print("  ✓ Slide 19 — KPI")
slide_20(); print("  ✓ Slide 20 — Merci")

out = os.path.join(BASE, "PRESENTATION_SECURITE_YNOV_2026.pptx")
prs.save(out)
print(f"\n✅ Fichier généré : {out}")
